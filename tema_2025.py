# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt, Cm
from pptx.chart.data import ChartData
from pptx.enum.text import MSO_ANCHOR
from PIL import Image
from datetime import datetime
import copy
import random
import requests
import os
from sheet2dict import Worksheet
from collections import defaultdict
from collections import Counter
import pandas as pd
import openpyxl
from openpyxl import Workbook
from openpyxl.styles import (
    PatternFill, Border, Side,
    Alignment, Font, GradientFill
)
from openpyxl.comments import Comment
from openpyxl import load_workbook
import warnings
import glob
import json
from itertools import zip_longest
import sys
import re
from openpyxl.worksheet.datavalidation import DataValidation
from dateutil import parser



print('====================================================================================================')
print('Добрый день, скрипт автоматизации реестра отклонений по ITSM и АС готов к работе, выберите действие:')
print('====================================================================================================')
print('**************************************************************************************************************')
print('Создание Анкет ИТ-услуг для оценки - введите цифру: 0 ')
print('Генерация реестра - ITSM + АС или ITSM -  введите цифру: 1')
print('Для подсчета ежеквартального инкремента после проведенного ЕОИТ в ДЗО начиная с - 2025г. - введите цифру: 2')
# print('Динамика изменений по ITSM + АС - введите цифру: 2')
# print('Генерация реестра - по АС введите цифру: 3')
print('**************************************************************************************************************')

# ==================================Общий обработчик ДАТ для 4-х вариантов ЕДК (OP, BO, MC/BC Critical) если они проставлены или вовсе не проставлены ===========================

# Вывод текущего квартала и года
def get_current_quarter():
    now = datetime.now()
    current_quarter = (now.month - 1) // 3 + 1
    current_year = now.year
    return f"{current_quarter}Q{current_year}"


# Пример использования
current_quarter = get_current_quarter()


# ================================================

# ====================Сравнение дат нужно для сравнение с текущей даты============================
def extract_quarter_year(q_str):
    """Извлекает квартал и год из строки формата 'QXYYYY'"""
    quarter, year = q_str.split('Q')
    return int(quarter), int(year)


def find_future_quarters(current_q, all_quarters):
    """Находит кварталы, которые больше текущего"""
    try:
        current_quarter, current_year = extract_quarter_year(current_q)
        future_quarters = []

        for q in all_quarters:
            try:
                quarter, year = extract_quarter_year(q)
                if (year > current_year) or (year == current_year and quarter > current_quarter):
                    future_quarters.append(q)
            except ValueError as e:
                print(f"Используется некорректная дата в ЕДК '{q}': {e}")
                continue

        return sorted(future_quarters)  # Возвращаем отсортированный список
    except ValueError as e:
        print(f"Ошибка обработки текущего квартала '{current_q}': {e}")
        return []


# ===============================================


# Словарь для преобразования русских/английских месяцев
MONTHS = {
    'янв': 1, 'фев': 2, 'мар': 3, 'апр': 4, 'май': 5, 'июн': 6,
    'июл': 7, 'авг': 8, 'сен': 9, 'окт': 10, 'ноя': 11, 'дек': 12,
    'jan': 1, 'feb': 2, 'mar': 3, 'apr': 4, 'may': 5, 'jun': 6,
    'jul': 7, 'aug': 8, 'sep': 9, 'oct': 10, 'nov': 11, 'dec': 12
}


def date_to_quarter_year(date_str):
    date_str = date_str.strip().lower()

    # Проверка на квартальные форматы (1кв.2025, q1 25, Q1'25)
    quarter_match = re.match(r'(?:q|кв)(\d)[\'\s\.]*(\d{2,4})', date_str) or \
                    re.match(r'(\d)\s*(?:кв|q)[\'\s\.]*(\d{2,4})', date_str) or \
                    re.match(r'(\d)\s*квартал\s*(\d{4})', date_str) or \
                    re.match(r'(\d{4})\s*(?:q|кв)\s*(\d)', date_str)  # 2025 q3

    if quarter_match:
        # Для формата "2025 q3" группы будут в другом порядке
        if quarter_match.re.pattern == r'(\d{4})\s*(?:q|кв)\s*(\d)':
            year = int(quarter_match.group(1))
            quarter = int(quarter_match.group(2))
        else:
            quarter = int(quarter_match.group(1))
            year = int(quarter_match.group(2))

        return f"{quarter}Q{year if year >= 100 else year + 2000}"

    # Проверка на текстовый месяц (июн.2025, Jun 2025, 15-июн-2025)
    month_match = re.match(r'([а-яa-z]+)\.?\s*(\d{4})', date_str) or \
                  re.match(r'(\d{1,2})-([а-яa-z]+)-(\d{4})', date_str) or \
                  re.match(r'(\d{1,2})\s+([а-яa-z]+)\s+(\d{4})', date_str)

    if month_match:
        groups = month_match.groups()
        if len(groups) == 2:  # формат "июн 2025"
            month_str = groups[0][:3]
            year = int(groups[1])
        else:  # формат "15-июн-2025" или "15 июн 2025"
            month_str = groups[1][:3]
            year = int(groups[2])

        if month_str in MONTHS:
            quarter = (MONTHS[month_str] - 1) // 3 + 1
            return f"{quarter}Q{year}"

    # Парсинг обычных дат (31.03.2025, 2025-07-31 и др.)
    try:
        parsed_date = parser.parse(date_str, dayfirst=True)
        quarter = (parsed_date.month - 1) // 3 + 1
        return f"{quarter}Q{parsed_date.year}"
    except ValueError:
        return None  # Неизвестный формат


def add_quarters(initial_date, quarters_to_add):
    # Парсим входную строку (формат "1Q2025" или "Q12025")
    if initial_date[0] == 'Q':  # формат Q12025
        quarter = int(initial_date[1])
        year = int(initial_date[2:])
    else:  # формат 1Q2025
        quarter = int(initial_date[0])
        year = int(initial_date[2:])

    # Вычисляем новые квартал и год
    total_quarters = quarter + quarters_to_add - 1
    new_year = year + total_quarters // 4
    new_quarter = (total_quarters % 4) + 1

    return f"{new_quarter}Q{new_year}"


# =====================================================================================================================
# ===============================================Выборка наименьшей даты из всех========================================================
def find_min_quarter(result_mas):
    min_year = float('inf')
    min_quarter = float('inf')
    min_quarter_str = None

    for q in result_mas:
        if q is None:  # Пропускаем None значения
            continue

        # Разбираем строку на квартал и год (поддерживаем оба формата: "2Q2025" и "Q22025")
        try:
            if 'Q' in q:
                parts = q.split('Q')
                if len(parts) == 2:
                    # Формат "2Q2025"
                    quarter = int(parts[0])
                    year = int(parts[1])
                else:
                    # Формат "Q22025"
                    quarter = int(q[1])
                    year = int(q[2:])

                # Сравниваем сначала год, потом квартал
                if (year < min_year) or (year == min_year and quarter < min_quarter):
                    min_year = year
                    min_quarter = quarter
                    min_quarter_str = f"{quarter}Q{year}"  # Сохраняем в едином формате
        except (ValueError, IndexError, AttributeError):
            continue  # Пропускаем некорректные строки

    return min_quarter_str


# Вызываем функцию и получаем минимальный квартал
# min_quarter = find_min_quarter(result_mas)
# print("Самый ранний квартал:", min_quarter)  # Выведет Q12025


# ===============================================Концовка обработчика Дат========================================================

dict_inn_company = {"9705118142":["Подразделения вне блоков B2C","Купер"],
                    "9701048328":["Подразделения вне блоков B2C","МегаМаркет"],
                    "7811554010":["Подразделения вне блоков B2C","Самокат"],
                    "7736322345":["Подразделения вне блоков B2C","СберЛогистика"],
                    "7730262964":["EdTech","СберОбразование"],
                    "7736316133":["EdTech","Школа 21"],
                    "7704865540":["E-Health","ЕАптека"],
                    "9710011437":["E-Health","СберЗдоровье"],
                    "9731065465":["E-Health","СберМедИИ"],
                    "9705124940":["GR","СберПраво"],
                    "7708328948":["Media & Ads","Звук"],
                    "7801445445":["Media & Ads","Звук Бизнес"],
                    "7814665871":["Media & Ads","ОККО"],
                    "7725243282":["Media & Ads","Рамблер"],
                    "7736319695":["Media & Ads","СберМаркетинг"],
                    "7736659589":["ДРПА","АктивБизнесКонсалт"],
                    "7736303529":["ДРПА","АктивБизнесТехнологии"],
                    "7736581290":["ДРПА","Сбербанк Капитал"],
                    "5405276278":["КИБ","2ГИС"],
                    "7736641983":["КИБ","Деловая среда"],
                    "9731062087":["КИБ","Дома"],
                    "7714843760":["КИБ","Инсейлс"],
                    "7709969870":["КИБ","Работа.ру"],
                    "7707308480":["КИБ","Сбер А"],
                    "7730269550":["КИБ","Сбер Бизнес Софт"],
                    "7730241227":["КИБ","СберАналитика"],
                    "7707009586":["КИБ","Сбербанк Лизинг"],
                    "7802754982":["КИБ","Сбербанк Факторинг"],
                    "7801392271":["КИБ","СберКорус"],
                    "7709688816":["КИБ","СберРешения"],
                    "7730262971":["КИБ","СберТаксФри"],
                    "7736612855":["КИБ","Стратеджи Партнерс Групп"],
                    "7727381792":["КИБ","Фьюэл-Ап"],
                    "7730261382":["КИБ","Цифровые решения регионов"],
                    "9709108748":["ЛК","Пульс"],
                    "7736128605":["ЛК","СберУниверситет"],
                    "7736249247":["Подразделения вне блоков B2C","ДомКлик"],
                    "9709054813":["Подразделения вне блоков B2C","СберАвто"],
                    "7736264044":["Подразделения вне блоков B2C","СберМобайл"],
                    "7704314221":["Подразделения вне блоков B2C","Ситидрайв"],
                    "9709078370":["Подразделения вне блоков B2C","Центр новых финансовых сервисов"],
                    "7702770003":["Развитие клиентского опыта B2C","СберСпасибо"],
                    "7710561081":["Риски","ОКБ"],
                    "7736324991":["Сервисы","БАРУС"],
                    "400014449":["Сервисы","Манжерок"],
                    "7730245060":["Сервисы","Медэксперт Плюс"],
                    "7729276546":["Сервисы","Московский городской Гольф Клуб"],
                    "9103007830":["Сервисы","Мрия"],
                    "7736663049":["Сервисы","СберСервис"],
                    "7708229993":["Сервисы","СовТех"],
                    "7720427871":["Сервисы","СТК"],
                    "9709073460":["Сеть продаж","СберМегаМаркетРитейл"],
                    "5024093941":["Строительство","Рублево-Архангельское"],
                    "9731026963":["Строительство","Смарт Констракшн"],
                    "6439098794":["Строительство","Инфотех Балаково"],
                    "7736632467":["Технологии","Сбербанк-Технологии"],
                    "7736279160":["Технологическое развитие","Cloud.ru"],
                    "9725045830":["Технологическое развитие","АВТОТЕХ"],
                    "7725745476":["Технологическое развитие","Живой Сайт"],
                    "7730253720":["Технологическое развитие","СалютДевайсы"],
                    "7805093681":["Технологическое развитие","ЦРТ"],
                    "7727718421":["Транзакционный банкинг B2C","Расчетные решения"],
                    "7750005860":["Транзакционный банкинг B2C","Расчетные решения НКО"],
                    "9702027017":["Транзакционный банкинг B2C","СберТройка"],
                    "9715225506":["Транзакционный банкинг B2C","Эвотор"],
                    "7750005725":["Транзакционный банкинг B2C","Юмани"],
                    "7725352740":["УБ","НПФ Сбербанка"],
                    "9725000621":["УБ","Пенсионные решения"],
                    "7736618039":["УБ","Рыночный Спецдепозитарий"],
                    "7744002123":["УБ","Сбербанк страхование жизни СК"],
                    "7706810747":["УБ","Сбербанк страхование СК"],
                    "7730257675":["УБ","Современные Фонды Недвижимости"],
                    "7706810730":["УБ","Страховой брокер Сбербанка"],
                    "7710183778":["УБ","Управляющая компания Первая"],
                    "7736252313":["ЦПНД","Цифровые технологии"]
                    }






automat_dzo = int(input())

if automat_dzo == 0:
    ws_itsm = Worksheet()
    book_znanii_as = {}
    print('!!! Чтобы сформировать анкеты по АС необходимо, чтобы "Анкета_объектов_оценки" находилась в общей папке(Опросник, База_знаний, папка AS)')
    print()
    script_dir = os.path.abspath(os.path.dirname(__file__))
    # adress_as = os.path.join(script_dir, 'AS')
    warnings.simplefilter("ignore")


    file_list = glob.glob(os.path.join(script_dir, 'no_hand', '*.xlsx'))
    mas_files_list = []
    for file in file_list:
        if 'Анкета_объектов' in file:
            mas_files_list.append(file)
            print('В папке найден файл Анкета_объектов_оценки:', *mas_files_list)
            break
    else:
        print('В папке отсутствует файл Анкета_объектов_оценки, скрипт не сформирует АС анкеты для оценки !!!')
        sys.exit()
    itsm_rezult = os.path.join(os.path.dirname(__file__), *mas_files_list)
    workbook_dtn = openpyxl.load_workbook(itsm_rezult)
    sheet = workbook_dtn.active


    # Ищем по всему Excel Колонку где имеется заголовок КОД!!!!=====================

    code_cell = None
    for row in sheet.iter_rows():
        for cell in row:
            if cell.value and "Код" in str(cell.value):
                code_cell = cell
                break
        if code_cell:
            break
    if code_cell:
        column_letter = code_cell.column_letter  # 'F'
        row_number = code_cell.row               # 3
        print(f"Буква столбца: {column_letter}, Номер строки: {row_number}")
        # print(f"Ячейка с 'Код' найдена: {code_cell.coordinate}")  # Например, "F3"

        first_column = sheet[f"{column_letter}"]
        count_dtn_code = 0
        for cell in first_column:
            if cell.value != 'Код':
                count_dtn_code += 1
            else:
                break

        ws_itsm.xlsx_to_dict(path=itsm_rezult, select_sheet='Реестр объектов обследования', data_only=True, skiprows=count_dtn_code)
    else:
        print("Ячейка с 'Код' не найдена!")


    for f in ws_itsm.sheet_items:
        juice_znaniya_as = f.setdefault('Код')
        juice_znaniya_type = f.setdefault('Тип объекта')
        juice_znaniya_answer_as = f.setdefault('Название')
        juice_znaniya_otklon_as = f.setdefault('Критичность сервиса')
        juice_znaniya_recomend_as = f.setdefault('Описание')


        #TODO: Добавление в словарь из базы знаний
        book_znanii_as.setdefault(juice_znaniya_as, [])
        book_znanii_as[juice_znaniya_as].append(juice_znaniya_type)
        book_znanii_as[juice_znaniya_as].append(juice_znaniya_answer_as)
        book_znanii_as[juice_znaniya_as].append(juice_znaniya_otklon_as)
        book_znanii_as[juice_znaniya_as].append(juice_znaniya_recomend_as)


    dtn_as_all = {}
    for i, v in book_znanii_as.items():
        if v[0] == 'Прикладной сервис':
            if i not in dtn_as_all:
                dtn_as_all[i] = []
            dtn_as_all[i].append(v[1:])

    real_name = input('Введите наименование компании: ')
    real_block = input('Введите название блока: ')
    real_exp_dtn = input('Эксперт по надежности (ДТН): ')
    real_exp_dka = input('Куратор по архитектуре (ДКА): ')

    def dtn_as(dtn_as_all, real_name, real_block, real_exp_dtn, real_exp_dka):
        criticals = {'Office Productivity':'OP',
                     'Business Operational':'BO',
                     'Business Critical':'BC',
                     'Mission Critical':'MC'}
        path = os.path.join(os.path.dirname(__file__), "no_hand", "Анкета - ИТ-услуга.xlsx")
        path_two = os.path.join(os.path.dirname(__file__), "no_hand", "Анкета - ИТ-для скрипта.xlsx")
        real_path = os.path.join(os.path.dirname(__file__), "AS")

        wb = openpyxl.load_workbook(path)
        wb_two = openpyxl.load_workbook(path_two)
        ws = wb["Характеристики ИТ-услуги"]
        ws_techno = wb["Технологии"]
        ws_techno_two = wb_two["Технологии"]
        ws_tech = wb["Компоненты ИТ-услуги"]
        ws_tech_two = wb_two["Компоненты ИТ-услуги"]

        #В этом словаре добавлены диапазоны для ячеек колонки C (Лист "Технология") из листа "Справочники"

        # validations = ws_techno.data_validations.dataValidation
        # validation_dict = {v.sqref: v for v in validations}

        mas_cell_tech = []
        validations_tech = ws_tech_two.data_validations.dataValidation
        validation_dict_tech = {str(v.sqref): v for v in validations_tech}
        validation_dict_tech_new = {}
        updated_validation_dict_tech = {}
        for row in range(1, 25):
            cell = f'A{row}'
            mas_cell_tech.append(cell)

        # Проходим по всем валидациям и добавляем каждую ячейку отдельно для Листа -----> Компоненты ИТ-услуги
        for validation in validations_tech:
            for i_cell in validation.cells:
                if 'Справочники' in validation.formula1:
                    validation_dict_tech_new[str(i_cell)] = validation

        for key, value in validation_dict_tech_new.items():
            # Извлекаем начальный и конечный индексы строки
            start_new, end_new = key.split(':')
            # Определяем первый индекс строки
            first_index = int(start_new[1:])
            # Создаем новый ключ как список ячеек
            new_key = [f'{start_new[0]}{i}' for i in range(first_index, int(end_new[1:]) + 1)]
            # Добавляем новое значение в словарь с новым ключом
            updated_validation_dict_tech[tuple(new_key)] = value
        all_cells_tech = set()
        for cell in validation_dict_tech.keys():
            if 'A' in cell:
                all_cells_tech.add(cell)
        result_techno = {', '.join(part.replace(':', ',') for part in item.split()) for item in all_cells_tech}

        # Преобразуем множество в строку
        result_techno_str = ','.join(result_techno)
        start, end = result_techno_str.split(',')
        # Определение минимального и максимального значений строк
        start_row = int(start[1:])
        end_row = int(end[1:])

        # Формирование множества всех ячеек между минимальным и максимальным значением строки
        cells_tech = {f'A{i}' for i in range(1, end_row + 1)}
        cells_tech_l = {f'L{i}' for i in range(1, end_row + 1)}

        for cell_address in cells_tech:
            for keys, formulas in updated_validation_dict_tech.items():
                if cell_address in keys and "Справочники" in formulas.formula1:
                    key_formula = formulas.formula1
                    cleaned_string = key_formula.replace('"', '')
                    dv = DataValidation(type="list", formula1=cleaned_string, allow_blank=True)
                    dv.add(ws_tech[cell_address])
                    ws_tech.add_data_validation(dv)

        for cell_address in cells_tech_l:
            for keys, formulas in updated_validation_dict_tech.items():
                if cell_address in keys and "Справочники" in formulas.formula1:
                    key_formula = formulas.formula1
                    cleaned_string = key_formula.replace('"', '')
                    dv = DataValidation(type="list", formula1=cleaned_string, allow_blank=True)
                    dv.add(ws_tech[cell_address])
                    ws_tech.add_data_validation(dv)

        #Для Листа Технологии

        mas_cell = []
        validations = ws_techno_two.data_validations.dataValidation
        validation_dict = {str(v.sqref): v for v in validations}

        for row in range(1, 115):
            cell = f'C{row}'
            mas_cell.append(cell)

        # Проходим по всем валидациям и добавляем каждую ячейку отдельно для Листа Технологии
        for validation in validations:
            for i_cell in validation.cells:
                if 'Справочники' in validation.formula1:
                    validation_dict[str(i_cell)] = validation
        all_cells = set()
        for cell in validation_dict.keys():
            if 'C' in cell:
                all_cells.add(cell)
        result = {', '.join(part.replace(':', ',') for part in item.split()) for item in all_cells}
        #Итоговое множество всех элементов у которых имеется диапазон
        results = set()
        for item in result:
            results.update(item.replace(' ', '').split(','))

        # Создаем словарь, где ключи - уникальные адреса ячеек, а значения - списки допустимых адресов
        validat = {}
        for address in results:
            if address not in validat:
                validat[address] = [address]
            else:
                continue

        for cell_address, valid_addresses in validat.items():
            for keys, formulas in validation_dict.items():

                if cell_address == keys and "Справочники" in formulas.formula1:
                    key_formula = formulas.formula1
                    cleaned_string = key_formula.replace('"', '')


                    dv = DataValidation(type="list", formula1=cleaned_string, allow_blank=True)
                    dv.add(ws_techno[cell_address])
                    ws_techno.add_data_validation(dv)
        # Получаем текущее время

        # Извлекаем текущий год

        for i, (key, values) in enumerate(dtn_as_all.items()):
            ws['C2'].value = real_name
            ws['C4'].value = datetime.now().year
            ws['C5'].value = real_block
            ws['C18'].value = real_exp_dtn
            ws['C19'].value = real_exp_dka
            ws['C3'].value = values[0][0]
            ws['C6'].value = key
            if values[0][1] in criticals:
                ws['C7'].value = criticals[values[0][1]]
            ws['C9'].value = values[0][2]

            name_ass = re.sub(r'[\\/:*?"<>|]', '', values[0][0])
            output_path = os.path.join(real_path, f'Анкета - ИТ-услуга-{name_ass}.xlsx')
            # Сохраняем и ЗАКРЫВАЕМ файл
            wb.save(output_path)
            wb.close()  # Важно!


    if __name__ == "__main__":
        dtn_as(dtn_as_all, real_name, real_block, real_exp_dtn, real_exp_dka)
#=======================================================================================================================
if automat_dzo == 1:
    ws = Worksheet()
    ws_two = Worksheet()
    ws_as = Worksheet()
    print()

    script_dir_all = os.path.abspath(os.path.dirname(__file__))
    mas_files_list = []
    keywords = ["Опросник", "База_знаний"]


    file_list = glob.glob(os.path.join(script_dir_all, '*.xlsx'))
    for file in file_list:
        if any(keyword in file for keyword in keywords):
            mas_files_list.append(file)


    oprosnik_files = [file for file in mas_files_list if 'Опросник' in file]
    baza_znaniy_files = [file for file in mas_files_list if 'База_знаний' in file]

    baza = os.path.join(os.path.dirname(__file__), *baza_znaniy_files)
    opros = os.path.join(os.path.dirname(__file__),   *oprosnik_files)
    script_dir = os.path.abspath(os.path.dirname(__file__))
    adress_as = os.path.join(script_dir, 'AS')
    print()
    print('-------- Формирую отчет --------')
    print()
    warnings.simplefilter("ignore")
    ws.xlsx_to_dict(path=baza, select_sheet='База')


    ws_two.xlsx_to_dict(path=opros,  select_sheet='Надёжность', skiprows=3, data_only=True)

    def check_version(filepath):
        workbook = openpyxl.load_workbook(filepath)
        worksheet = workbook['Надёжность']
        first_row = worksheet[1]
        mas_ver_dtn = []
        for cell in first_row:
            if cell.value is not None:
                mas_ver_dtn.append(cell.value)
        # Получение ключа страховщика
        insurer_key = next(item for item in mas_ver_dtn if item not in ['ОПРОСНИК\nпо оценке зрелости ITSM-процессов', 'версия опросника:'])

        # Получение ключа версии опросника
        pattern = r'\d+\.\d+\.\d+\.\d+\.\d+'
        version_key = next(item for item in mas_ver_dtn if re.match(pattern, str(item)))
        # Создание нового словаря
        result = {
            'ОПРОСНИК\nпо оценке зрелости ITSM-процессов': insurer_key,
            'версия опросника:': version_key
        }

        expected_version = '6.1.'

        if expected_version not in result['версия опросника:']:
            print(f"Ожидаемая версия должна быть: {expected_version}, а в Опроснике используется версия: {result['версия опросника:']}. Работа скрипта прерывается.")
            exit()

    if __name__ == "__main__":
        filepath = opros
        check_version(filepath)

    workbook_name = openpyxl.load_workbook(opros)
    worksheet = workbook_name['Надёжность']
    first_row = worksheet[1]
    mas_ver_dtn = []
    for cell in first_row:
        if cell.value is not None:
            mas_ver_dtn.append(cell.value)
    # Получение ключа страховщика
    insurer_key = next(item for item in mas_ver_dtn if item not in ['ОПРОСНИК\nпо оценке зрелости ITSM-процессов', 'версия опросника:'])
    # Получение ключа версии опросника
    pattern = r'\d+\.\d+\.\d+\.\d+\.\d+'
    version_key = next(item for item in mas_ver_dtn if re.match(pattern, str(item)))
    # Создание нового словаря
    inn_index = mas_ver_dtn.index('ИНН:')
    result = {
        'ОПРОСНИК\nпо оценке зрелости ITSM-процессов': insurer_key,
        'версия опросника:': version_key,
        'ИНН:': mas_ver_dtn[inn_index + 1]
    }

    #Вытаскиваю критичность из Excel
    second_row = worksheet[3]
    mas_profile_dtn = []
    for cell in second_row:
        if cell.value is not None:
            mas_profile_dtn.append(cell.value)

    # Получение ключа страховщика
    insurer_key_profile = next((item for item in mas_profile_dtn if item not in ['Каким максимальным уровнем критичности обладает ИТ-сервис в компании?', 'ФИО ответственного:']), None)

    # Получение значения ФИО
    fio_dtn = mas_profile_dtn[-1]
    result_profile = {
        'Каким максимальным уровнем критичности обладает ИТ-сервис в компании?': insurer_key_profile,
        'ФИО ответственного:': fio_dtn
    }


    if result_profile['Каким максимальным уровнем критичности обладает ИТ-сервис в компании?'] not in ['Нет подобных ИТ-сервисов', 'OP', 'BO', 'BC и выше']:
        print(f'В файле не указана критичность ИТ-сервис в компании, скрипт работать не будет!!!')
        sys.exit()
    else:

        for k, v in dict_inn_company.items():
            if str(result['ИНН:']) == k:
                name_dtn = v[1]
                name_company = v[1] + '_' +result_profile['Каким максимальным уровнем критичности обладает ИТ-сервис в компании?'] + '_ЕРО.xlsx'
                name_coms = name_company.replace('"', '').replace('«', '').replace('»', '')  #Если в файле  Опросника название компании в кавычках «»
                save = os.path.join(os.path.dirname(__file__), name_coms)
                break

            else:
                name_dtn = result['ОПРОСНИК\nпо оценке зрелости ITSM-процессов']
                name_company = result['ОПРОСНИК\nпо оценке зрелости ITSM-процессов'] + '_' +result_profile['Каким максимальным уровнем критичности обладает ИТ-сервис в компании?'] + '_ЕРО.xlsx'
                name_coms = name_company.replace('"', '').replace('«', '').replace('»', '')  #Если в файле  Опросника название компании в кавычках «»
                save = os.path.join(os.path.dirname(__file__), name_coms)


        # name_company = result['ОПРОСНИК\nпо оценке зрелости ITSM-процессов'] + '_' +result_profile['Каким максимальным уровнем критичности обладает ИТ-сервис в компании?'] + '_ЕРО.xlsx'
        # name_coms = name_company.replace('"', '').replace('«', '').replace('»', '')  #Если в файле  Опросника название компании в кавычках «»
        # save = os.path.join(os.path.dirname(__file__), name_coms)



    #TODO  Начало обработки по анкетам АС(ИТ-технология) если такая папка существует вообщем создаю функцию (2-функции: my_function, find_and_open_excel )
    def my_function():
        as_dict = {}
        as_dict_norisk = {}
        book_znanii_as = {}
        script_dir = os.path.abspath(os.path.dirname(__file__))
        adress_as = os.path.join(script_dir, 'AS')
        ws_as.xlsx_to_dict(path=baza, select_sheet='База АС')
        if not os.path.isdir(adress_as):
            print(f'Папка {adress_as}  с анкетами по ИТ-услугам не найдена, продолжаю обработку по ITSM процессам  ')
            return
        else:
            for filename in os.listdir(adress_as):
                if filename.endswith(".xlsx"):
                    file_path = os.path.join(adress_as, filename)
                    try:
                        wb = openpyxl.load_workbook(file_path, data_only=True)
                        # print(f'Открываем файл - {filename}')
                        ws = wb['Титульный']
                        ws_new_as = wb['Технологии']
                        print(f'Работаем с файлом - {filename} ---> АС - {ws["C4"].value}')

                        if ws['C4'].value != 0:
                            as_dict[ws['C4'].value] = []
                            for row in range(2, ws_new_as.max_row + 1):
                                value_a = ws_new_as.cell(row, 1).value  # Столбец A
                                value_b = ws_new_as.cell(row, 3).value  # Столбец C
                                value_d = ws_new_as.cell(row, 6).value  # Столбец F

                                value_rekomend = ws_new_as.cell(row, 8).value

                                # value_h = ws_new_as.cell(row, 8).value  # Столбец H
                                value_e = ws_new_as.cell(row, 16).value  # Столбец P
                                value_n = ws_new_as.cell(row, 14).value  # Столбец N
                                value_ps = ws['C8'].value #Добавляем из Титульного листа "КОД ОБЪЕКТА ИССЛЕДОВАНИЯ"

                                values = [value_a, value_b, value_d, value_e, value_ps, value_n, value_rekomend] #value_h,
                                as_dict[ws['C4'].value].append(values)

                        else:
                            print(f'Не заполнено поле -- Название ИТ-услуги -- в файле {filename}  ')
                            break
                    except Exception as e:
                        print(f"Ошибка при открытии файла '{file_path}': {e}")


        for key, value in as_dict.items():
            as_dict_norisk[key] = []
            for sublist in value:
                if not any(x in sublist for x in ['Нет риска', 'Нет отклонения']):
                    if 'Неприменимо' not in sublist:
                        if sublist[2] and sublist[3] != None:
                            # if None not in sublist:
                            as_dict_norisk[key].append(sublist)


        for f in ws_as.sheet_items:

            juice_znaniya_as = f.setdefault('Номер вопроса')
            juice_znaniya_answer_as = f.setdefault('Ответ')
            juice_znaniya_otklon_as = f.setdefault('Риск')
            juice_znaniya_recomend_as = f.setdefault('Рекомендация')
            # juice_critical_as = f.setdefault('Уровень критичности отклонения')

            #TODO: Добавление в словарь из базы знаний
            book_znanii_as.setdefault(juice_znaniya_as, [])
            book_znanii_as[juice_znaniya_as].append(juice_znaniya_answer_as)
            book_znanii_as[juice_znaniya_as].append(juice_znaniya_otklon_as)
            book_znanii_as[juice_znaniya_as].append(juice_znaniya_recomend_as)
            # book_znanii_as[juice_znaniya_as].append(juice_critical_as)

        clean_book_as = {key: value for key, value in book_znanii_as.items() if key != 'x'}

        dict_3 = {}
        dict_tire = {} # Добавляем в этот словарь все значения убирая "-"(тире из Базы знаний)

        for key_2, values_2 in as_dict_norisk.items():
            for sublist_2 in values_2:

                key_1 = str(sublist_2[0])
                value_2 = str(sublist_2[1])
                value_3 = sublist_2[1]
                if key_1 in clean_book_as:
                    for i, item in enumerate(clean_book_as[key_1]):
                        if value_2 in item:
                            description = clean_book_as[key_1][i + 1] if i + 1 < len(clean_book_as[key_1]) else ""  # Проверка границы
                            recommendation = clean_book_as[key_1][i + 2] if i + 2 < len(clean_book_as[key_1]) else "" # Проверка границы

                            if sublist_2[-1] == None:
                                dict_3.setdefault(key_2, []).append([key_1, value_2, description, recommendation, sublist_2[2], sublist_2[3], sublist_2[4], sublist_2[5]])
                            else:
                                dict_3.setdefault(key_2, []).append([key_1, value_2, description, sublist_2[-1], sublist_2[2], sublist_2[3], sublist_2[4], sublist_2[5]])

                        if (isinstance(value_3, (int, float)) and  (isinstance(item, str))):
                            description = clean_book_as[key_1][i + 1] if i + 1 < len(clean_book_as[key_1]) else ""  # Проверка границы
                            recommendation = clean_book_as[key_1][i + 2] if i + 2 < len(clean_book_as[key_1]) else "" # Проверка границы

                            # dict_3.setdefault(key_2, []).append([key_1, value_2, description, recommendation, sublist_2[2], sublist_2[3], sublist_2[4], sublist_2[5]])
                            if sublist_2[-1] == None:
                                dict_3.setdefault(key_2, []).append([key_1, value_2, description, recommendation, sublist_2[2], sublist_2[3], sublist_2[4], sublist_2[5]])
                            else:
                                dict_3.setdefault(key_2, []).append([key_1, value_2, description, sublist_2[-1], sublist_2[2], sublist_2[3], sublist_2[4], sublist_2[5]])
                            break

        for key, value in dict_3.items():
            new_value = []
            for sublist in value:
                new_sublist = []
                for item in sublist:
                    if isinstance(item, str) and ' — ' in item:
                        parts = item.split(' — ')
                        target_part = parts[1].capitalize()  # Сохранить нужную часть
                        new_sublist.extend(parts)
                        new_sublist[new_sublist.index(parts[1])] = target_part  # Заменить
                    else:
                        new_sublist.append(item)
                new_value.append(new_sublist)
            dict_tire[key] = new_value
        #Меняем знаяения букв в словаре A, B, C
        dict_kritik = {'A':'Критичный', 'B':'Высокий', 'C':'Умеренный'}

        for key, value in dict_tire.items():
            for sublist in value:
                for i, item in enumerate(sublist):
                    if item in dict_kritik:
                        sublist[i] = dict_kritik[item]

        return dict_tire
    if __name__ == "__main__":
        my_function()


    def find_and_open_excel(adress_as):
        as_dict = {}
        if not os.path.isdir(adress_as):
            print(f"Ошибка: папка '{adress_as}' не найдена.")
            return
        else:
            mas_anketa = []
            for filename in os.listdir(adress_as):
                if filename.endswith(".xlsx"):
                    mas_anketa.append(filename)

            for anketa in mas_anketa:
                file_path = os.path.join(adress_as, anketa)
                # file_path = os.path.join(adress_as, (k for k in mas_anketa))
                try:
                    wb = openpyxl.load_workbook(file_path, data_only=True)
                    # print(f'Открываем файл - {filename}')
                    ws = wb['Титульный']
                    print(f'Работаем с файлом - {anketa} ---> АС - {ws["C4"].value}')

                    if ws['C4'].value != 0:
                        try:
                            as_dict[ws['C4'].value] = [ws['C5'].value, round(ws['D21'].value, 2), round(ws['F21'].value, 2)]
                        except:
                            as_dict[ws['C4'].value] = [ws['C5'].value, round(ws['D21'].value, 2), 0]

                    else:
                        print(f'Не заполненно поле -- Название ИТ-услуги -- в файле {anketa}  ')
                        sys.exit()
                except Exception as e:
                    print(f"Ошибка при открытии файла '{file_path}': {e}")
            # else:
            #     print(f'В папке {adress_as} отсутствуют анкеты по ИТ-услугам, продолжаю обработку по ITSM процессам  ')
            #     return

        return as_dict
    rezult_as = find_and_open_excel(adress_as)


    # print(rezult_as)
    # print(my_function())


    # def my_vnd_as():
    #     as_dict = {}
    #     as_dict_norisk = {}
    #     book_znanii_as_op = {}
    #     book_znanii_as_bo = {}
    #     book_znanii_as_bc = {}
    #     book_znanii_as_mc = {}
    #     counts_as = {}
    #     script_dir = os.path.abspath(os.path.dirname(__file__))
    #     adress_as = os.path.join(script_dir, 'AS')
    #     ws_as.xlsx_to_dict(path=baza, select_sheet='База АС')
    #     if not os.path.isdir(adress_as):
    #         print(f'Папка {adress_as}  с анкетами по ИТ-услугам не найдена, продолжаю обработку по ITSM процессам  ')
    #         return
    #     else:
    #         for filename in os.listdir(adress_as):
    #             if filename.endswith(".xlsx"):
    #                 file_path = os.path.join(adress_as, filename)
    #                 try:
    #                     wb = openpyxl.load_workbook(file_path, data_only=True)
    #                     # print(f'Открываем файл - {filename}')
    #                     ws = wb['Титульный']
    #                     ws_new_as = wb['Технологии']
    #
    #
    #                     if ws['C4'].value != 0:
    #                         as_dict[ws['C4'].value] = []
    #                         for row in range(2, ws_new_as.max_row + 1):
    #                             value_a = ws_new_as.cell(row, 1).value  # Столбец A
    #                             # value_b = ws_new_as.cell(row, 3).value  # Столбец C
    #                             value_d = ws_new_as.cell(row, 6).value  # Столбец F
    #                             # value_h = ws_new_as.cell(row, 8).value  # Столбец H
    #                             # value_e = ws_new_as.cell(row, 16).value  # Столбец P
    #                             # value_n = ws_new_as.cell(row, 14).value  # Столбец N
    #                             value_ps = ws['C5'].value #Добавляем из Титульного листа "КОД ОБЪЕКТА ИССЛЕДОВАНИЯ"
    #                             values = [value_ps, value_a, value_d]
    #                             as_dict[ws['C4'].value].append(values)
    #
    #                     else:
    #                         print(f'Не заполнено поле -- Название ИТ-услуги -- в файле {filename}  ')
    #                         break
    #                 except Exception as e:
    #                     print(f"Ошибка при открытии файла '{file_path}': {e}")
    #
    #
    #
    #     # print(as_dict)
    #     for k, v in as_dict.items():
    #         # print(k, v[0][0])
    #         if v[0][0] == 'OP':
    #             for f in ws_as.sheet_items:
    #                 juice_znaniya_as = f.setdefault('Коды для уровней отклонения')
    #                 juice_znaniya_answer_as = f.setdefault('OP')
    #                 juice_znaniya_otklon_as = f.setdefault('Уровень готовности надёжности OP')
    #
    #                 # Инициализация списка для ключа, если его еще нет
    #                 if k not in book_znanii_as_op:
    #                     book_znanii_as_op[k] = []
    #                 # Добавление нового подсписка с тремя элементами
    #                 book_znanii_as_op[k].append([
    #                     juice_znaniya_as,
    #                     juice_znaniya_answer_as,
    #                     juice_znaniya_otklon_as
    #                 ])
    #         if v[0][0] == 'BO':
    #             for f in ws_as.sheet_items:
    #                 juice_znaniya_as = f.setdefault('Коды для уровней отклонения')
    #                 juice_znaniya_answer_as = f.setdefault('BO')
    #                 juice_znaniya_otklon_as = f.setdefault('Уровень готовности надёжности BO')
    #
    #                 # Инициализация списка для ключа, если его еще нет
    #                 if k not in book_znanii_as_bo:
    #                     book_znanii_as_bo[k] = []
    #
    #                 # Добавление нового подсписка с тремя элементами
    #                 book_znanii_as_bo[k].append([
    #                     juice_znaniya_as,
    #                     juice_znaniya_answer_as,
    #                     juice_znaniya_otklon_as
    #                 ])
    #
    #
    #
    #     # Удаление пустых значений (если нужно)
    #     for k in book_znanii_as_op:
    #         book_znanii_as_op[k] = [item for item in book_znanii_as_op[k] if all(item)]
    #     # Инициализация словаря для хранения результатов
    #     result_dict = {}
    #     # Проходим по каждому ключу в book_znanii_as_op
    #     for key, value_list in book_znanii_as_op.items():
    #         # Инициализация счетчиков для текущего ключа
    #         result_dict[key] = {
    #             'Критичный на 3:': [0, 0],  # [общее количество, выполнено]
    #             'Высокий на 3:': [0, 0],
    #             'Умеренный на 3:': [0, 0],
    #             'Критичный на 4:': [0, 0],
    #             'Высокий на 4:': [0, 0],
    #             'Умеренный на 4:': [0, 0]
    #         }
    #
    #         # Подсчет общего количества для каждого уровня и категории
    #         for item in value_list:
    #             code, category, level = item
    #             result_key = f"{category} на {level}:"
    #             if result_key in result_dict[key]:
    #                 result_dict[key][result_key][0] += 1  # Увеличиваем общее количество
    #
    #         # Подсчет выполненных элементов из as_dict
    #         if key in as_dict:
    #             for as_item in as_dict[key]:
    #                 if len(as_item) >= 3 and as_item[2] == 'Нет риска':
    #                     code = as_item[1]
    #                     # Ищем код в book_znanii_as_op
    #                     for book_item in value_list:
    #                         if str(book_item[0]) == str(code):
    #                             category = book_item[1]
    #                             level = book_item[2]
    #                             result_key = f"{category} на {level}:"
    #                             if result_key in result_dict[key]:
    #                                 result_dict[key][result_key][1] += 1  # Увеличиваем количество выполненных
    #
    #     # Вывод результата
    #     for key, value in result_dict.items():
    #         print(f"{key}: {value}")
    #
    #
    #     for key, values in as_dict.items():
    #         for i, v in book_znanii_as_op.items():
    #             if i == key:
    #                 # counts_as[key] = {'Критичный': 0, 'Высокий': 0, 'Умеренный': 0}
    #
    #                 pass
    #
    #
    #     return as_dict
    # if __name__ == "__main__":
    #     my_vnd_as()



    # print(my_function())
    counts_as = {}
    for key, values in my_function().items():
        for value in values:
            risk_level = value[-4]
            # print(risk_level)
            if key not in counts_as:
                counts_as[key] = {'Критичный': 0, 'Высокий': 0, 'Умеренный': 0}
            counts_as[key][risk_level] += 1
    # Словарь data_new выводит {'Address Processing': {'Критичный': 1, 'Высокий': 74, 'Умеренный': 10}}
    # общее количество отклонений по требованию где стоит нет или нет ответа(None)!
    data_new_count = {}
    for key in counts_as:
        data_new_count[key] = [rezult_as[key][0]]
        data_new_count[key].extend(list(counts_as[key].values()))


    #Функция нужна для подсчета всего требований по АС - это нужно для расчета Индекса и готовности 3-го и 4-го уровня надежности

    def func_as_artefakt(adress_as):
        new_result = {}
        new_result_result_treb = {}

        if not os.path.isdir(adress_as):
            print(f"Ошибка: папка '{adress_as}' не найдена.")
            return

        for filename in os.listdir(adress_as):
            if filename.endswith(".xlsx"):
                file_path = os.path.join(adress_as, filename)
                try:
                    wb = openpyxl.load_workbook(file_path, data_only=True)
                    ws = wb['Технологии']
                    wm = wb['Титульный']
                    new_result[wm['C4'].value] = []
                    new_result_result_treb[wm['C4'].value] = []
                    # Определяем диапазоны ячеек
                    range_1 = 'E10:E18'
                    range_3 = 'G10:G18'
                    range_2 = 'F10:F18'
                    range_4 = 'H10:H18'


                    # Функция для получения списка значений из диапазона
                    def get_values_from_range(wm, range_str):
                        values = []
                        for row in wm[range_str]:
                            for cell in row:
                                # Обрабатываем три случая: '-', None и другие значения
                                if cell.value == '-' or cell.value is None:
                                    values.append(0)
                                else:
                                    values.append(cell.value)
                        return values
                    # Получаем списки значений для обоих диапазонов
                    values_1 = get_values_from_range(wm, range_1)
                    values_2 = get_values_from_range(wm, range_3)
                    values_3 = get_values_from_range(wm, range_2)
                    values_4 = get_values_from_range(wm, range_4)
                    # Суммы значений для диапазонов
                    total_sum_1 = sum(values_1)
                    total_sum_2 = sum(values_2)
                    total_sum_3 = sum(values_3)
                    total_sum_4 = sum(values_4)

                    # Добавляем суммарные значения в результат
                    new_result_result_treb[wm['C4'].value].append((
                        total_sum_1, total_sum_3, total_sum_2, total_sum_4
                    ))

                except Exception as e:
                    print(f"Произошла ошибка при обработке файла {filename}: {e}")

        return new_result_result_treb
    all_treb_as = func_as_artefakt(adress_as)

    # print(all_treb_as) ## {'Address Processing': [(38.0, 2, 11.0, 1.5)]}
    # print(data_new_count) ##{'Address Processing': ['BO', 0, 74, 10]}

 #------------------------------------------------------------------------------------------------------------------


    #Словари для Опросника и Знаний
    book_opros = {}
    book_znanii = {}

    dict_itsm = {'ITSM-INC': 'Управление технологическими инцидентами',
                 'ITSM-CATM': 'Управление каталогом сервисов',
                 'ITSM-SLM': 'Управление уровнями сервисов',
                 'ITSM-PRB': 'Управление технологическими проблемами',
                 'ITSM-ERM': 'Управление технологическими рисками',
                 'ITSM-ICHG': 'Управление инфраструктурными изменениями',
                 'ITSM-CONT': 'Управление непрерывностью технологий',
                 'ITSM-AVL': 'Управление доступностью технологий',
                 'ITSM-M&E': 'Управление технологическими событиями и мониторингом',
                 'ITSM-RFF': 'Управление технологическими запросами на обслуживание',
                 'ITSM-CFG': 'Управление конфигурациями',
                 'ITSM-UPD': 'Управление обновлением ПО',
                 'ITSM-SD': 'Управление технологическими обращениями',
                 'ITSM-CAP': 'Управление мощностями технологий',
                 'ITSM-MNTW': 'Управление регламентными работами'
                 }

    dict_itsm_dtn = {'ITSM-CFG': '5492.CFG.',
                     'ITSM-CATM': '5492.CATM.',
                     'ITSM-SLM': '5492.SLM.',
                     'ITSM-INC': '5492.INC.',
                     'ITSM-PRB': '5492.PRB.',
                     'ITSM-ERM': '5492.ERM.',
                     'ITSM-RFF': '5492.RFF.',
                     'ITSM-ICHG': '5492.ICHG.',
                     'ITSM-UPD': '5492.UPD.',
                     'ITSM-AVL': '5492.AVL.',
                     'ITSM-CONT': '5492.CONT.',
                     'ITSM-M&E': '5492.M&E.',
                     'ITSM-CAP': '5492.CAP.',
                     'ITSM-SD': '5492.SD.',
                     'ITSM-MNTW': '5492.MNTW.'

                     }
    #TODO: Собираем информацию с Опросника согласно полям!
    for b in ws_two.sheet_items:

        juice_id = b.setdefault('ID\nвопроса')

        juice_urov_ur_got = b.setdefault('Уровень готовности в зависимости от критичности АС')
        juice_ur_otklon = b.setdefault('Уровень отклонения')
        juice_answer = b.setdefault('Ответ с учётом артефакта \n(авто)')

        book_opros.setdefault(juice_id, [])
        book_opros[juice_id].append(juice_urov_ur_got)
        book_opros[juice_id].append(juice_ur_otklon)
        book_opros[juice_id].append(juice_answer)
    removed_value = book_opros.pop('-')  # Удаляем ключ '-'

    # ----- При парсинге Опросника если ключь ID Вопроса с ITSM то он мапит с добавлением префикса DTN- нужно для ЕРО в ЦТК-----
    # Создаём новый словарь с нужными ключами
    updated_book_opros = {}
    for key, value in book_opros.items():
        if not key.startswith('DTN-'):
            new_key = 'DTN-' + key
        else:
            new_key = key
        updated_book_opros[new_key] = value

    # Перезаписываем исходный словарь (если нужно)
    book_opros = updated_book_opros

    #TODO: Собираем информацию с файла Базы знаний согласно полям!
    for f in ws.sheet_items:

        juice_znaniya = f.setdefault('Код вопроса')
        juice_znaniya_otklon = f.setdefault('Выявленное отклонение')
        juice_znaniya_posledstvie = f.setdefault('Негативные последствия от проявления отклонения')
        juice_znaniya_recomend = f.setdefault('Рекомендации по устранению отклонения')
        juice_artefakt = f.setdefault('Артефакт')

        book_znanii.setdefault(juice_znaniya, [])
        book_znanii[juice_znaniya].append([juice_znaniya_otklon])
        book_znanii[juice_znaniya].append([juice_znaniya_posledstvie])
        book_znanii[juice_znaniya].append([juice_znaniya_recomend + ' Aртeфaкт: ' + juice_artefakt])



    # Создаем defaultdict, где значения по умолчанию — пустые словари
    result_dicts = defaultdict(dict)
    #Расчет по отклонениям в рамках каждого вопроса
    prefixes = [
        'ITSM-INC', 'ITSM-CATM', 'ITSM-CFG', 'ITSM-SLM', 'ITSM-PRB', 'ITSM-ERM',
        'ITSM-RFF', 'ITSM-ICHG', 'ITSM-UPD', 'ITSM-AVL', 'ITSM-M&E', 'ITSM-CONT',
        'ITSM-CAP', 'ITSM-SD', 'ITSM-MNTW']


    for key, value in book_opros.items():
        if value[0] != 'ЦТК':
            for prefix in prefixes:
                if prefix in key:
                    result_dicts[prefix][key] = value
                    break

    # Теперь у нас есть словари для каждого префикса
    dict_inc = result_dicts['ITSM-INC']
    dict_catm = result_dicts['ITSM-CATM']
    dict_cfg = result_dicts['ITSM-CFG']
    dict_slm = result_dicts['ITSM-SLM']
    dict_prb = result_dicts['ITSM-PRB']
    dict_erm = result_dicts['ITSM-ERM']
    dict_rff = result_dicts['ITSM-RFF']
    dict_ichg = result_dicts['ITSM-ICHG']
    dict_upd = result_dicts['ITSM-UPD']
    dict_avl = result_dicts['ITSM-AVL']
    dict_me = result_dicts['ITSM-M&E']
    dict_cont = result_dicts['ITSM-CONT']
    dict_cap = result_dicts['ITSM-CAP']
    dict_mntw = result_dicts['ITSM-MNTW']
    dict_sd = result_dicts['ITSM-SD']


    # dict_tree_four_assesment_itsm = {
    #     'Всего критичных': 0,
    #     'Выполнено критичных': 0,
    #     'Всего высоких': 0,
    #     'Выполнено высоких': 0,
    #     'Всего умеренных': 0,
    #     'Выполнено умеренных': 0,
    # }
    #
    # for i, d in result_dicts.items():
    #     for k, v in d.items():
    #         if '3' in v and 'Критичный' in v:
    #             dict_tree_four_assesment_itsm['Всего критичных'] += 1
    #         if '3' in v and 'Критичный' in v and 'да' in v:
    #             dict_tree_four_assesment_itsm['Выполнено критичных'] += 1
    #         if '3' in v and 'Высокий' in v:
    #             dict_tree_four_assesment_itsm['Всего высоких'] += 1
    #         if '3' in v and 'Высокий' in v and 'да' in v:
    #             dict_tree_four_assesment_itsm['Выполнено высоких'] += 1
    #         if '3' in v and 'Умеренный' in v:
    #             dict_tree_four_assesment_itsm['Всего умеренных'] += 1
    #         if '3' in v and 'Умеренный' in v and 'да' in v:
    #             dict_tree_four_assesment_itsm['Выполнено умеренных'] += 1
            # if '4' in v and 'Критичный' in v:
            #     dict_tree_four_assesment_itsm['Всего критичных'] += 1
            # if '4' in v and 'Критичный' in v and 'да' in v:
            #     dict_tree_four_assesment_itsm['Выполнено критичных'] += 1
            # if '4' in v and 'Высокий' in v:
            #     dict_tree_four_assesment_itsm['Всего высоких'] += 1
            # if '4' in v and 'Высокий' in v and 'да' in v:
            #     dict_tree_four_assesment_itsm['Выполнено высоких'] += 1
            # if '4' in v and 'Умеренный' in v:
            #     dict_tree_four_assesment_itsm['Всего умеренных'] += 1
            # if '4' in v and 'Умеренный' in v and 'да' in v:
            #     dict_tree_four_assesment_itsm['Выполнено умеренных'] += 1

    #Собираем массив с теми отклонениями которые получили ИТОГОВЫЙ ответ "НЕТ" нужны для формирования информации из базы знаний
    dict_all_answer_no = []

    dict_inc_vnd = {'ITSM-INC': {'Всего требований': 0,
                                 'Всего требований на 3': 0,
                                 'Всего выполнено требований на 3': 0,
                                 'Всего требований на 4': 0,
                                 'Всего выполнено требований на 4': 0}}

    dict_catm_vnd = {'ITSM-CATM': {'Всего требований': 0,
                                   'Всего требований на 3': 0,
                                   'Всего выполнено требований на 3': 0,
                                   'Всего требований на 4': 0,
                                   'Всего выполнено требований на 4': 0}}

    dict_cfg_vnd = {'ITSM-CFG': {'Всего требований': 0,
                                 'Всего требований на 3': 0,
                                 'Всего выполнено требований на 3': 0,
                                 'Всего требований на 4': 0,
                                 'Всего выполнено требований на 4': 0}}


    dict_slm_vnd = {'ITSM-SLM': {'Всего требований': 0,
                                 'Всего требований на 3': 0,
                                 'Всего выполнено требований на 3': 0,
                                 'Всего требований на 4': 0,
                                 'Всего выполнено требований на 4': 0}}

    dict_prb_vnd = {'ITSM-PRB': {'Всего требований': 0,
                                 'Всего требований на 3': 0,
                                 'Всего выполнено требований на 3': 0,
                                 'Всего требований на 4': 0,
                                 'Всего выполнено требований на 4': 0}}


    dict_erm_vnd = {'ITSM-ERM': {'Всего требований': 0,
                                 'Всего требований на 3': 0,
                                 'Всего выполнено требований на 3': 0,
                                 'Всего требований на 4': 0,
                                 'Всего выполнено требований на 4': 0}}


    dict_rff_vnd = {'ITSM-RFF': {'Всего требований': 0,
                                 'Всего требований на 3': 0,
                                 'Всего выполнено требований на 3': 0,
                                 'Всего требований на 4': 0,
                                 'Всего выполнено требований на 4': 0}}

    dict_ichg_vnd = {'ITSM-ICHG': {'Всего требований': 0,
                                   'Всего требований на 3': 0,
                                   'Всего выполнено требований на 3': 0,
                                   'Всего требований на 4': 0,
                                   'Всего выполнено требований на 4': 0}}

    dict_upd_vnd = {'ITSM-UPD': {'Всего требований': 0,
                                 'Всего требований на 3': 0,
                                 'Всего выполнено требований на 3': 0,
                                 'Всего требований на 4': 0,
                                 'Всего выполнено требований на 4': 0}}


    dict_avl_vnd = {'ITSM-AVL': {'Всего требований': 0,
                                 'Всего требований на 3': 0,
                                 'Всего выполнено требований на 3': 0,
                                 'Всего требований на 4': 0,
                                 'Всего выполнено требований на 4': 0}}


    dict_me_vnd = {'ITSM-M&E': {'Всего требований': 0,
                               'Всего требований на 3': 0,
                               'Всего выполнено требований на 3': 0,
                               'Всего требований на 4': 0,
                               'Всего выполнено требований на 4': 0}}

    dict_cont_vnd = {'ITSM-CONT': {'Всего требований': 0,
                                   'Всего требований на 3': 0,
                                   'Всего выполнено требований на 3': 0,
                                   'Всего требований на 4': 0,
                                   'Всего выполнено требований на 4': 0}}

    dict_cap_vnd = {'ITSM-CAP': {'Всего требований': 0,
                                 'Всего требований на 3': 0,
                                 'Всего выполнено требований на 3': 0,
                                 'Всего требований на 4': 0,
                                 'Всего выполнено требований на 4': 0}}

    dict_sd_vnd = {'ITSM-SD': {'Всего требований': 0,
                               'Всего требований на 3': 0,
                               'Всего выполнено требований на 3': 0,
                               'Всего требований на 4': 0,
                               'Всего выполнено требований на 4': 0}}

    dict_mntw_vnd = {'ITSM-MNTW': {'Всего требований': 0,
                                   'Всего требований на 3': 0,
                                   'Всего выполнено требований на 3': 0,
                                   'Всего требований на 4': 0,
                                   'Всего выполнено требований на 4': 0}}


    # ********************************************** ITSM-INC ********************************************************************
    # Расчет требований по процессу inc для 3-го уровня готовности
    for k, v in dict_inc.items():
        if v[2] == 'нет':
            dict_all_answer_no.append(k)
        if v[0] == '3':
            dict_inc_vnd['ITSM-INC']['Всего требований на 3'] += 1
        if v[0] == '4':
            dict_inc_vnd['ITSM-INC']['Всего требований на 4'] += 1
        if v[2] == 'да' and v[0] == '3':
            dict_inc_vnd['ITSM-INC']['Всего выполнено требований на 3'] += 1
        if v[2] == 'да' and v[0] == '4':
            dict_inc_vnd['ITSM-INC']['Всего выполнено требований на 4'] += 1
    dict_inc_vnd['ITSM-INC']['Всего требований'] = dict_inc_vnd['ITSM-INC']['Всего требований на 3'] + dict_inc_vnd['ITSM-INC']['Всего требований на 4']


    # ********************************************** ITSM-CATM ********************************************************************
    # Расчет требований по процессу CATM для 3-го уровня готовности
    for k, v in dict_catm.items():
        if v[2] == 'нет':
            dict_all_answer_no.append(k)
        if v[0] == '3':
            dict_catm_vnd['ITSM-CATM']['Всего требований на 3'] += 1
        if v[0] == '4':
            dict_catm_vnd['ITSM-CATM']['Всего требований на 4'] += 1
        if v[2] == 'да' and v[0] == '3':
            dict_catm_vnd['ITSM-CATM']['Всего выполнено требований на 3'] += 1
        if v[2] == 'да' and v[0] == '4':
            dict_catm_vnd['ITSM-CATM']['Всего выполнено требований на 4'] += 1
    dict_catm_vnd['ITSM-CATM']['Всего требований'] = dict_catm_vnd['ITSM-CATM']['Всего требований на 3'] + dict_catm_vnd['ITSM-CATM']['Всего требований на 4']


    # ********************************************** ITSM-CFG ********************************************************************
    # Расчет требований по процессу inc для 3-го уровня готовности
    for k, v in dict_cfg.items():
        if v[2] == 'нет':
            dict_all_answer_no.append(k)
        if v[0] == '3':
            dict_cfg_vnd['ITSM-CFG']['Всего требований на 3'] += 1
        if v[0] == '4':
            dict_cfg_vnd['ITSM-CFG']['Всего требований на 4'] += 1
        if v[2] == 'да' and v[0] == '3':
            dict_cfg_vnd['ITSM-CFG']['Всего выполнено требований на 3'] += 1
        if v[2] == 'да' and v[0] == '4':
            dict_cfg_vnd['ITSM-CFG']['Всего выполнено требований на 4'] += 1
    dict_cfg_vnd['ITSM-CFG']['Всего требований'] = dict_cfg_vnd['ITSM-CFG']['Всего требований на 3'] + dict_cfg_vnd['ITSM-CFG']['Всего требований на 4']


    # ********************************************** ITSM-SLM ********************************************************************
    # Расчет требований по процессу inc для 3-го уровня готовности
    for k, v in dict_slm.items():
        if v[2] == 'нет':
            dict_all_answer_no.append(k)
        if v[0] == '3':
            dict_slm_vnd['ITSM-SLM']['Всего требований на 3'] += 1
        if v[0] == '4':
            dict_slm_vnd['ITSM-SLM']['Всего требований на 4'] += 1
        if v[2] == 'да' and v[0] == '3':
            dict_slm_vnd['ITSM-SLM']['Всего выполнено требований на 3'] += 1
        if v[2] == 'да' and v[0] == '4':
            dict_slm_vnd['ITSM-SLM']['Всего выполнено требований на 4'] += 1
    dict_slm_vnd['ITSM-SLM']['Всего требований'] = dict_slm_vnd['ITSM-SLM']['Всего требований на 3'] + dict_slm_vnd['ITSM-SLM']['Всего требований на 4']


    # ********************************************** ITSM-PRB ********************************************************************
    # Расчет требований по процессу inc для 3-го уровня готовности
    for k, v in dict_prb.items():
        if v[2] == 'нет':
            dict_all_answer_no.append(k)
        if v[0] == '3':
            dict_prb_vnd['ITSM-PRB']['Всего требований на 3'] += 1
        if v[0] == '4':
            dict_prb_vnd['ITSM-PRB']['Всего требований на 4'] += 1
        if v[2] == 'да' and v[0] == '3':
            dict_prb_vnd['ITSM-PRB']['Всего выполнено требований на 3'] += 1
        if v[2] == 'да' and v[0] == '4':
            dict_prb_vnd['ITSM-PRB']['Всего выполнено требований на 4'] += 1
    dict_prb_vnd['ITSM-PRB']['Всего требований'] = dict_prb_vnd['ITSM-PRB']['Всего требований на 3'] + dict_prb_vnd['ITSM-PRB']['Всего требований на 4']


    # ********************************************** ITSM-ERM ********************************************************************
    # Расчет требований по процессу inc для 3-го уровня готовности
    for k, v in dict_erm.items():
        if v[2] == 'нет':
            dict_all_answer_no.append(k)
        if v[0] == '3':
            dict_erm_vnd['ITSM-ERM']['Всего требований на 3'] += 1
        if v[0] == '4':
            dict_erm_vnd['ITSM-ERM']['Всего требований на 4'] += 1
        if v[2] == 'да' and v[0] == '3':
            dict_erm_vnd['ITSM-ERM']['Всего выполнено требований на 3'] += 1
        if v[2] == 'да' and v[0] == '4':
            dict_erm_vnd['ITSM-ERM']['Всего выполнено требований на 4'] += 1
    dict_erm_vnd['ITSM-ERM']['Всего требований'] = dict_erm_vnd['ITSM-ERM']['Всего требований на 3'] + dict_erm_vnd['ITSM-ERM']['Всего требований на 4']


    # ********************************************** ITSM-RFF ********************************************************************
    # Расчет требований по процессу inc для 3-го уровня готовности
    for k, v in dict_rff.items():
        if v[2] == 'нет':
            dict_all_answer_no.append(k)
        if v[0] == '3':
            dict_rff_vnd['ITSM-RFF']['Всего требований на 3'] += 1
        if v[0] == '4':
            dict_rff_vnd['ITSM-RFF']['Всего требований на 4'] += 1
        if v[2] == 'да' and v[0] == '3':
            dict_rff_vnd['ITSM-RFF']['Всего выполнено требований на 3'] += 1
        if v[2] == 'да' and v[0] == '4':
            dict_rff_vnd['ITSM-RFF']['Всего выполнено требований на 4'] += 1
    dict_rff_vnd['ITSM-RFF']['Всего требований'] = dict_rff_vnd['ITSM-RFF']['Всего требований на 3'] + dict_rff_vnd['ITSM-RFF']['Всего требований на 4']


    # ********************************************** ITSM-ICHG ********************************************************************
    # Расчет требований по процессу inc для 3-го уровня готовности
    for k, v in dict_ichg.items():
        if v[2] == 'нет':
            dict_all_answer_no.append(k)
        if v[0] == '3':
            dict_ichg_vnd['ITSM-ICHG']['Всего требований на 3'] += 1
        if v[0] == '4':
            dict_ichg_vnd['ITSM-ICHG']['Всего требований на 4'] += 1
        if v[2] == 'да' and v[0] == '3':
            dict_ichg_vnd['ITSM-ICHG']['Всего выполнено требований на 3'] += 1
        if v[2] == 'да' and v[0] == '4':
            dict_ichg_vnd['ITSM-ICHG']['Всего выполнено требований на 4'] += 1
    dict_ichg_vnd['ITSM-ICHG']['Всего требований'] = dict_ichg_vnd['ITSM-ICHG']['Всего требований на 3'] + dict_ichg_vnd['ITSM-ICHG']['Всего требований на 4']


    # ********************************************** ITSM-UPD ********************************************************************
    # Расчет требований по процессу inc для 3-го уровня готовности
    for k, v in dict_upd.items():
        if v[2] == 'нет':
            dict_all_answer_no.append(k)
        if v[0] == '3':
            dict_upd_vnd['ITSM-UPD']['Всего требований на 3'] += 1
        if v[0] == '4':
            dict_upd_vnd['ITSM-UPD']['Всего требований на 4'] += 1
        if v[2] == 'да' and v[0] == '3':
            dict_upd_vnd['ITSM-UPD']['Всего выполнено требований на 3'] += 1
        if v[2] == 'да' and v[0] == '4':
            dict_upd_vnd['ITSM-UPD']['Всего выполнено требований на 4'] += 1
    dict_upd_vnd['ITSM-UPD']['Всего требований'] = dict_upd_vnd['ITSM-UPD']['Всего требований на 3'] + dict_upd_vnd['ITSM-UPD']['Всего требований на 4']


    # ********************************************** ITSM-AVL ********************************************************************
    # Расчет требований по процессу inc для 3-го уровня готовности
    for k, v in dict_avl.items():
        if v[2] == 'нет':
            dict_all_answer_no.append(k)
        if v[0] == '3':
            dict_avl_vnd['ITSM-AVL']['Всего требований на 3'] += 1
        if v[0] == '4':
            dict_avl_vnd['ITSM-AVL']['Всего требований на 4'] += 1
        if v[2] == 'да' and v[0] == '3':
            dict_avl_vnd['ITSM-AVL']['Всего выполнено требований на 3'] += 1
        if v[2] == 'да' and v[0] == '4':
            dict_avl_vnd['ITSM-AVL']['Всего выполнено требований на 4'] += 1
    dict_avl_vnd['ITSM-AVL']['Всего требований'] = dict_avl_vnd['ITSM-AVL']['Всего требований на 3'] + dict_avl_vnd['ITSM-AVL']['Всего требований на 4']


    # ********************************************** ITSM-M&E ********************************************************************
    # Расчет требований по процессу inc для 3-го уровня готовности
    for k, v in dict_me.items():
        if v[2] == 'нет':
            dict_all_answer_no.append(k)
        if v[0] == '3':
            dict_me_vnd['ITSM-M&E']['Всего требований на 3'] += 1
        if v[0] == '4':
            dict_me_vnd['ITSM-M&E']['Всего требований на 4'] += 1
        if v[2] == 'да' and v[0] == '3':
            dict_me_vnd['ITSM-M&E']['Всего выполнено требований на 3'] += 1
        if v[2] == 'да' and v[0] == '4':
            dict_me_vnd['ITSM-M&E']['Всего выполнено требований на 4'] += 1
    dict_me_vnd['ITSM-M&E']['Всего требований'] = dict_me_vnd['ITSM-M&E']['Всего требований на 3'] + dict_me_vnd['ITSM-M&E']['Всего требований на 4']


    # ********************************************** ITSM-CONT ********************************************************************
    # Расчет требований по процессу inc для 3-го уровня готовности
    for k, v in dict_cont.items():
        if v[2] == 'нет':
            dict_all_answer_no.append(k)
        if v[0] == '3':
            dict_cont_vnd['ITSM-CONT']['Всего требований на 3'] += 1
        if v[0] == '4':
            dict_cont_vnd['ITSM-CONT']['Всего требований на 4'] += 1
        if v[2] == 'да' and v[0] == '3':
            dict_cont_vnd['ITSM-CONT']['Всего выполнено требований на 3'] += 1
        if v[2] == 'да' and v[0] == '4':
            dict_cont_vnd['ITSM-CONT']['Всего выполнено требований на 4'] += 1
    dict_cont_vnd['ITSM-CONT']['Всего требований'] = dict_cont_vnd['ITSM-CONT']['Всего требований на 3'] + dict_cont_vnd['ITSM-CONT']['Всего требований на 4']


    # ********************************************** ITSM-CAP ********************************************************************
    # Расчет требований по процессу inc для 3-го уровня готовности
    for k, v in dict_cap.items():
        if v[2] == 'нет':
            dict_all_answer_no.append(k)
        if v[0] == '3':
            dict_cap_vnd['ITSM-CAP']['Всего требований на 3'] += 1
        if v[0] == '4':
            dict_cap_vnd['ITSM-CAP']['Всего требований на 4'] += 1
        if v[2] == 'да' and v[0] == '3':
            dict_cap_vnd['ITSM-CAP']['Всего выполнено требований на 3'] += 1
        if v[2] == 'да' and v[0] == '4':
            dict_cap_vnd['ITSM-CAP']['Всего выполнено требований на 4'] += 1
    dict_cap_vnd['ITSM-CAP']['Всего требований'] = dict_cap_vnd['ITSM-CAP']['Всего требований на 3'] + dict_cap_vnd['ITSM-CAP']['Всего требований на 4']


    # ********************************************** ITSM-MNTW ********************************************************************
    # Расчет требований по процессу inc для 3-го уровня готовности
    for k, v in dict_mntw.items():
        if v[2] == 'нет':
            dict_all_answer_no.append(k)
        if v[0] == '3':
            dict_mntw_vnd['ITSM-MNTW']['Всего требований на 3'] += 1
        if v[0] == '4':
            dict_mntw_vnd['ITSM-MNTW']['Всего требований на 4'] += 1
        if v[2] == 'да' and v[0] == '3':
            dict_mntw_vnd['ITSM-MNTW']['Всего выполнено требований на 3'] += 1
        if v[2] == 'да' and v[0] == '4':
            dict_mntw_vnd['ITSM-MNTW']['Всего выполнено требований на 4'] += 1
    dict_mntw_vnd['ITSM-MNTW']['Всего требований'] = dict_mntw_vnd['ITSM-MNTW']['Всего требований на 3'] + dict_mntw_vnd['ITSM-MNTW']['Всего требований на 4']


    # ********************************************** ITSM-SD ********************************************************************
    # Расчет требований по процессу inc для 3-го уровня готовности
    for k, v in dict_sd.items():
        if v[2] == 'нет':
            dict_all_answer_no.append(k)
        if v[0] == '3':
            dict_sd_vnd['ITSM-SD']['Всего требований на 3'] += 1
        if v[0] == '4':
            dict_sd_vnd['ITSM-SD']['Всего требований на 4'] += 1
        if v[2] == 'да' and v[0] == '3':
            dict_sd_vnd['ITSM-SD']['Всего выполнено требований на 3'] += 1
        if v[2] == 'да' and v[0] == '4':
            dict_sd_vnd['ITSM-SD']['Всего выполнено требований на 4'] += 1
    dict_sd_vnd['ITSM-SD']['Всего требований'] = dict_sd_vnd['ITSM-SD']['Всего требований на 3'] + dict_sd_vnd['ITSM-SD']['Всего требований на 4']

    #Словарь где собраны все отклонения из базы знаний для формирования итогового ЕРО.
    dict_all_baza = {}


    for keys, values in book_znanii.items():
        if keys in dict_all_answer_no:
            dict_all_baza[keys] = values


    # print(book_znanii)
    # print(dict_all_baza)


    #Поиск и удаление \xa0 (пробел в Unicode) из строк, функция рекурсивно проходит через словарь, списки и строки
    def replace_non_breaking_spaces(d):
        if isinstance(d, dict):
            return {k: replace_non_breaking_spaces(v) for k, v in d.items()}
        elif isinstance(d, list):
            return [replace_non_breaking_spaces(item) for item in d]
        elif isinstance(d, str):
            return d.replace('\xa0', ' ')#.replace('\n', '')
        else:
            return d


    result_my_dict = replace_non_breaking_spaces(dict_all_baza)

    # Все словари в переменной result_dicts_itsm
    result_dicts_itsm = {'ITSM-INC': dict_inc,
                    'ITSM-CATM': dict_catm,
                    'ITSM-CFG': dict_cfg,
                    'ITSM-SLM': dict_slm,
                    'ITSM-PRB': dict_prb,
                    'ITSM-ERM': dict_erm,
                    'ITSM-RFF': dict_rff,
                    'ITSM-ICHG': dict_ichg,
                    'ITSM-UPD': dict_upd,
                    'ITSM-AVL': dict_avl,
                    'ITSM-M&E': dict_me,
                    'ITSM-CONT': dict_cont,
                    'ITSM-CAP': dict_cap,
                    'ITSM-MNTW': dict_mntw,
                    'ITSM-SD': dict_sd
                    }

    # Объединяем все словари в один
    combined_dict = {}
    for key, value in result_dicts_itsm.items():
        combined_dict.update(value)



    # Создаем пустой словарь для хранения результатов
    combined_dict_itsm = {}

    # Проходимся по внутренним словарям и суммируем значения
    for key in dict_inc_vnd['ITSM-INC']:
        value = (
                dict_inc_vnd['ITSM-INC'].get(key, 0) +
                dict_catm_vnd['ITSM-CATM'].get(key, 0) +
                dict_cfg_vnd['ITSM-CFG'].get(key, 0) +
                dict_slm_vnd['ITSM-SLM'].get(key, 0) +
                dict_prb_vnd['ITSM-PRB'].get(key, 0) +
                dict_erm_vnd['ITSM-ERM'].get(key, 0) +
                dict_rff_vnd['ITSM-RFF'].get(key, 0) +
                dict_ichg_vnd['ITSM-ICHG'].get(key, 0) +
                dict_upd_vnd['ITSM-UPD'].get(key, 0) +
                dict_avl_vnd['ITSM-AVL'].get(key, 0) +
                dict_me_vnd['ITSM-M&E'].get(key, 0) +
                dict_cont_vnd['ITSM-CONT'].get(key, 0))
                # dict_cap_vnd['ITSM-CAP'].get(key, 0) +
                # dict_mntw_vnd['ITSM-MNTW'].get(key, 0) +
                # dict_sd_vnd['ITSM-SD'].get(key, 0))

        combined_dict_itsm[key] = value  #Подсчет всех требований для 12 - ти ПРОЦЕССОВ для расчета индекса надежности и 3, 4-го уровней надежности!!!

        # Создаем новый словарь, который будет объединять все словари - нужно для вывода на листе Надежности
        combined_dict_all_proc = {}
        combined_dict_all_proc.update(dict_inc_vnd)
        combined_dict_all_proc.update(dict_catm_vnd)
        combined_dict_all_proc.update(dict_cfg_vnd)
        combined_dict_all_proc.update(dict_slm_vnd)
        combined_dict_all_proc.update(dict_prb_vnd)
        combined_dict_all_proc.update(dict_erm_vnd)
        combined_dict_all_proc.update(dict_rff_vnd)
        combined_dict_all_proc.update(dict_ichg_vnd)
        combined_dict_all_proc.update(dict_upd_vnd)
        combined_dict_all_proc.update(dict_avl_vnd)
        combined_dict_all_proc.update(dict_me_vnd)
        combined_dict_all_proc.update(dict_cont_vnd)
        combined_dict_all_proc.update(dict_cap_vnd)
        combined_dict_all_proc.update(dict_mntw_vnd)
        combined_dict_all_proc.update(dict_sd_vnd)


    #Нужно добавить в общий словарь для вывода 1 и 2-е значение каждого отклонения согласно ключу
    #Это нужно для вывода в Exccel ({'ITSM-INC-01': ['3', 'Критичный', 'да'])
    for id, vd in result_my_dict.items():
        for kd, md in combined_dict.items():
            if id == kd:
                result_my_dict[id].append([md[0]])
                result_my_dict[id].append([md[1]])



    # print(result_my_dict)
    #Подсчет всех отклонений по критичности result_all_itsm --->{'ITSM': {'Критичный': 84, 'Высокий': 89, 'Умеренный': 23}} у которых стоит НЕТ!!!
    counts_itsm = {}
    for key, values in result_my_dict.items():
        risk_level = ''.join(values[-1])

        if key not in counts_itsm:
            counts_itsm[key] = {'Критичный': 0, 'Высокий': 0, 'Умеренный': 0, 'Рекомендация на развитие':0}
            counts_itsm[key][risk_level] += 1

    result_all_itsm = {'ITSM': {}}
    # Проходим по каждому элементу исходного словаря
    for key in counts_itsm.keys():
        for status, value in counts_itsm[key].items():
            # Если статус еще не добавлен в итоговый словарь, добавляем ег1
            # с нулевым значением
            if status not in result_all_itsm['ITSM']:
                result_all_itsm['ITSM'][status] = 0
            # Суммируем значение статуса
            result_all_itsm['ITSM'][status] += value


    wb = Workbook()
    ws = wb.active
    sheet = wb.active

    ws.title = 'Единая дорожная карта'
    ws_help = wb.create_sheet(title='Справочник')
    second_workshet = wb.create_sheet(title='Показатели надежности')
    five_workshet = wb.create_sheet(title='Детальная статистика по АС')
    six_workshet = wb.create_sheet(title='Тепловая карта для 3-го уровня')
    seven_workshet = wb.create_sheet(title='Тепловая карта для 4-го уровня')



    # Определяем стиль границы
    thin_border = Border(left=Side(style='thin'),
                         right=Side(style='thin'),
                         top=Side(style='thin'),
                         bottom=Side(style='thin'))

    #Изменяем ячейки по высоте на втором активном листе
    second_workshet.row_dimensions[1].height = 50

    #Изменяем ячейки по высоте на первом листе
    ws.row_dimensions[1].height = 60


    #Статусы для листа "СПРАВОЧНИК"----------------------------------------------------------------------------------------------
    stata_help = {'Статус отклонения': ['(ДЗО) В работе - анализ',
                                        '(ДЗО) В работе – план устранения на проверку ПАО',
                                        '(ДЗО) В работе - реализация',
                                        '(ДЗО) Устранено - на проверку ПАО',
                                        '(ПАО) Устранено - подтверждено ПАО',
                                        '(ДЗО) Исключение из реестра - на проверку ПАО',
                                        '(ПАО) Исключение из реестра - подтверждено ПАО',
                                        '(ДЗО) Принятие рисков - на утверждение ПАО',
                                        '(ПАО) Принятие рисков - утверждено ПАО']}

    about_stata = {'Описание статуса':['ДЗО анализирует найденное отклонение и рекомендации',
                                       'ДЗО проработало необходимые мероприятия для устранения отклонения, ставит ДЗО',
                                       'ПАО подтверждает мероприятия, проработанные ДЗО для устранения отклонения, ДЗО реализует мероприятия',
                                       'ДЗО устранило отклонение, статус ставит ДЗО',
                                       'ПАО согласно с тем, что отклонение считается устраненным, статус ставит ДЗО',
                                       'ДЗО несогласно с отклонением и предлагает его исключить из списка на устранение, статус ставит ДЗО',
                                       'ПАО согласно с тем, что данное отклонение исключается из списка на устранение, статус ставит ПАО',
                                       'ДЗО принимает риски, связанные с отклонением и не будет его устранять, статус ставит ДЗО',
                                       'ПАО согласно с тем, что ДЗО принимает риски, связанные с отклонением и не будет его устранять, статус ставит ПАО']}


    # Заголовки
    ws_help['P1'].value = 'Статус отклонения'
    ws_help['Q1'].value = 'Описание статуса'

    # Форматируем заголовки
    ws_help['P1'].font = Font(name='Calibri', size=11, bold=True)
    ws_help['Q1'].font = Font(name='Calibri', size=11, bold=True)
    ws_help['P1'].font = Font(name='Calibri', size=12, bold=True)
    ws_help.column_dimensions['P'].width = 35
    ws_help['Q1'].font = Font(name='Calibri', size=12, bold=True)


    # Заполняем статусы отклонений
    start_row_zero = 2
    for status in stata_help['Статус отклонения']:
        cell = ws_help.cell(row=start_row_zero, column=16)
        cell.value = status
        cell.font = Font(name='Calibri', size=11, bold=False)
        start_row_zero += 1

    # Заполняем описания статусов
    start_row_zero = 2
    for description in about_stata['Описание статуса']:
        cell = ws_help.cell(row=start_row_zero, column=17)
        cell.value = description
        cell.font = Font(name='Calibri', size=11, bold=False)
        start_row_zero += 1

    #----------------------------------------Конец вывода статусов отклонений----------------------------------------

    # Скрываем второй лист
    ws_help.sheet_state = 'hidden'  # Возможные состояния: visible, hidden, veryHidden
    dv = DataValidation(type="list", formula1='=Справочник!$P$2:$P$10', allow_blank=True)
    ws.add_data_validation(dv)


    #Изменяем ячейки по ширине во второй вкладке ("Показатели надежности")

    second_workshet['A1'].font = Font(name='Calibri', size=12, bold=True)
    second_workshet.column_dimensions['A'].width = 70



    second_workshet['B1'].font = Font(name='Calibri',size=12, bold=True)
    second_workshet.column_dimensions['B'].width = 30

    second_workshet['C1'].font = Font(name='Calibri',size=12, bold=True)
    second_workshet.column_dimensions['C'].width = 30

    second_workshet['D1'].font = Font(name='Calibri',size=12, bold=True)
    second_workshet.column_dimensions['D'].width = 20


    second_workshet['E1'].font = Font(name='Calibri',size=12, bold=True)
    second_workshet.column_dimensions['E'].width = 20

    second_workshet['F1'].font = Font(name='Calibri',size=12, bold=True)
    second_workshet.column_dimensions['F'].width = 20

    second_workshet['G1'].font = Font(name='Calibri',size=12, bold=True)
    second_workshet.column_dimensions['G'].width = 20

    second_workshet['H1'].font = Font(name='Calibri',size=12, bold=True)
    second_workshet.column_dimensions['H'].width = 30

    second_workshet['I1'].font = Font(name='Calibri',size=12, bold=True)
    second_workshet.column_dimensions['I'].width = 30



    ws.column_dimensions['A'].width = 25
    ws['A1'].font = Font(name='Calibri', size=8, bold=True)
    ws.column_dimensions['B'].width = 24
    ws['B1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['C'].width = 20
    ws['C1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['D'].width = 25
    ws['D1'].font = Font(name='Calibri', size=8, bold=True)
    ws.column_dimensions['E'].width = 31
    ws['E1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['F'].width = 18
    ws['F1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['G'].width = 10
    ws['G1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['H'].width = 10
    ws['H1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['I'].width = 16
    ws['I1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['J'].width = 16
    ws['J1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['K'].width = 18
    ws['K1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['L'].width = 25
    ws['L1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['M'].width = 50
    ws['M1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['N'].width = 50
    ws['N1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['O'].width = 14
    ws['O1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['P'].width = 17
    ws['P1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['Q'].width = 20
    ws['Q1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['R'].width = 22
    ws['R1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['S'].width = 17
    ws['S1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['T'].width = 17
    ws['T1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['U'].width = 14
    ws['U1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['V'].width = 14
    ws['V1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['W'].width = 35
    ws['W1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['X'].width = 14
    ws['X1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['Y'].width = 22
    ws['Y1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['Z'].width = 22
    ws['Z1'].font = Font(name='Calibri',size=8, bold=True)



    #Изменяем ячейки по ширине в 4-й вкладке ("Детальная статитстика по АС")

    five_workshet['A1'].font = Font(name='Calibri', size=12, bold=True)
    five_workshet.column_dimensions['A'].width = 70
    five_workshet['B1'].font = Font(name='Calibri',size=12, bold=True)
    five_workshet.column_dimensions['B'].width = 30
    five_workshet['C1'].font = Font(name='Calibri',size=12, bold=True)
    five_workshet.column_dimensions['C'].width = 30
    five_workshet['D1'].font = Font(name='Calibri',size=12, bold=True)
    five_workshet.column_dimensions['D'].width = 30
    five_workshet['E1'].font = Font(name='Calibri',size=12, bold=True)
    five_workshet.column_dimensions['E'].width = 30


    six_workshet['A1'].font = Font(name='Calibri', size=12, bold=True)
    six_workshet.column_dimensions['A'].width = 17
    six_workshet['B1'].font = Font(name='Calibri',size=12, bold=True)
    six_workshet.column_dimensions['B'].width = 17
    six_workshet['C1'].font = Font(name='Calibri',size=12, bold=True)
    six_workshet.column_dimensions['C'].width = 17
    six_workshet['D1'].font = Font(name='Calibri',size=12, bold=True)
    six_workshet.column_dimensions['D'].width = 17
    six_workshet['E1'].font = Font(name='Calibri',size=12, bold=True)
    six_workshet.column_dimensions['E'].width = 17
    six_workshet['F1'].font = Font(name='Calibri', size=12, bold=True)
    six_workshet.column_dimensions['F'].width = 17
    six_workshet['G1'].font = Font(name='Calibri',size=12, bold=True)
    six_workshet.column_dimensions['G'].width = 17
    six_workshet['H1'].font = Font(name='Calibri',size=12, bold=True)
    six_workshet.column_dimensions['H'].width = 17
    six_workshet['I1'].font = Font(name='Calibri',size=12, bold=True)
    six_workshet.column_dimensions['I'].width = 17
    six_workshet['J1'].font = Font(name='Calibri',size=12, bold=True)
    six_workshet.column_dimensions['J'].width = 17
    six_workshet['K1'].font = Font(name='Calibri', size=12, bold=True)
    six_workshet.column_dimensions['K'].width = 17
    six_workshet['L1'].font = Font(name='Calibri',size=12, bold=True)
    six_workshet.column_dimensions['L'].width = 17
    six_workshet['M1'].font = Font(name='Calibri',size=12, bold=True)
    six_workshet.column_dimensions['M'].width = 17
    six_workshet['N1'].font = Font(name='Calibri',size=12, bold=True)
    six_workshet.column_dimensions['N'].width = 17
    six_workshet['O1'].font = Font(name='Calibri',size=12, bold=True)
    six_workshet.column_dimensions['O'].width = 17
    six_workshet['P1'].font = Font(name='Calibri',size=12, bold=True)
    six_workshet.column_dimensions['P'].width = 17
    six_workshet['Q1'].font = Font(name='Calibri',size=12, bold=True)
    six_workshet.column_dimensions['Q'].width = 17
    six_workshet['R1'].font = Font(name='Calibri',size=12, bold=True)
    six_workshet.column_dimensions['R'].width = 17
    six_workshet['S1'].font = Font(name='Calibri',size=12, bold=True)
    six_workshet.column_dimensions['S'].width = 17
    six_workshet['T1'].font = Font(name='Calibri',size=12, bold=True)
    six_workshet.column_dimensions['T'].width = 17



    seven_workshet['A1'].font = Font(name='Calibri', size=12, bold=True)
    seven_workshet.column_dimensions['A'].width = 17
    seven_workshet['B1'].font = Font(name='Calibri',size=12, bold=True)
    seven_workshet.column_dimensions['B'].width = 17
    seven_workshet['C1'].font = Font(name='Calibri',size=12, bold=True)
    seven_workshet.column_dimensions['C'].width = 17
    seven_workshet['D1'].font = Font(name='Calibri',size=12, bold=True)
    seven_workshet.column_dimensions['D'].width = 17
    seven_workshet['E1'].font = Font(name='Calibri',size=12, bold=True)
    seven_workshet.column_dimensions['E'].width = 17
    seven_workshet['F1'].font = Font(name='Calibri', size=12, bold=True)
    seven_workshet.column_dimensions['F'].width = 17
    seven_workshet['G1'].font = Font(name='Calibri',size=12, bold=True)
    seven_workshet.column_dimensions['G'].width = 17
    seven_workshet['H1'].font = Font(name='Calibri',size=12, bold=True)
    seven_workshet.column_dimensions['H'].width = 17
    seven_workshet['I1'].font = Font(name='Calibri',size=12, bold=True)
    seven_workshet.column_dimensions['I'].width = 17
    seven_workshet['J1'].font = Font(name='Calibri',size=12, bold=True)
    seven_workshet.column_dimensions['J'].width = 17
    seven_workshet['K1'].font = Font(name='Calibri', size=12, bold=True)
    seven_workshet.column_dimensions['K'].width = 17
    seven_workshet['L1'].font = Font(name='Calibri',size=12, bold=True)
    seven_workshet.column_dimensions['L'].width = 17
    seven_workshet['M1'].font = Font(name='Calibri',size=12, bold=True)
    seven_workshet.column_dimensions['M'].width = 17
    seven_workshet['N1'].font = Font(name='Calibri',size=12, bold=True)
    seven_workshet.column_dimensions['N'].width = 17
    seven_workshet['O1'].font = Font(name='Calibri',size=12, bold=True)
    seven_workshet.column_dimensions['O'].width = 17
    seven_workshet['P1'].font = Font(name='Calibri',size=12, bold=True)
    seven_workshet.column_dimensions['P'].width = 17
    seven_workshet['Q1'].font = Font(name='Calibri',size=12, bold=True)
    seven_workshet.column_dimensions['Q'].width = 17
    seven_workshet['R1'].font = Font(name='Calibri',size=12, bold=True)
    seven_workshet.column_dimensions['R'].width = 17
    seven_workshet['S1'].font = Font(name='Calibri',size=12, bold=True)
    seven_workshet.column_dimensions['S'].width = 17
    seven_workshet['T1'].font = Font(name='Calibri',size=12, bold=True)
    seven_workshet.column_dimensions['T'].width = 17


    second_workshet['A1'] = 'Процесс'
    second_workshet['B1'] = 'Всего требований'
    second_workshet['C1'] = 'Требований на 3-й уровень'
    second_workshet['D1'] = 'Выполнено требований на 3-й уровень'
    second_workshet['E1'] = '%, достижения 3-го уровня'
    second_workshet['F1'] = 'Требований на 4-й уровень'
    second_workshet['G1'] = 'Выполнено требований на 4-й уровень'
    second_workshet['H1'] = '%, достижения 4-го уровня'
    # second_workshet['I1'] = 'Примечание'

    # Цвет для Критического уровня риска(подкрашиваем ячейки)
    high_risk_fill = PatternFill(start_color="FF9900", end_color="FF9900", fill_type="solid")
    high_profile_fill = PatternFill(start_color="9bbb59", end_color="9bbb59", fill_type="solid")


    five_workshet['A1'] = 'Название АС'
    five_workshet['B1'] = 'Критичность'
    five_workshet['C1'] = 'Уровень отклонения: Критичный'
    five_workshet['D1'] = 'Уровень отклонения: Высокий'
    five_workshet['E1'] = 'Уровень отклонения: Умеренный'

    six_workshet['A1'] = 'Компания'
    six_workshet['B1'] = 'Общий уровень надежности'
    six_workshet['C1'] = 'Уровень технологической надежности'
    six_workshet['D1'] = 'Управление технологическими инцидентами'
    six_workshet['E1'] = 'Управление каталогом сервисов'
    six_workshet['F1'] = 'Управление уровнями сервисов'
    six_workshet['G1'] = 'Управление технологическими проблемами'
    six_workshet['H1'] = 'Управление технологическими рисками'
    six_workshet['I1'] = 'Управление инфраструктурными изменениями'
    six_workshet['J1'] = 'Управление непрерывностью технологий'
    six_workshet['K1'] = 'Управление доступностью технологий'
    six_workshet['L1'] = 'Управление технологическими событиями и мониторингом'
    six_workshet['M1'] = 'Управление технологическими запросами на обслуживание'
    six_workshet['N1'] = 'Управление конфигурациями'
    six_workshet['O1'] = 'Управление обновлением ПО'
    six_workshet['P1'] = 'Ключевые проблемы'
    six_workshet['Q1'] = 'Динамические показатели ITSM'
    six_workshet['R1'] = 'Динамические показатели по технологиям'
    six_workshet['S1'] = 'Уровень технологической надежности'
    six_workshet['T1'] = 'Индекс надежности(%)'

    seven_workshet['A1'] = 'Компания'
    seven_workshet['B1'] = 'Общий уровень надежности'
    seven_workshet['C1'] = 'Уровень технологической надежности'
    seven_workshet['D1'] = 'Управление технологическими инцидентами'
    seven_workshet['E1'] = 'Управление каталогом сервисов'
    seven_workshet['F1'] = 'Управление уровнями сервисов'
    seven_workshet['G1'] = 'Управление технологическими проблемами'
    seven_workshet['H1'] = 'Управление технологическими рисками'
    seven_workshet['I1'] = 'Управление инфраструктурными изменениями'
    seven_workshet['J1'] = 'Управление непрерывностью технологий'
    seven_workshet['K1'] = 'Управление доступностью технологий'
    seven_workshet['L1'] = 'Управление технологическими событиями и мониторингом'
    seven_workshet['M1'] = 'Управление технологическими запросами на обслуживание'
    seven_workshet['N1'] = 'Управление конфигурациями'
    seven_workshet['O1'] = 'Управление обновлением ПО'
    seven_workshet['P1'] = 'Ключевые проблемы'
    seven_workshet['Q1'] = 'Динамические показатели ITSM'
    seven_workshet['R1'] = 'Динамические показатели по технологиям'
    seven_workshet['S1'] = 'Уровень технологической надежности'
    seven_workshet['T1'] = 'Индекс надежности(%)'



    ws['A1'] = 'Код задачи/отклонения'
    ws['B1'] = 'Вид объекта 1 уровень'
    ws['C1'] = 'Вид объекта 2 уровень'
    ws['D1'] = 'Вид объекта 3 уровень'
    ws['E1'] = 'Категория задачи / отклонения'
    ws['F1'] = 'Дополнительный комментарий по задаче / выявленному отклонению'
    ws['G1'] = '№'
    ws['H1'] = 'Инициатор'
    ws['I1'] = 'Основание'
    ws['J1'] = 'Код вопроса ЕОИТ'
    ws['K1'] = 'Код объекта'
    ws['L1'] = 'Объект'
    ws['M1'] = 'Задача / выявленное отклонение'
    # ws['N1'] = 'Негативные последствия'
    ws['N1'] = 'Рекомендованные мероприятия'
    ws['O1'] = 'Уровень принятия решения'
    ws['P1'] = 'Вхождение в уровни готовности'
    ws['Q1'] = 'Уровень критичности'

    ws['R1'] = 'Запланированное мероприятие'
    ws['S1'] = 'Код мероприятия из трекера компании Группы (при использовании трекера)'
    ws['T1'] = 'Дата начала мероприятий'
    ws['U1'] = 'Дата окончания мероприятий'
    ws['V1'] = 'Ответственный'
    ws['W1'] = 'Статус'
    ws['X1'] = 'Комментарии компании Группы'
    ws['Y1'] = 'Комментарии инициатора'


    # Записываем данные на первую страницу отклонения по ITSM и AS
    start_i = 0
    count = 0



    for i, (key, values) in enumerate(result_my_dict.items(), start=2):
        for id, m in dict_itsm.items():
            if id in key:

                ws[f'L{i}'] = m
                ws[f'K{i}'] = id
        ws[f'I{i}'] = "ДТН"
        ws[f'B{i}'] = "Процессы"
        ws[f'O{i}'] = "Уровень 1"
        dv.add(ws[f'W{i}']) # Добаввляем стату в колонку Y
        ws[f'G{i}'] = (int(f'{i}') - 1)
        ws[f'J{i}'] = key     #Начинаем писать данные со столбца --- F --- в ЕОИТ
        ws[f'M{i}'] = ''.join(values[0])
        # ws[f'N{i}'] = ''.join(values[1])   ----> Негативные последствия убрали из ЕОИТ
        ws[f'N{i}'] = ''.join(values[2])
        uroven = ''.join(values[3])

        if uroven == '3':
            ws[f'P{i}'] = 'Надежность 3'
        else:
            ws[f'P{i}'] = 'Надежность 4'
        ws[f'Q{i}'] = ''.join(values[4])

        if 'ITSM-CAP' in key or 'ITSM-SD' in key or 'ITSM-MNTW' in key:
            ws[f'Q{i}'] = 'Рекомендация на развитие'
            ws[f'P{i}'] = 'Надежность 4'


        #Изменяем шрифт у активных полей

        ws[f'B{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'C{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'E{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'F{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'G{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'H{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'I{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'J{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'K{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'L{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'M{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'N{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'O{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'P{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'Q{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'R{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'S{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'Y{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'X{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'W{i}'].font = Font(name='Calibri', size=8, bold=False)

        #Изменяем расположение активных полей

        ws[f'B{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'C{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'E{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'F{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'G{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'H{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'I{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'J{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'K{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'L{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'M{i}'].alignment = Alignment(horizontal="left", vertical="top", wrapText=True)
        ws[f'N{i}'].alignment = Alignment(horizontal="left", vertical="top", wrapText=True)
        ws[f'O{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'P{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'S{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'Q{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'R{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'Y{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'X{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'W{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        start_i = i


    for i, (key, values) in enumerate(my_function().items(), start=start_i):
        for j in values:
            count += 1
            ws[f'G{start_i+count}'] = (int(f'{start_i+count}') - 1)
            ws[f'I{start_i+count}'] = "ДТН"
            ws[f'B{start_i+count}'] = "АС"
            ws[f'O{start_i+count}'] = "Уровень 1"
            ws[f'J{start_i+count}'] = j[0]
            dv.add(ws[f'W{start_i+count}']) # Добаввляем стату в колонку X
            ws[f'K{start_i+count}'] = j[7]

            ws[f'L{start_i+count}'] = key
            ws[f'M{start_i+count}'] = j[2]
            # ws[f'N{start_i+count}'] = j[3] ----> Негативные последствия убрали из ЕОИТ
            ws[f'N{start_i+count}'] = j[4]
            try:
                if j[8] == 3:

                    ws[f'P{start_i+count}'] = 'Надежность 3'
                else:
                    ws[f'P{start_i+count}'] = 'Надежность 4'
                ws[f'Q{start_i+count}'] = j[5]
            except IndexError:
                # print(j)
                # обработать случай, когда элемента нет
                pass

            #Изменяем шрифт у активных полей
            ws[f'B{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'C{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'E{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'F{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'G{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'H{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'I{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'J{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'K{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'L{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'M{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'N{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'O{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'P{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'Q{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'R{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'S{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'Y{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'X{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'W{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)

            #Изменяем расположение активных полей
            ws[f'B{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'C{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'E{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'F{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'G{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'H{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'I{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'J{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'K{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'L{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'M{start_i+count}'].alignment = Alignment(horizontal="left", vertical="top", wrapText=True)
            ws[f'N{start_i+count}'].alignment = Alignment(horizontal="left", vertical="top", wrapText=True)
            ws[f'O{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'P{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'S{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'Q{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'R{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'Y{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'X{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'W{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)

    #Записываем данные на вторую страницу лист Надежность
    # print(combined_dict_all_proc)
    # print(all_treb_as) ## {'Address Processing': [(38.0, 2, 11.0, 1.5)]}
    # print(data_new_count) ##{'Address Processing': ['BO', 0, 74, 10]}
    # print(all_treb_as)

    def transform_data(input_data):
        output = {'Всего требований': 0, 'Всего требований на 3': 0, 'Всего выполнено требований на 3': 0, 'Всего требований на 4': 0, 'Всего выполнено требований на 4': 0}
        for process, values in input_data.items():
            for value_tuple in values:
                critical = value_tuple[0]
                critical_yes = value_tuple[1]
                high = value_tuple[2]
                high_yes = value_tuple[3]

                output['Всего требований на 3'] += critical
                output['Всего выполнено требований на 3'] += critical_yes
                output['Всего требований на 4'] += high
                output['Всего выполнено требований на 4'] += high_yes
                output['Всего требований'] = output['Всего требований на 3'] + output['Всего требований на 4']
        return output

        # Преобразование данных
    output_data = transform_data(all_treb_as)

    #ITSM
    #{'Всего требований': 223, 'Всего требований на 3': 158, 'Всего выполнено требований на 3': 21, 'Всего требований на 4': 65, 'Всего выполнено требований на 4': 0}
    # print(combined_dict_itsm)
    #AS
    #{'Всего требований': 103.0, 'Всего требований на 3': 81.0, 'Всего выполнено требований на 3': 4, 'Всего требований на 4': 22.0, 'Всего выполнено требований на 4': 3.0}
    # print(output_data)

    # Объединение словарей
    combined_dict_it_as = {}

    for key in combined_dict_itsm:
        combined_dict_it_as[key] = combined_dict_itsm.get(key, 0) + output_data.get(key, 0)


    index_nadejnost = (combined_dict_it_as['Всего выполнено требований на 3'] +
                       combined_dict_it_as['Всего выполнено требований на 4']) / (combined_dict_it_as['Всего требований на 3'] +
                                                                                  combined_dict_it_as['Всего требований на 4']) * 100


    second_workshet['B22'] = str(round(index_nadejnost, 0)) + '%'


    three_uroven_nadejnosti = combined_dict_it_as['Всего выполнено требований на 3']/combined_dict_it_as['Всего требований на 3'] + 2
    four_uroven_nadejnosti = combined_dict_it_as['Всего выполнено требований на 4']/combined_dict_it_as['Всего требований на 4'] + 3
    second_workshet['A25'] = '3-й уровень готовности  надежности'
    second_workshet['A26'] = '4-й уровень готовности  надежности'
    second_workshet['B25'] = float(round(three_uroven_nadejnosti, 2))
    second_workshet['B26'] = float(round(four_uroven_nadejnosti, 2))
    second_workshet['B20'] = float(second_workshet['B25'].value)



    for i, (key, values) in enumerate(dict_itsm.items(), start=2):
        if key not in ['ITSM-CAP', 'ITSM-SD', 'ITSM-MNTW']:
            second_workshet[f'A{i}'] = values
            second_workshet[f'B{i}'] = combined_dict_all_proc[key]['Всего требований']
            second_workshet[f'C{i}'] = combined_dict_all_proc[key]['Всего требований на 3']
            second_workshet[f'D{i}'] = combined_dict_all_proc[key]['Всего выполнено требований на 3']

            if combined_dict_all_proc[key]['Всего требований на 3'] == 0:
                second_workshet[f'E{i}'] = 0
            else:
                second_workshet[f'E{i}'] = f"{int(round(second_workshet[f'D{i}'].value / second_workshet[f'C{i}'].value * 100))}%"
                # second_workshet[f'E{i}'] = str(round(second_workshet[f'D{i}'].value / second_workshet[f'C{i}'].value, 2) * 100) + '%'
            second_workshet[f'F{i}'] = combined_dict_all_proc[key]['Всего требований на 4']
            second_workshet[f'G{i}'] = combined_dict_all_proc[key]['Всего выполнено требований на 4']
            if combined_dict_all_proc[key]['Всего требований на 4'] == 0:
                second_workshet[f'H{i}'] = 0
            else:
                # second_workshet[f'H{i}'] = str(round(second_workshet[f'G{i}'].value / second_workshet[f'F{i}'].value, 2) * 100) + '%'
                second_workshet[f'H{i}'] = f"{int(round(second_workshet[f'G{i}'].value / second_workshet[f'F{i}'].value * 100))}%"




        second_workshet[f'B{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        second_workshet[f'C{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        second_workshet[f'D{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        second_workshet[f'E{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        second_workshet[f'F{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        second_workshet[f'G{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        second_workshet[f'H{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        second_workshet[f'I{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)


        second_workshet['A16'] = 'Выбран профиль'
        second_workshet['B16'].fill = high_profile_fill
        second_workshet['B16'] = result_profile['Каким максимальным уровнем критичности обладает ИТ-сервис в компании?']
        second_workshet['A17'] = 'Динамические показатели'
        second_workshet['A18'] = 'ITSM'
        second_workshet['B17'] = 'Критичное'
        # second_workshet['B18'] = str(dict_tree_four_assesment_itsm['Всего критичных']) + '/' + str(dict_tree_four_assesment_itsm['Выполнено критичных'])
        second_workshet['B18'] = str(result_all_itsm['ITSM']['Критичный']) + ' / ' + '0'
        second_workshet['C17'] = 'Высокое'
        # second_workshet['C18'] = str(dict_tree_four_assesment_itsm['Всего высоких']) + '/' + str(dict_tree_four_assesment_itsm['Выполнено высоких'])
        second_workshet['C18'] = str(result_all_itsm['ITSM']['Высокий']) + ' / ' + '0'
        second_workshet['D17'] = 'Умеренное'
        # second_workshet['D18'] = str(dict_tree_four_assesment_itsm['Всего умеренных']) + '/' + str(dict_tree_four_assesment_itsm['Выполнено умеренных'])
        second_workshet['D18'] = str(result_all_itsm['ITSM']['Умеренный'])  + ' / ' + '0'
        second_workshet['A20'] = 'На слайд, в графу "Уровень надёжности" ='
        second_workshet['A22'] = 'Индекс надежности для ЦТК и тепловой карты'


        #РАСЧЕТЫ ИНДЕКСА НАДЕЖНОСТИ, А ТАКЖЕ УРОВНЯ ГОТОВНОСТИ ПО УРОВНЯМ

        # print(combined_dict_all_proc)
        # print(all_treb_as) ## {'Address Processing': [(38.0, 2, 11.0, 1.5)]}
        # print(combined_dict_itsm)

        # print(combined_dict_all_proc)



    # выравниваем ячейки, только имеющие значения
    for cell in ws[1]:
        if cell.value:

            cell.alignment = Alignment(horizontal="center", vertical="center")
            ws['C1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws['D1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws['E1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws['F1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws['G1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws['H1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws['I1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws['J1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws['K1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws['L1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws['N1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws['Q1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws['P1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws['O1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws['R1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws['S1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws['T1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws['U1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws['V1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws['W1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws['X1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws['Y1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws['Z1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws['AA1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws['AB1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws['AG1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws['AH1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws['AC1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)


   # Нужно выполнить проверку если АС отсутствует в момент сборки ЕРО,
   # т.к. в момент оценки могут быть выполнены все отклонения!!!


    dict_reestr_all_as = {}
    for key, values in all_treb_as.items():
        if key not in data_new_count:
            dict_reestr_all_as[key] = [rezult_as[key][0], 0, 0, 0]
        else:
            dict_reestr_all_as[key] = data_new_count[key]
    # print(dict_reestr_all_as)
    # print(all_treb_as) ----> {'Address Processing': [(38.0, 2, 11.0, 1.5)], 'Testing': [(43.0, 43.0, 11.0, 11.0)]}
    # print(data_new_count) ----> {'Address Processing': ['BO', 0, 74, 10]}
    # print(rezult_as) ----> {'Address Processing': ['BO', 2.05, 3.14], 'Testing': ['BC', 3, 4]}
    # print(dict_reestr_all_as) ----> {'Address Processing': ['BO', 0, 74, 10], 'Testing': ['BC', 0, 0, 0]}


    #Заполняем ЛИСТ ИТ-ТЕХНОЛОГИЙ
    count_as = 0
    for i, (key, values) in enumerate(dict_reestr_all_as.items(), start=2):
        count_as += 1

        five_workshet[f'A{i}'] = key
        five_workshet[f'B{i}'] = values[0]
        five_workshet[f'C{i}'] = values[1]
        five_workshet[f'D{i}'] = values[2]
        five_workshet[f'E{i}'] = values[3]


        five_workshet[f'A{i}'].alignment = Alignment(horizontal="left", vertical="center", wrapText=True)
        five_workshet[f'B{i}'].alignment = Alignment(horizontal="center", vertical="top", wrapText=True)
        five_workshet[f'C{i}'].alignment = Alignment(horizontal="center", vertical="top", wrapText=True)
        five_workshet[f'D{i}'].alignment = Alignment(horizontal="center", vertical="top", wrapText=True)
        five_workshet[f'E{i}'].alignment = Alignment(horizontal="center", vertical="top", wrapText=True)

        six_workshet[f'B{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        six_workshet[f'C{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        six_workshet[f'D{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        six_workshet[f'E{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        six_workshet[f'F{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        six_workshet[f'G{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        six_workshet[f'H{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        six_workshet[f'I{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        six_workshet[f'J{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        six_workshet[f'K{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        six_workshet[f'L{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        six_workshet[f'M{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        six_workshet[f'N{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        six_workshet[f'O{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        six_workshet[f'P{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)

        seven_workshet[f'B{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        seven_workshet[f'C{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        seven_workshet[f'D{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        seven_workshet[f'E{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        seven_workshet[f'F{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        seven_workshet[f'G{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        seven_workshet[f'H{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        seven_workshet[f'I{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        seven_workshet[f'J{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        seven_workshet[f'K{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        seven_workshet[f'L{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        seven_workshet[f'M{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        seven_workshet[f'N{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        seven_workshet[f'O{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        seven_workshet[f'P{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)

        five_workshet[f'A{i}'].font = Font(name='Calibri', size=11, bold=False)
        five_workshet[f'E{i}'].font = Font(name='Calibri', size=11, bold=False)
        five_workshet[f'B{i}'].font = Font(name='Calibri', size=11, bold=False)
        five_workshet[f'C{i}'].font = Font(name='Calibri', size=11, bold=False)
        five_workshet[f'D{i}'].font = Font(name='Calibri', size=11, bold=False)
        five_workshet[f'E{i}'].font = Font(name='Calibri', size=11, bold=False)


    try:
        uroven_tech_nadejnosti_three = (output_data['Всего выполнено требований на 3']) / (output_data['Всего требований на 3'])  + 2
    except ZeroDivisionError:
        uroven_tech_nadejnosti_three = 2

    try:
        uroven_tech_nadejnosti_four = (output_data['Всего выполнено требований на 4']) / (output_data['Всего требований на 4']) + 3

    except ZeroDivisionError:
        uroven_tech_nadejnosti_four = 3


    five_workshet['A1'] = 'Название АС'
    five_workshet['B1'] = 'Критичность'
    five_workshet['C1'] = 'Уровень отклонения: Критичный'
    five_workshet['D1'] = 'Уровень отклонения: Высокий'
    five_workshet['E1'] = 'Уровень отклонения: Умеренный'
    five_workshet[f'B{count_as + 8}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    five_workshet[f'A{count_as + 4}'] = 'Динамические показатели'
    five_workshet[f'A{count_as + 5}'] = 'АС'
    five_workshet[f'A{count_as + 7}'] = 'Доп. показатели / инфо'
    five_workshet[f'A{count_as + 8}'] = 'На слайд, в графу "Уровень технологической надежности" ='



    #Убрал, проверку пока все компании не дойдут до 3-го уровня надежности!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!
    if output_data['Всего выполнено требований на 3'] == output_data['Всего требований на 3'] and output_data['Всего требований на 3'] != 0:
        five_workshet[f'B{count_as + 8}'] = str(round(uroven_tech_nadejnosti_four, 2))

    else:
        five_workshet[f'B{count_as + 8}'] = str(round(uroven_tech_nadejnosti_three, 2))

    five_workshet[f'B{count_as + 4}'] = 'Критичное'
    five_workshet[f'C{count_as + 4}'] = 'Высокое'
    five_workshet[f'D{count_as + 4}'] = 'Умеренное'
    five_workshet[f'A{count_as + 8}'].font = Font(name='Calibri', size=11, bold=True)
    five_workshet[f'B{count_as + 8}'].font = Font(name='Calibri', size=11, bold=True)


    # Эта часть кода  для Подсчета ОБЩЕГО КОЛИЧЕСТВА ОТКЛОНЕНИЙ ПО ВСЕМ АС ===================> на странице Детальная Статистика
    # Итоговый словарь

    new_statistics_as = {}

    # Инициализация счетчиков
    critical = [0]
    high = [0]
    moderate = [0]


    # Проходим по каждому элементу исходного словаря
    for key, values in data_new_count.items():
        # Пропускаем элементы, которые не содержат '/' в первом элементе
        # Критичные
        critical.append(values[1])
        high.append(values[2])
        moderate.append(values[3])

    # Заполнение итогового словаря
    new_statistics_as['Критичное'] = sum(critical)
    new_statistics_as['Высокое'] = sum(high)
    new_statistics_as['Умеренное'] = sum(moderate)

    five_workshet[f"B{count_as + 5}"] = str(new_statistics_as['Критичное']) + ' / ' + '0'
    five_workshet[f"C{count_as + 5}"] = str(new_statistics_as['Высокое']) + ' / ' + '0'
    five_workshet[f"D{count_as + 5}"] = str(new_statistics_as['Умеренное']) + ' / ' + '0'



    # Функция для замены % в значении, если оно есть
    def remove_percent(value):
        if isinstance(value, str):  # Проверяем, является ли значение строкой
            return value.replace('%', '')  # Удаляем символ %
        return value  # Если это не строка, возвращаем значение как есть

    six_workshet['A2'] = name_dtn
    six_workshet['B2'] = second_workshet['B25'].value
    if output_data['Всего выполнено требований на 3'] == output_data['Всего требований на 3']:
        six_workshet['C2'] = '3.0'
    else:
        value = round(uroven_tech_nadejnosti_three, 2)
        six_workshet['C2'] = f"{value:.2f}".replace('.', ',')  # Заменяем запятую на точку

    six_workshet['D2'] = remove_percent(second_workshet['E2'].value)
    six_workshet['E2'] = remove_percent(second_workshet['E3'].value)
    six_workshet['F2'] = remove_percent(second_workshet['E4'].value)
    six_workshet['G2'] = remove_percent(second_workshet['E5'].value)


    if dict_erm_vnd['ITSM-ERM']['Всего требований на 3'] != 0:
        six_workshet['H2'] = remove_percent(second_workshet['E6'].value)
    else:
        six_workshet['H2'] = '-'

    six_workshet['I2'] = remove_percent(second_workshet['E7'].value)

    if dict_cont_vnd['ITSM-CONT']['Всего требований на 3'] != 0:
        six_workshet['J2'] = remove_percent(second_workshet['E8'].value)
    else:
        six_workshet['J2'] = '-'

    if dict_avl_vnd['ITSM-AVL']['Всего требований на 3'] != 0:
        six_workshet['K2'] = remove_percent(second_workshet['E9'].value)
    else:
        six_workshet['K2'] = '-'
    six_workshet['L2'] = remove_percent(second_workshet['E10'].value)
    six_workshet['M2'] = remove_percent(second_workshet['E11'].value)
    six_workshet['N2'] = remove_percent(second_workshet['E12'].value)

    if dict_upd_vnd['ITSM-UPD']['Всего требований на 3'] != 0:
        six_workshet['O2'] = remove_percent(second_workshet['E13'].value)
    else:
        six_workshet['O2'] = '-'


    seven_workshet['A2'] = name_dtn
    seven_workshet['B2'] = second_workshet['B26'].value
    seven_workshet_one = round(uroven_tech_nadejnosti_four, 2)
    seven_workshet['C2'] = f"{seven_workshet_one:.2f}".replace('.', ',')  # Заменяем запятую на точку
    seven_workshet['D2'] = remove_percent(second_workshet['H2'].value)
    seven_workshet['E2'] = remove_percent(second_workshet['H3'].value)
    seven_workshet['F2'] = remove_percent(second_workshet['H4'].value)
    seven_workshet['G2'] = remove_percent(second_workshet['H5'].value)
    seven_workshet['H2'] = remove_percent(second_workshet['H6'].value)
    seven_workshet['I2'] = remove_percent(second_workshet['H7'].value)
    seven_workshet['J2'] = remove_percent(second_workshet['H8'].value)
    seven_workshet['K2'] = remove_percent(second_workshet['H9'].value)
    seven_workshet['L2'] = remove_percent(second_workshet['H10'].value)
    seven_workshet['M2'] = remove_percent(second_workshet['H11'].value)
    seven_workshet['N2'] = remove_percent(second_workshet['H12'].value)
    seven_workshet['O2'] = remove_percent(second_workshet['H13'].value)


    for cell in five_workshet[1]:
        if cell.value:
            # Устанавливаем выравнивание для соответствующих ячеек
            five_workshet[f"B{count_as + 4}"].alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
            five_workshet[f"C{count_as + 4}"].alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
            five_workshet[f"D{count_as + 4}"].alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
            five_workshet[f"B{count_as + 5}"].alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
            five_workshet[f"C{count_as + 5}"].alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
            five_workshet[f"D{count_as + 5}"].alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)


    #-----------------Записываю результаты в Тепловую карту показания технологической надежности для АС--------------------
    six_workshet['R2'] = (
            str(new_statistics_as['Критичное']) + ' / ' + '0' + ', ' +
            str(new_statistics_as['Высокое']) + ' / ' + '0' + ', ' +
            str(new_statistics_as['Умеренное']) + ' / ' + '0'
    )

    six_workshet['Q2'] = (
            str(result_all_itsm['ITSM']['Критичный']) + ' / ' + '0' + ', ' +
            str(result_all_itsm['ITSM']['Высокий']) + ' / ' + '0' + ', ' +
            str(result_all_itsm['ITSM']['Умеренный']) + ' / ' + '0'

    )
    six_workshet['T2'] = second_workshet['B22'].value

    critical_value = second_workshet['B16'].value
    # Запись с проверкой допустимых значений
    if critical_value in ['OP', 'BO', 'BC и выше']:  # Если значение из допустимого списка
        six_workshet['S2'] = critical_value
        seven_workshet['S2'] = critical_value
    elif critical_value is not None:  # Если есть значение, но недопустимое
        six_workshet['S2'] = f"Некорректное значение: {critical_value}. Введите OP, BO, BC и выше"
        seven_workshet['S2'] = f"Некорректное значение: {critical_value}. Введите OP, BO, BC и выше"
    else:  # Если значение отсутствует
        six_workshet['S2'] = "Введите Критичность вручную (OP, BO, BC и выше)"
        seven_workshet['S2'] = "Введите Критичность вручную (OP, BO, BC и выше)"


    seven_workshet['R2'] = (
            str(new_statistics_as['Критичное']) + ' / ' + '0' + ', ' +
            str(new_statistics_as['Высокое']) + ' / ' + '0' + ', ' +
            str(new_statistics_as['Умеренное']) + ' / ' + '0'
    )

    seven_workshet['Q2'] = (
            str(result_all_itsm['ITSM']['Критичный']) + ' / ' + '0' + ', ' +
            str(result_all_itsm['ITSM']['Высокий']) + ' / ' + '0' + ', ' +
            str(result_all_itsm['ITSM']['Умеренный']) + ' / ' + '0'

    )

    seven_workshet['T2'] = second_workshet['B22'].value



    six_workshet['P2'] = 'Ключевые проблемы (максимум 5 шт.), выбираем вручную, каждая отделяется друг от друга точкой с запятой. (Для удобства в ячейке  используем alt+Enter)'
    seven_workshet['P2'] = 'Ключевые проблемы (максимум 5 шт.), выбираем вручную, каждая отделяется друг от друга точкой с запятой. (Для удобства в ячейке  используем alt+Enter)'



    # выравниваем ячейки, только имеющие значения
    for cell in second_workshet[1]:
        if cell.value:

            cell.alignment = Alignment(horizontal="center", vertical="center")

            #Изменяем шрифт у активных полей
            second_workshet['A20'].font = Font(name='Calibri', size=11, bold=True)
            second_workshet['A21'].font = Font(name='Calibri', size=11, bold=True)
            second_workshet['A22'].font = Font(name='Calibri', size=11, bold=True)
            second_workshet['A16'].font = Font(name='Calibri', size=11, bold=True)
            second_workshet['A23'].font = Font(name='Calibri', size=11, bold=True)
            second_workshet['A24'].font = Font(name='Calibri', size=11, bold=True)
            second_workshet['A25'].font = Font(name='Calibri', size=11, bold=True)
            second_workshet['A26'].font = Font(name='Calibri', size=11, bold=True)
            second_workshet['A27'].font = Font(name='Calibri', size=11, bold=True)
            second_workshet['B20'].font = Font(name='Calibri', size=11, bold=True)
            second_workshet['B21'].font = Font(name='Calibri', size=11, bold=True)
            second_workshet['B22'].font = Font(name='Calibri', size=11, bold=True)
            second_workshet['B16'].font = Font(name='Calibri', size=11, bold=True)
            second_workshet['B22'].font = Font(name='Calibri', size=11, bold=True)
            second_workshet['B23'].font = Font(name='Calibri', size=11, bold=True)
            second_workshet['B24'].font = Font(name='Calibri', size=11, bold=True)
            second_workshet['B25'].font = Font(name='Calibri', size=11, bold=True)
            second_workshet['B26'].font = Font(name='Calibri', size=11, bold=True)
            second_workshet['B27'].font = Font(name='Calibri', size=11, bold=True)



            second_workshet['B1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['C1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['D1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['E1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['F1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['G1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['H1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['A16'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['B16'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['B17'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['B18'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['B20'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['B22'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['B25'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['B26'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['C17'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['C18'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['D17'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['D18'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)


            five_workshet['A1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            five_workshet['B1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            five_workshet['C1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            five_workshet['D1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            five_workshet['E1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)


            six_workshet['A1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['B1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['C1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['D1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['E1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['F1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['G1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['H1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['I1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['J1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['K1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['L1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['M1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['N1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['O1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['P1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['Q1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['R1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['S1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['T1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['A2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['B2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['C2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['D2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['E2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['F2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['G2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['H2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['I2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['J2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['K2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['L2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['M2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['N2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['O2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['P2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['Q2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['R2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['S2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            six_workshet['T2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)

            seven_workshet['A1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['B1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['C1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['D1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['E1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['F1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['G1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['H1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['I1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['J1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['K1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['L1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['M1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['N1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['O1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['P1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['Q1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['R1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['S1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['T1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['A2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['B2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['C2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['D2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['E2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['F2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['G2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['H2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['I2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['J2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['K2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['L2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['M2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['N2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['O2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['P2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['Q2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['R2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['S2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            seven_workshet['T2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)



    #Для прорисовки таблицы на листе ---===Показатели Надежности===---
    for row in second_workshet['A17':'D18']:
        for cell in row:
            cell.border = thin_border

    for row in five_workshet[f'A{count_as + 4}':f'D{count_as + 5}']:
        for cell in row:
            cell.border = thin_border

    # Проверяем существование файла
    # Проверяем существование файла
    # if os.path.exists(save + '\\' + f'{name_company}' + '_ЕРО.xlsx'):

    while True:
        usa_present = input('Создать отчет ЕРО + слайд? - "y", Создать только отчет ЕРО - "n"): ').strip().lower()
        if usa_present in ('y', 'n'):
            break
        print("Пожалуйста, введите 'y' для да или 'n' для нет.")


    # Теперь usa_present гарантированно содержит 'y' или 'n'

    if usa_present == 'y':
        #=======================================================================================================================

        #=============================================Генерация слайда из GigaCHAT===========================================================
        # Словарь с API ключами


        try:

            dict_giga_reliz = {"ura_giga_api": ["4e41d892-18ed-427a-83c8-cd76cc5df693", "YWNhMzE5ZWUtYzdjOC00MjE0LTg0MWItMmU3ZGQ2NDQyYjcwOmJkYzY4OWY5LThkZDMtNDA2ZC05MGNjLTFmMTM3ODBjYzg3ZA=="],
                               "rafail_giga_api": ["f3e39c44-940b-41bc-bc48-94705025d7de", "ZTQ5MTk2YzktYmNmYi00YjIwLThiNmEtMWFjMGZjZjljYWRlOjkzZTYwYmJiLTA1MmUtNGEyNy1iMzY4LWM2MzJhYTk3ZGEyZQ=="],
                               "dima_giga_api": ["8c09c32a-c57a-4f7b-b396-855c12a2bdee", "NWU5ZWQ3YTgtN2UxNi00M2I4LTkwNGUtNDA2Y2QxNjBkZDUwOmQ2ODM5ZWVjLTNjMjctNDM3Ni05MTZjLWNjYzFkY2EyNzZhMw=="],
                               "koly_giga_api": ["3f5d31ea-237f-43b9-b96a-0e49518cba18", "NGUzNGUyZGYtYmM5Ni00NTM5LTljMWYtMmFmOTJjZjIyMjVmOmI2YTQwY2E3LWIzM2ItNDM5Yi05Yjk3LTlhYTBkMzk5NmIyMA== "],
                               "igor_giga_api": ["34c3aeb9-4211-4a11-ad49-2a151b3ca0c4", "OWRhNjc1ZTEtYmVkMS00MTk3LTkwYzgtNzE4Y2M2ZTYxNTNiOmNiZmIyOTg4LThjZGYtNGYxNy04NWFhLTA0MjYzOTVkOWI1Ng=="]}

            # Выбираем случайную пару значений
            random_key = random.choice(list(dict_giga_reliz.keys()))
            rq_uid_value = dict_giga_reliz[random_key][0]
            authorization_value = dict_giga_reliz[random_key][1]


            url = "https://sm-auth-sd.prom-88-89-apps.ocp-geo.ocp.sigma.sbrf.ru/api/v2/oauth"
            payload = {
                'scope': 'GIGACHAT_API_PERS'
            }
            headers = {
                'Content-Type': 'application/x-www-form-urlencoded',
                'Accept': 'application/json',
                'RqUID': f'{rq_uid_value}',
                'Authorization': f'Basic {authorization_value}'
            }

            response = requests.request("POST", url, headers=headers, data=payload, verify=False)

            if response.status_code != 200:
                # Выводим сообщение об ошибке, но НЕ прерываем выполнение
                error_msg = f">>>>>>>>>>>>>>>>>>>>>>>>>>>>>>> Закончились токены у пользователя: {random_key}.<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<"
                print(error_msg)  # Вывод в консоль
                # Пропускаем текущую итерацию или используем fallback-логику
                final_answer = ""

            else:
                response_data_keys = response.json()
            # print(response_data_keys['access_token'])


            url = "https://gigachat.devices.sberbank.ru/api/v1/chat/completions"

            books_all_treb = []
            for i, (key, values) in enumerate(result_my_dict.items(), start=1):
                lm_otklon = ''.join(values[0])
                uroven = ''.join(values[3])
                if uroven == '3':
                    books_all_treb.append(lm_otklon)


            # Этап 1: Анализ отклонений
            step1_payload = {
                "model": "GigaChat",
                "messages": [
                    {
                        "role": "system",
                        "content": f"Ты эксперт по ITSM процессам, поэтому тщательно проанализируй каждое отклонение из списка: {books_all_treb}."

                    }
                ]
            }
            step1_headers = {
                'Content-Type': 'application/json',
                'Accept': 'application/json',
                'Authorization': f'Bearer {response_data_keys["access_token"]}'
            }
            step1_response = requests.post(url, headers=step1_headers, data=json.dumps(step1_payload), verify=False)
            step1_data = step1_response.json()
            analysis_results = step1_data['choices'][0]['message']['content'].strip()

            # Этап 2: Классификация отклонений
            step2_payload = {
                "model": "GigaChat",
                "messages": [
                    {
                        "role": "system",
                        "content": f"Основываясь на анализе отклонений ({analysis_results}), классифицируй отклонения по важности выполнения."

                    }
                    ]
                }
            step2_response = requests.post(url, headers=step1_headers, data=json.dumps(step2_payload), verify=False)
            step2_data = step2_response.json()
            classification_results = step2_data['choices'][0]['message']['content'].strip()

            # Этап 3: Выбор критичных отклонений
            step3_payload = {
                "model": "GigaChat",
                "messages": [
                    {
                        "role": "system",
                        "content": f"Выпиши по одному самому критическому отклонению для каждого процесса из списка({classification_results}). "
                                   f"Критерии отбора:\n"
                                   f"- нарушение базовой функции процесса;\n"
                                   f"- делает процесс нефункциональным;\n"
                                   f"- критически влияет на безопасность/доступность ИТ-услуг.\n"
                    }
                ]
            }
            step3_response = requests.post(url, headers=step1_headers, data=json.dumps(step3_payload), verify=False)
            step3_data = step3_response.json()
            critical_deviations = step3_data['choices'][0]['message']['content'].strip()

            # Этап 4: Формирование итоговых утверждений
            step4_payload = {
                "model": "GigaChat",
                "messages": [
                    {
                        "role": "system",
                        "content": f"""
                        Напечатайте список из ровно 5 критических отклонений, выбранных из списка {critical_deviations}, придерживаясь жёстких правил:

                        
                        Правила оформления:
                        1. Никаких предварительных пояснений, заголовков ("Критические отклонения:", "Список критических отклонений:" и др.).
                        2. Только чистые предложения, перечисляемые через "; ".
                        3. Каждое предложение начинать с заглавной буквы.
                        4. Без маркеров (-, •, *) и цифровых обозначений (1., 2.).
                        5. Минимальная длина каждого отклонения — одно короткое предложение.
                        
                        Пример:  
                        Нет учета рисков; Нет оценки показателей качества услуг; Нет резервного копирования; Нет автоматизации поддержки пользователей; Нет утвержденных стандартов  
                        
                        Если формат нарушен — ответ не засчитывается!
                        """
                            # f"Сформируй 5 независимых отклонений из {critical_deviations}, без описания  и пояснениячтобы каждое отклонение описывало одно фундаментальное нарушение в одном из процессов."
                    }
                ]
            }
            step4_response = requests.post(url, headers=step1_headers, data=json.dumps(step4_payload), verify=False)
            step4_data = step4_response.json()
            final_answer = step4_data['choices'][0]['message']['content'].strip()

            # Пост-обработка для удаления начального дефиса (с пробелом или без)
            if final_answer.startswith("- "):  # Сначала проверяем дефис с пробелом
                final_answer = final_answer[2:]
            elif final_answer.startswith("-"):  # Затем проверяем одиночный дефис
                final_answer = final_answer[1:]

            # print(final_answer)

            # Очищаем от других нежелательных символов
            final_answer = (
                final_answer.replace('"', '')
                    .replace('•', '')
                    .replace('№', '')
                    .replace('#', '')
                    .strip()  # Убираем пробелы по краям
            )
            final_answer = re.sub(r'\([^)]*\)', '', final_answer)  # Удаление всего в скобках
            final_answer = '; '.join([s.strip().capitalize() for s in final_answer.split(';')])

        except Exception as e:
            final_answer = ""


        # jjj=input('++++++++')
        # six_workshet['P2'] = output_string
        # seven_workshet['P2'] = output_string

        #============================================= !!!!!!КОНЕЦ Генерации из GigaCHAT!!!!!!!!!!!!!===========================================================


        if os.path.exists(save):
            k = f'{name_company}'
            print(f"Файл '{k}' уже существует.")
            # print(f"Файл '{save}'\Отчет_по_ITSM_AS.xlsx' уже существует.")
            # Запрашиваем у пользователя, хочет ли он перезаписать файл
            while True:
                perazapis = input("Вы хотите перезаписать файл Отчет_по_ITSM_AS.xlsx? (y/n): ")
                if perazapis == 'n':
                    print(f"Файл '{k} не перезаписан.")
                    break
                elif perazapis == 'y':
                    wb.save(save)
                    print(f"Файл '{k} перезаписан.")
                    print('=====  Отчет сформирован  =====')
                    break
                else:
                    print("Неверный ответ. Пожалуйста, введите 'y' для перезаписи или 'n' для отказа.")
        else:
            wb.save(save)
            print('=====  Отчет сформирован  =====')




        #============================================ Начало Генерация слайда ======================================================================

        print('====================================================================================================')
        print('====================================================================================================')
        print('******************************************')
        print('Генерация слайда из ТЕПЛОВОЙ КАРТЫ')
        print('******************************************')
        #=======================================================================================================================

        def apply_color_to_ppt(text_run, percent):
            """Задаёт цвет текста в PowerPoint на основе процента."""
            r = int(255 * (100 - percent) / 100)
            g = int(255 * percent / 100)
            text_run.font.color.rgb = RGBColor(r, g, 0)
            text_run.font.bold = True

        ws_two = Worksheet()
        book_opros = {}
        dict_assesment = {}
        script_dir_all = os.path.abspath(os.path.dirname(__file__))
        mas_files_list = []
        keywords = ["Объединяющие элементы"]

        file_list = glob.glob(os.path.join(script_dir_all, '*.xlsx'))
        for file in file_list:
            if any(keyword in file for keyword in keywords):
                mas_files_list.append(file)

        oprosnik_files = [file for file in mas_files_list if 'Объединяющие элементы' in file]
        opros = os.path.join(os.path.dirname(__file__),  *oprosnik_files)

        print()
        print('-------- Формирую слайд --------')
        print()
        warnings.simplefilter("ignore")

        ws_two.xlsx_to_dict(path=opros,  select_sheet='Объединяющие элементы', data_only=True)
        #TODO: Собираем информацию с Опросника согласно полям!
        for b in ws_two.sheet_items:
            juice_one = b.setdefault('Компания')
            juice_two = b.setdefault('Блок / Индустрия (сокращенно)')
            # juice_three = b.setdefault('Общий уровень надежности')
            # juice_four = b.setdefault('Уровень технологической надежности АС')
            # juice_five = b.setdefault('Управление технологическими инцидентами')
            # juice_six = b.setdefault('Управление каталогом сервисов')
            # juice_seven = b.setdefault('Управление уровнями сервисов')
            # juice_eight = b.setdefault('Управление технологическими проблемами')
            # juice_nine = b.setdefault('Управление технологическими рисками')
            # juice_ten = b.setdefault('Управление инфраструктурными изменениями')
            # juice_eleven = b.setdefault('Управление непрерывностью технологий')
            # juice_twelve = b.setdefault('Управление доступностью технологий')
            # juice_thirteen = b.setdefault('Управление техн. событиями и мониторингом')
            # juice_fourteen = b.setdefault('Управление технологическими ЗНО')
            # juice_fifteen = b.setdefault('Управление конфигурациями')
            # juice_sixteen = b.setdefault('Управление обновлением ПО')
            juice_seventeen = b.setdefault('Подключение к SberRadar (%)')
            juice_eighteen = b.setdefault('Передача на опердашбоард Sber911 (%)')
            juice_nineteen = b.setdefault('Вкючение авторегистрации')
            juice_twenty = b.setdefault('Подключение ОЭ Sber911 (%)')
            juice_twentyone = b.setdefault('Подключение ОЭ SberITSM (%)')
            juice_twentytwo = b.setdefault('Исполнение поручений Комиссии по инцидентам')
            juice_twentythree = b.setdefault('Статус по соглашению надежности')
            # juice_twentyfour = b.setdefault('Ключевые проблемы')
            # juice_twentyfive = b.setdefault('Динамические показатели ITSM')
            # juice_twentysix = b.setdefault('Динамические показатели по технологиям')
            # juice_twentyseven = b.setdefault('Уровень технологической надежности')
            juice_twentyeight = b.setdefault('Количество ВПИ за квартал')



            book_opros.setdefault(juice_one, [])
            book_opros[juice_one].append(juice_two)
            # book_opros[juice_one].append(juice_three)
            # book_opros[juice_one].append(juice_four)
            # book_opros[juice_one].append(juice_five)
            # book_opros[juice_one].append(juice_six)
            # book_opros[juice_one].append(juice_seven)
            # book_opros[juice_one].append(juice_eight)
            # book_opros[juice_one].append(juice_nine)
            # book_opros[juice_one].append(juice_ten)
            # book_opros[juice_one].append(juice_eleven)
            # book_opros[juice_one].append(juice_twelve)
            # book_opros[juice_one].append(juice_thirteen)
            # book_opros[juice_one].append(juice_fourteen)
            # book_opros[juice_one].append(juice_fifteen)
            # book_opros[juice_one].append(juice_sixteen)
            book_opros[juice_one].append(juice_seventeen)
            book_opros[juice_one].append(juice_eighteen)
            book_opros[juice_one].append(juice_nineteen)
            book_opros[juice_one].append(juice_twenty)
            book_opros[juice_one].append(juice_twentyone)
            book_opros[juice_one].append(juice_twentytwo)
            book_opros[juice_one].append(juice_twentythree)
            # book_opros[juice_one].append(juice_twentyfour)
            # book_opros[juice_one].append(juice_twentyfive)
            # book_opros[juice_one].append(juice_twentysix)
            # book_opros[juice_one].append(juice_twentyseven)
            book_opros[juice_one].append(juice_twentyeight)

        # print(book_opros)

        if six_workshet['A2'].value in book_opros:
            print('Компания соответствует названию в файле ====== ОБЪЕДИНЯЩИЕ ЭЛЕМЕНТЫ ========')
            # print(book_opros[six_workshet['A2'].value])
            #Находим значение в ячейке C2 и переводим в float, т.к. было число в ячейке 2,5 и получаем 2.5
            def to_float(vr):
                if isinstance(vr, str):
                    return float(vr.replace(',', '.'))
                return float(vr)  # сработает для int, float, numpy.float64 и т.д.
            vals_c = to_float(six_workshet['C2'].value)

            dict_assesment.setdefault(six_workshet['A2'].value, [])
            dict_assesment[six_workshet['A2'].value].append(book_opros[six_workshet['A2'].value][0])
            dict_assesment[six_workshet['A2'].value].append(str(six_workshet['B2'].value))
            dict_assesment[six_workshet['A2'].value].append(str(vals_c))
            dict_assesment[six_workshet['A2'].value].append(six_workshet['D2'].value)
            dict_assesment[six_workshet['A2'].value].append(six_workshet['E2'].value)
            dict_assesment[six_workshet['A2'].value].append(six_workshet['F2'].value)
            dict_assesment[six_workshet['A2'].value].append(six_workshet['G2'].value)
            dict_assesment[six_workshet['A2'].value].append(six_workshet['H2'].value)
            dict_assesment[six_workshet['A2'].value].append(six_workshet['I2'].value)
            dict_assesment[six_workshet['A2'].value].append(six_workshet['J2'].value)
            dict_assesment[six_workshet['A2'].value].append(six_workshet['K2'].value)
            dict_assesment[six_workshet['A2'].value].append(six_workshet['L2'].value)
            dict_assesment[six_workshet['A2'].value].append(six_workshet['M2'].value)
            dict_assesment[six_workshet['A2'].value].append(six_workshet['N2'].value)
            dict_assesment[six_workshet['A2'].value].append(six_workshet['O2'].value)
            dict_assesment[six_workshet['A2'].value].append(book_opros[six_workshet['A2'].value][1])
            dict_assesment[six_workshet['A2'].value].append(book_opros[six_workshet['A2'].value][2])
            dict_assesment[six_workshet['A2'].value].append(book_opros[six_workshet['A2'].value][3])
            dict_assesment[six_workshet['A2'].value].append(book_opros[six_workshet['A2'].value][4])
            dict_assesment[six_workshet['A2'].value].append(book_opros[six_workshet['A2'].value][5])
            dict_assesment[six_workshet['A2'].value].append(book_opros[six_workshet['A2'].value][6])
            dict_assesment[six_workshet['A2'].value].append(book_opros[six_workshet['A2'].value][7])


            if final_answer == '' or final_answer is None:
                dict_assesment[six_workshet['A2'].value].append("Проблема 1; Проблема 2; Проблема 3; Проблема 4; Проблема 5")
            else:
                dict_assesment[six_workshet['A2'].value].append(final_answer)
            dict_assesment[six_workshet['A2'].value].append(six_workshet['Q2'].value)
            dict_assesment[six_workshet['A2'].value].append(six_workshet['R2'].value)
            dict_assesment[six_workshet['A2'].value].append(six_workshet['S2'].value)
            dict_assesment[six_workshet['A2'].value].append(book_opros[six_workshet['A2'].value][8])


        else:
            print('Компания не соответствует названию в файле ====== ОБЪЕДИНЯЩИЕ ЭЛЕМЕНТЫ ========')
            sys.exit(0)

            for key in dict_assesment:
                if isinstance(dict_assesment[key], list):  # Проверяем, что значение — список
                    dict_assesment[key] = [
                        item.replace('\n', '') if isinstance(item, str) else item
                        for item in dict_assesment[key]
                    ]


        # print(dict_assesment)
        # for key in book_opros:
        #     if isinstance(book_opros[key], list):  # Проверяем, что значение — список
        #         book_opros[key] = [
        #             item.replace('\n', '') if isinstance(item, str) else item
        #             for item in book_opros[key]
        #         ]


        #=============================================СОЗДАНИЕ СЛАЙДА===========================================================


        dict_blocks = {"EdTech": "Индустрия EdTech",
                       "E-Health": "Индустрия E-Health",
                       "GRC": "Блок «GR, правовые вопросы, комплаенс и ДЗО»",
                       "Media & Ads": "Индустрия «Media & Ads»",
                       "SBI": "Блок «Sberbank International»",
                       "Виртуальное": "«Виртуальное ДЗО»",
                       "ДДиРС B2C": "Департамент данных и рекомендательных систем B2C",
                       "ДЗиС": "Департамент «Занять и сберегать»",
                       "ДЗОК": "Департамент «Заботы о клиентах»",
                       "ДРПА": "Департамент по работе с проблемными активами",
                       "КИБ": "Блок «КИБ»",
                       "ЛиК": "Блок «Люди и Культура»",
                       "Подразделения вне блоков B2C": "Подразделения вне блоков B2C",
                       "Развитие клиентского опыта B2C": "Блок «Развитие клиентского опыта B2C»",
                       "Риски": "Блок «Риски»",
                       "Сервисы": "Блок «Сервисы»",
                       "Сеть продаж": "Блок «Сеть продаж»",
                       "Строительство": "Блок «Строительство»",
                       "Технологии": "Блок «Технологии»",
                       "ТР": "Блок «Технологическое развитие»",
                       "Транзакционный банкинг B2C": "Блок «Транзакционный банкинг B2C»",
                       "УБ": "Блок «Управление благосостоянием»",
                       "ЦПНБ": "Центр Перспективных Направлений Бизнеса"
                       }


        # proverka_vyborki = input('------------ТЕСТ------------ ')

        for key, value in dict_assesment.items():
            problems_text = value[22]
            # Проверяем тип данных
            if isinstance(problems_text, str):
                # Если это строка - разделяем по точкам с запятой
                problems_list = problems_text.split(';')
                problems_list = [problem.strip() for problem in problems_list]
            elif isinstance(problems_text, list):
                # Если это уже список - просто используем его
                problems_list = problems_text
            else:
                # Если ни строка, ни список - создаем пустой список
                problems_list = []

            # print(problems_list)  # Для отладки

            element = value[23]  # '121 / 121, 117 / 117, 26 / 26'
            parts = element.split(',')  # Делим строку по запятым

            # Проверяем, что частей достаточно
            if len(parts) >= 3:
                result_kritik_itsm = parts[0].strip()  # '121 / 121'
                result_visok_itsm = parts[1].strip()  # '117 / 117'
                result_umeren_itsm = parts[2].strip()  # '26 / 26'
            else:
                result_kritik_itsm = "0 / 0"
                result_visok_itsm = "0 / 0"
                result_umeren_itsm = "0 / 0"


            element_techno = value[24]
            parts_techno = element_techno.split(',')  # Делим строку по запятым

            # Проверяем, что частей достаточно
            if len(parts_techno) >= 3:
                result_kritik_as = parts_techno[0].strip()  # '121 / 121'
                result_visok_as = parts_techno[1].strip()  # '117 / 117'
                result_umeren_as = parts_techno[2].strip()  # '26 / 26'
            else:
                result_kritik_as = "0 / 0"
                result_visok_as = "0 / 0"
                result_umeren_as = "0 / 0"


            # Получаем значения для замены из словаря
            # Функция для обработки всех текстовых элементов
            def replace_text_in_shape(shape, data):
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        for key, value in data.items():
                            pattern = rf'\b({re.escape(key)})\b'
                            matches = re.findall(pattern, paragraph.text)
                            if matches:
                                for run in paragraph.runs:
                                    if key in run.text:
                                        try:
                                            # Сохраняем оригинальное форматирование
                                            original_font_name = run.font.name
                                            original_font_size = run.font.size
                                            original_font_color = run.font.color

                                            # Заменяем текст
                                            run.text = re.sub(pattern, value, run.text)
                                            # Восстанавливаем форматирование
                                            run.font.name = original_font_name
                                            run.font.size = original_font_size
                                            if original_font_color and original_font_color.type == 1:  # MSO_COLOR_TYPE.RGB
                                                run.font.color.rgb = original_font_color.rgb
                                        except Exception as e:
                                            print(f'Ошибка при замене текста: {e}')

                # Обработка таблиц
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            replace_text_in_cell(cell, data)
            # Специальная функция для работы с ячейками таблицы
            def replace_text_in_cell(cell, data):
                if cell.text_frame:
                    text_frame = cell.text_frame
                    for paragraph in text_frame.paragraphs:
                        for key, value in data.items():
                            pattern = rf'\b({re.escape(key)})\b'
                            matches = re.findall(pattern, paragraph.text)
                            if matches:
                                for run in paragraph.runs:
                                    if key in run.text:
                                        try:
                                            # Сохраняем оригинальное форматирование
                                            original_font_name = run.font.name
                                            original_font_size = run.font.size
                                            original_font_color = run.font.color

                                            # Заменяем текст
                                            run.text = re.sub(pattern, value, run.text)

                                            # Восстанавливаем форматирование
                                            run.font.name = original_font_name
                                            run.font.size = original_font_size
                                            if original_font_color and original_font_color.type == 1:  # MSO_COLOR_TYPE.RGB
                                                run.font.color.rgb = original_font_color.rgb
                                        except Exception as e:
                                            print(f'Ошибка при замене текста: {e}')
            # Открываем существующий шаблон PowerPoint



            open_prezent_slide = os.path.join(os.path.dirname(__file__), 'no_hand', 'template.pptx')
            presentation = Presentation(open_prezent_slide)

            # Определяем данные для вставки


            data = {
                'name': 'name_dzo',
                'subtitle ': 'name_dzo_two',
                'data_as': 'text_dzo',
                'num': 'all_number',
                'assesment':'Оценка',
                'itas': 'Оценка всех ас',
                'block':'Блок',
                'itsmkall':'Критичные отклонения по ITSM',
                'itsmuall':'Умеренные отклонения по ITSM',
                'itshight':'Высокие отклонения по ITSM',
                'askall':'Критичные отклонения по АС',
                'asvall':'Высокие отклонения  по АС',
                'asuall':'Умеренные отклонения по АС',
                'answer1':'Критичное отклонение 1',
                'answer2':'Критичное отклонение 2',
                'answer3':'Критичное отклонение 3',
                'answer4':'Критичное отклонение 4',
                'answer5':'Критичное отклонение 5',
                'q_ozenka':'Оценка',
                'est_profile':'Profile',
                'sberradar': 'Подключение к SberRadar',
                'sber911': 'Подключение к опердашбоард Sber911',
                'raschet1':'Расчетная мощность АС',
                'oesber911': 'Подключение ОЭ Sber911',
                'sberitsm':'Подключение ОЭ SberITSM',
                'sogdtn': 'Статус по соглашению надежности',
                'poruchenie': 'Исполнение поручений Комиссии по инцидентам',
                'inzidenton':'Вкючение авторегистрации инцидентов',
                'vpi': 'Количество ВПИ за квартал'
            }

            data['name'] = key
            if value[0] in dict_blocks:
                data['block'] = dict_blocks[value[0]]
            else:
                data['block'] = value[0]

            # for id_block, name_block in dict_blocks.items():
            #     print(id_block,'=========',value[0])
            #
            #     if id_block == value[0]:
            #         data['block'] = name_block
            #     else:
            #         data['block'] = value[0]

            data['subtitle'] = key
            data['itsmkall'] = result_kritik_itsm
            data['itshight'] = result_visok_itsm
            data['itsmuall'] = result_umeren_itsm

            # # Получаем текущую дату
            # current_date = datetime.date.today()
            # # Определяем название папки на основе текущей даты
            # quarter = (current_date.month - 1) // 3 + 1
            # year = current_date.year
            # quarter_folder = f"Q{quarter} {year}"

            # =============Поиск и удаление картинки на презентации если количество отклонений будет меньше или вовсе не будет=================
            # Открываем презентацию
            # Перебираем все слайды


            # Путь к изображению
            NO_HAND_IMAGE = os.path.abspath(os.path.join('no_hand', 'picture.png'))

            # Проверка, что файл существует
            if not os.path.exists(NO_HAND_IMAGE):
                raise FileNotFoundError(f"Файл не найден: {NO_HAND_IMAGE}")

            # Словарь с метками и путями к изображениям
            data_image = {
                'img1': NO_HAND_IMAGE,
                'img2': NO_HAND_IMAGE,
                'img3': NO_HAND_IMAGE,
                'img4': NO_HAND_IMAGE,
                'img5': NO_HAND_IMAGE
            }

            slide_zero = presentation.slides[0]

            # Проходим по всем фигурам на слайде
            for shape in slide_zero.shapes:
                if shape.has_table:
                    table = shape.table
                    table_left = shape.left
                    table_top = shape.top

                    # Сначала обрабатываем все проблемы
                    problems_list = [p.strip() for p in problems_text.split(';') if p.strip()]

                    # Создаем словарь только для существующих проблем
                    data_image = {}
                    for i in range(1, 6):
                        if i <= len(problems_list):
                            data[f'answer{i}'] = problems_list[i-1]
                            data_image[f'img{i}'] = NO_HAND_IMAGE

                        else:
                            data[f'answer{i}'] = ""

                    # Теперь проходим по ячейкам таблицы
                    for row_idx, row in enumerate(table.rows):
                        for col_idx, cell in enumerate(row.cells):
                            cell_text = cell.text.strip()

                            # Вставляем изображение только если есть соответствующая проблема
                            if cell_text in data_image:
                                cell_left = table_left + sum(table.columns[i].width for i in range(col_idx))
                                cell_top = table_top + sum(table.rows[i].height for i in range(row_idx))
                                cell_width = table.columns[col_idx].width
                                cell_height = table.rows[row_idx].height

                                cell.text = ""  # Очищаем ячейку
                                cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Центрируем содержимое


                                # Открываем изображение и получаем его размеры
                                # Размер изображения — 50% от ширины ячейки

                                img_width = cell_width * 0.5
                                # Рассчитываем высоту изображения, сохраняя пропорции

                                with Image.open(data_image[cell_text]) as img:
                                    original_width, original_height = img.size
                                    aspect_ratio = original_height / original_width
                                    img_height = img_width * aspect_ratio

                                # Добавляем изображение по центру
                                pic = slide_zero.shapes.add_picture(
                                    data_image[cell_text],
                                    cell_left + (cell_width - img_width) / 2,  # Центр по X
                                    cell_top + (cell_height - img_height) / 2,  # Центр по Y
                                    width=img_width,
                                    height=img_height
                                )
                            elif cell_text.startswith('img') and cell_text not in data_image:
                                # Очищаем ячейки, где должны быть изображения, но проблем нет
                                cell.text = ""
                                # Устанавливаем вертикальное выравнивание по центру для ячейки

            #*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*

            #*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*

            #=================================Прямоугольник 34(Изменение динамики по Уровням)====================================

            def apply_color_to_text(shape, target_texts, colors):
                if not shape.has_text_frame:
                    return

                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text = run.text.strip().lower()
                        for i, target_text in enumerate(target_texts):
                            if target_text.lower() in text.replace(" ", "").replace("\n", "").replace("\t", ""):
                                # Создаём новый цветовой объект и присваиваем его
                                run.font.color.rgb = colors[i]
                                # print(f"Цвет изменён для текста '{target_text}' в фигуре '{shape.name}'.")

            # Целевые тексты и их цвета
            target_texts = ["assessment", "num"]
            # Преобразуем значение в float один раз
            try:
                value_float = float(value[1])
            except (ValueError, IndexError):
                value_float = 0.0  # Значение по умолчанию при ошибке

            # Инициализируем значения по умолчанию (желтый)
            data['num'] = f"{value_float:.2f}"  # Форматируем с двумя знаками после запятой
            data['assessment'] = '«Оцененный» -'
            colors = [RGBColor(238, 237, 25), RGBColor(238, 237, 25)]
            # Проверяем условия в правильном порядке (от строгих к более общим)
            if value[1] == '0':
                data['num'] = '0.00'
                data['assessment'] = '«Черный ящик» -'
                colors = [RGBColor(227, 18, 19), RGBColor(227, 18, 19)]  # Красный
            elif value_float < 2:
                data['assessment'] = '«Декларативный» -'
                colors = [RGBColor(233, 178, 26), RGBColor(233, 178, 26)]  # Оранжевый
            elif 2 <= value_float < 3:
                data['assessment'] = '«Оцененный» -'  # Оставляем по умолчанию
                colors = [RGBColor(238, 237, 25), RGBColor(238, 237, 25)]  # Желтый
            elif 3 <= value_float < 4:
                data['assessment'] = '«Готовность к эксплуатации» -'
                colors = [RGBColor(130, 187, 72), RGBColor(130, 187, 72)]  # Зеленый
            elif value_float == 4:
                data['assessment'] = '«Готовность к развитию» -'
                colors = [RGBColor(17, 160, 81), RGBColor(17, 160, 81)]  # Темно-зеленый



            # Применяем изменения
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.name == "Прямоугольник 34":
                        # print(f"Найдена фигура: {shape.name}")
                        apply_color_to_text(shape, target_texts, colors)




            #=============================================================================

            # Проверяем, что value[2] существует и непустой


            if len(value) > 2:
                str_value = value[2].strip()  # Удаляем пробелы

                if not str_value:  # Если строка пустая после strip()
                    # print(key)
                    data['itas'] = ""
                else:
                    try:
                        itas_value = float(str_value)
                        data['itas'] = f"{itas_value:.2f}"
                    except ValueError:
                        data['itas'] = "0.00"
            else:
                data['itas'] = "0.00"  # Если value слишком короткий

            data['askall'] = result_kritik_as
            data['asvall'] = result_visok_as
            data['asuall'] = result_umeren_as
            # data['data_as'] = data_assesment
            # data['q_ozenka'] = quarter_folder
            if value[25]:
                data['est_profile'] = '(' + str(value[25]) + ')'
            else:
                data['est_profile'] = value[25]
            data['sogdtn'] = value[21]



            if value[15]:  # Проверяем, что value[15] не пустая строка
                try:
                    data['sberradar'] = str(round(float(value[15]))) + '%'
                except ValueError:
                    data['sberradar'] = value[15]
            else:
                data['sberradar'] = '0%'


            if value[16]:  # Проверяем, что value[16] не пустая строка
                try:
                    data['sber911'] = str(round(float(value[16]))) + '%'
                except ValueError:
                    data['sber911'] = value[16]
            else:
                data['sber911'] = '0%'


            data['inzidenton'] = value[17]


            if value[18]:  # Проверяем, что value[18] не пустая строка
                try:
                    data['oesber911'] = str(round(float(value[18]))) + '%'
                except ValueError:
                    data['oesber911'] = value[18]
            else:
                data['oesber911'] = '0%'


            if value[19]:  # Проверяем, что value[19] не пустая строка
                try:
                    data['sberitsm'] = str(value[19])
                except ValueError:
                    data['sberitsm'] = value[19]
            else:
                data['sberitsm'] = '0'


            if value[20]:  # Проверяем, что value[20] не пустая строка
                try:
                    data['poruchenie'] = str(round(float(value[20]))) + '%'
                except ValueError:
                    data['poruchenie'] = value[20]
            else:
                data['poruchenie'] = '0%'



            if value[26]: # Проверяем, что value[26] не пустая строка
                try:
                    data['vpi'] = (value[26])
                except ValueError:
                    data['vpi'] = value[26]
            else:
                data['vpi'] = '0'




            if value[1] == '0':
                data['num'] = '0.00'
                data['assessment'] = '«Черный ящик» -'
                colors = [RGBColor(227, 18, 19), RGBColor(227, 18, 19)]  # Красный





            # Находим первый слайд
            slide = presentation.slides[0]

            # Обрабатываем все фигуры на слайде
            for shape in slide.shapes:
                replace_text_in_shape(shape, data)

            # Исходный словарь
            dict_old = {
                'Управление технологическими инцидентами': '0%',
                'Управление каталогом сервисов': '0%',
                'Управление уровнями сервисов': '0%',
                'Управление технологическими проблемами': '0%',
                'Управление технологическими рисками': '0%',
                'Управление инфраструктурными изменениями': '0%',
                'Управление непрерывностью технологий': '0%',
                'Управление доступностью технологий': '0%',
                'Управление техн. событиями и мониторингом': '0%',
                'Управление ТЗНО': '0%',
                'Управление конфигурациями': '0%',
                'Управление обновлением ПО': '0%'
            }

            dict_old['Управление технологическими инцидентами'] = value[3]
            dict_old['Управление каталогом сервисов'] = value[4]
            dict_old['Управление уровнями сервисов'] = value[5]
            dict_old['Управление технологическими проблемами'] = value[6]
            dict_old['Управление технологическими рисками'] = value[7]
            dict_old['Управление инфраструктурными изменениями'] = value[8]
            dict_old['Управление непрерывностью технологий'] = value[9]
            dict_old['Управление доступностью технологий'] = value[10]
            dict_old['Управление техн. событиями и мониторингом'] = value[11]
            dict_old['Управление ТЗНО'] = value[12]
            dict_old['Управление конфигурациями'] = value[13]
            dict_old['Управление обновлением ПО'] = value[14]

            # Фильтруем словарь, удаляя категории с отсутствующими данными
            filtered_dict = {k: v for k, v in dict_old.items() if v != '-'}

            # Поиск первой диаграммы на слайде
            chart_shape = None
            for shape in slide.shapes:
                if shape.has_chart:
                    chart_shape = shape
                    chart = shape.chart
                    break

            if chart:
                # Создаем новый объект ChartData
                # Настройка размеров диаграммы перед обновлением данных
                chart_data = ChartData()

                # Извлекаем категории из отфильтрованного словаря
                categories = list(filtered_dict.keys())
                values = []

                # Функция для безопасного преобразования значения в число
                def safe_convert_to_decimal(value):
                    if value is None:
                        return 0  # Возвращаем 0, если значение None
                    if not isinstance(value, str):
                        return 0  # Возвращаем 0, если значение не строка
                    try:
                        # Пытаемся преобразовать проценты в десятичные числа
                        return float(value.strip('%')) / 100
                    except ValueError:
                        return 0  # Возвращаем 0, если преобразование не удалось

                # Преобразуем значения в числовой формат
                for category in categories:
                    valurt = filtered_dict[category]
                    valurt = safe_convert_to_decimal(valurt)  # Безопасное преобразование
                    values.append(valurt)

                # Теперь categories и values готовы для использования в ChartData
                # print("Категории:", categories)
                # print("Значения:", values)
                # print(f'Слайд компании {key} - ГОТОВ!')

                # Заполняем данные для новой диаграммы
                chart_data.categories = categories
                chart_data.add_series('', values)  # Название серии можно изменить

                # Применяем новые данные к диаграмме
                chart.replace_data(chart_data)

                # Настраиваем формат меток данных для отображения в процентах
                for series in chart.series:
                    series.data_labels.show_value = True  # Показываем значения на диаграмме
                    series.data_labels.number_format = '0%'  # Устанавливаем формат процентов

                # Вычисляем, насколько нужно сжать диаграмму
                original_height = chart_shape.height  # Исходная высота диаграммы
                original_categories_count = len(dict_old)  # Исходное количество категорий
                filtered_categories_count = len(filtered_dict)  # Количество категорий после фильтрации

                # Если количество категорий уменьшилось, сжимаем диаграмму
                if filtered_categories_count < original_categories_count:
                    # Вычисляем коэффициент сжатия
                    compression_ratio = filtered_categories_count / original_categories_count
                    # Уменьшаем высоту диаграммы
                    new_height = int(original_height * compression_ratio)
                    chart_shape.height = new_height



            else:
                print("Диаграмма не найдена.")



            def safe_format_value(value):
                """Безопасно форматирует значение с цветом, возвращая всегда кортеж (text, color)"""
                try:
                    # Обработка None и пустых значений
                    if value is None or str(value).strip() == '':
                        return "-", RGBColor(128, 128, 128)  # Серый для пустых значений

                    str_val = str(value).strip()


                    # Специальные текстовые случаи
                    special_cases = {
                        'исключен': ('Исключен', RGBColor(0, 176, 80)),  # Зеленый
                        'включено': ('Включено', RGBColor(0, 176, 80)),  # Зеленый
                        'подписано': ('Подписано', RGBColor(0, 176, 80)),  # Зеленый
                        'да': ('Да', RGBColor(0, 176, 80)),
                        'yes': ('Yes', RGBColor(0, 176, 80)),
                        'нет': ('Нет', RGBColor(255, 0, 0)),  # Красный
                        'не подписано': ('Не подписано', RGBColor(255, 0, 0)),  # Красный
                        'отказ': ('Отказ', RGBColor(255, 0, 0)),  # Красный
                        'не подключен': ('Не подключен', RGBColor(255, 0, 0)),  # Красный
                        'отсутствует': ('Отсутствует', RGBColor(255, 0, 0)),  # Красный
                        'no': ('No', RGBColor(255, 0, 0)),
                        'n/a': ('N/A', RGBColor(128, 128, 128)),  # Серый
                        'долгое подписание': ('Долгое подписание', RGBColor(255, 255, 0)),  # Желтый
                        'в процессе': ('В процессе', RGBColor(255, 255, 0)),  # Желтый
                        '«оцененный»': ('«Оцененный»', RGBColor(255, 255, 0)),  # Желтый
                        'в работе': ('В работе', RGBColor(255, 255, 0)),  # Желтый



                        # 'aaa': ('aaa', RGBColor(255, 0, 0)),  # Красный

                        '-': ('-', RGBColor(128, 128, 128)),  # Серый
                    }

                    lower_val = str_val.lower()
                    if lower_val in special_cases:
                        return special_cases[lower_val]

                    # Обработка процентов
                    if '%' in str_val:
                        num_str = str_val.replace('%', '')
                    else:
                        num_str = str_val


                    num = float(num_str)
                    num = max(0, min(100, num))  # Ограничиваем 0-100

                    # Градиент цветов
                    if num <= 50:
                        # От красного (0%) до желтого (50%)
                        red = 255
                        green = int(255 * (num / 50))
                        blue = 0
                    else:
                        # От желтого (50%) до зеленого (100%)
                        red = int(255 * (1 - (num - 50) / 50))
                        # green = int(255 - (255 - 176) * ((num - 50) / 50))  # Корректируем формулу для green
                        green = int(round(255 - (255 - 176) * ((num - 50) / 50)))  # Корректируем формулу для green
                        blue = 0

                        # Убедимся, что при num = 100 green = 176
                    if num == 100:
                        red = 0
                        green = 176
                        blue = 80

                    return f"{int(round(num))}%", RGBColor(red, green, blue)

                except Exception as e:
                    # Если что-то пошло не так, возвращаем оригинальное значение без цвета
                    return str(value), None


            def apply_to_cell(cell, value):
                """Безопасное применение значения к ячейке с обработкой ошибок"""
                try:
                    # Получаем текст и цвет
                    text, color = safe_format_value(value)

                    # Очищаем ячейку
                    text_frame = cell.text_frame
                    text_frame.clear()

                    # Устанавливаем текст
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = text

                    # Устанавливаем цвет если есть
                    if color:
                        run.font.color.rgb = color

                    run.font.name = "SB Sans Text"
                    run.font.bold = False
                    run.font.size = Pt(11)         # Устанавливаем размер шрифта


                except Exception as e:
                    print(f"Ошибка при форматировании ячейки: {str(e)}")
                    cell.text = str(value)  # Просто вставляем текст как ест

            # Форматируем значение
            # text, color = safe_format_value(data['sber911'])


            def safe_format_sberitsm(value):
                """Безопасно форматирует значение с цветом, возвращая всегда кортеж (text, color)"""
                str_val = str(value).strip() if value is not None else ''

                # Значение по умолчанию (если ни одно условие не сработает)
                color = RGBColor(0, 0, 0)  # Черный по умолчанию
                display_text = str_val  # По умолчанию исходное значение

                try:
                    if str_val == 'Исключен':
                        color = RGBColor(17, 160, 81)  # Темно-зеленый
                        display_text = 'Исключен'
                    else:
                        int_val = int(str_val) if str_val else 0  # Преобразуем в int, если не пусто

                        if int_val == 0:
                            color = RGBColor(227, 18, 19)  # Красный
                        elif int_val == 1:
                            color = RGBColor(233, 178, 26)  # Оранжевый
                        elif int_val == 2:
                            color = RGBColor(238, 237, 25)  # Желтый
                        elif int_val == 3:
                            color = RGBColor(130, 187, 72)  # Зеленый
                        elif int_val == 4:
                            color = RGBColor(17, 160, 81)  # Темно-зеленый

                        display_text = str(int_val)  # Форматируем как строку без лишних символов

                except (ValueError, TypeError):
                    # Если преобразование не удалось, оставляем значения по умолчанию
                    pass

                return display_text, color


            def apply_to_cell_vpi_sberitsm(cell, value):
                """Безопасное применение значения к ячейке с обработкой ошибок"""
                try:
                    # Получаем текст и цвет
                    text, color = safe_format_sberitsm(value)
                    # Очищаем ячейку
                    text_frame = cell.text_frame
                    text_frame.clear()
                    # Устанавливаем текст
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = text
                    # Устанавливаем цвет если есть
                    if color:
                        run.font.color.rgb = color
                    run.font.name = "SB Sans Text"
                    run.font.bold = False
                    run.font.size = Pt(11)         # Устанавливаем размер шрифта

                except Exception as e:
                    print(f"Ошибка при форматировании ячейки: {str(e)}")
                    cell.text = str(value)  # Просто вставляем текст как есть





            # Две нижние фуннкции для таблицы 2 а именнно для ------  КОЛИЧЕСТВА впи ЗА КВАРТАЛ  ---------
            def safe_format_vpi(value):
                """Безопасно форматирует значение с цветом, возвращая всегда кортеж (text, color)"""
                str_val = str(value).strip()

                if int(str_val) == 0:
                    color = RGBColor(0, 176, 80)
                else:
                    color = RGBColor(255, 0, 0)

                return f"{int(str_val)}", color  # Возвращаем текст и объект RGBColor


            def apply_to_cell_vpi(cell, value):
                """Безопасное применение значения к ячейке с обработкой ошибок"""
                try:
                    # Получаем текст и цвет
                    text, color = safe_format_vpi(value)
                    # Очищаем ячейку
                    text_frame = cell.text_frame
                    text_frame.clear()
                    # Устанавливаем текст
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = text
                    # Устанавливаем цвет если есть
                    if color:
                        run.font.color.rgb = color
                    run.font.name = "SB Sans Text"
                    run.font.bold = False
                    run.font.size = Pt(11)         # Устанавливаем размер шрифта

                except Exception as e:
                    print(f"Ошибка при форматировании ячейки: {str(e)}")
                    cell.text = str(value)  # Просто вставляем текст как есть



            for shape in slide.shapes:
                if shape.has_table and shape.name == "Таблица 2":
                    target_table = shape.table
                    apply_to_cell(target_table.cell(2, 1), data['sogdtn'])
                    apply_to_cell(target_table.cell(3, 1), data['sberradar'])
                    apply_to_cell(target_table.cell(4, 1), data['sber911'])
                    apply_to_cell(target_table.cell(5, 1), data['inzidenton'])
                    apply_to_cell(target_table.cell(6, 1), data['oesber911'])
                    apply_to_cell_vpi_sberitsm(target_table.cell(7, 1), data['sberitsm'])
                    apply_to_cell(target_table.cell(8, 1), data['poruchenie'])
                    apply_to_cell_vpi(target_table.cell(9, 1), data['vpi'])


            # open_prezent_slide = os.path.join(os.path.dirname(__file__), 'no_hand', 'template.pptx')
            # presentation = Presentation(open_prezent_slide)

            # Функция для изменения цвета текста в зависимости от значения ___ Значение ITAS___ в Уровень технологической надежности
            def apply_color_to_table_cell(table, row_idx, col_idx, color):
                """Изменяет цвет текста в указанной ячейке таблицы"""
                cell = table.cell(row_idx, col_idx)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = color
            # print(key, value[2])
            # Обработка ITAS

            if len(value) > 2 and value[2].strip():

                try:
                    itas_value = float(value[2])
                    data['itas'] = f"{itas_value:.2f}"

                    # Найдем таблицу один раз
                    table_found = None
                    for slide in presentation.slides:
                        for shape in slide.shapes:
                            if shape.has_table and shape.name == "Таблица 2":
                                table_found = shape.table
                                break
                        if table_found:
                            break

                    if table_found:
                        if itas_value < 2:
                            apply_color_to_table_cell(table_found, 1, 1, RGBColor(255, 0, 0))  # Красный
                        elif 1 < itas_value < 3:
                            apply_color_to_table_cell(table_found, 1, 1, RGBColor(255, 255, 0))  # Желтый
                        elif 3 <= itas_value < 4:
                            apply_color_to_table_cell(table_found, 1, 1, RGBColor(130, 187, 72))  # Зеленый
                        elif itas_value == 4:
                            apply_color_to_table_cell(table_found, 1, 1, RGBColor(17, 160, 81))  # Темно-зеленый

                except ValueError:
                    data['itas'] = "0.00"
                    # Можно также установить цвет по умолчанию при ошибке
                    # apply_color_to_table_cell(table_found, 1, 1, DEFAULT_COLOR)

            # # Сохраняем презентацию
            saves = os.path.join(os.path.dirname(__file__), f'{key}' + '.pptx')
            presentation.save(saves)


    else:
        if os.path.exists(save):
            k = f'{name_company}'
            print(f"Файл '{k}' уже существует.")
            # print(f"Файл '{save}'\Отчет_по_ITSM_AS.xlsx' уже существует.")
            # Запрашиваем у пользователя, хочет ли он перезаписать файл
            while True:
                perazapis = input("Вы хотите перезаписать файл Отчет_по_ITSM_AS.xlsx? (y/n): ")
                if perazapis == 'n':
                    print(f"Файл '{k} не перезаписан.")
                    break
                elif perazapis == 'y':
                    wb.save(save)
                    print(f"Файл '{k} перезаписан.")
                    print('=====  Отчет сформирован  =====')
                    break
                else:
                    print("Неверный ответ. Пожалуйста, введите 'y' для перезаписи или 'n' для отказа.")
        else:
            wb.save(save)
            print('=====  Отчет сформирован  =====')


#============================================Конец Генерация слайда ========================================================================







#=======================================================================================================================
if automat_dzo == 2:
    ws = Worksheet()
    ws_two = Worksheet()
    ws_as = Worksheet()
    ws_itsm = Worksheet()
    ws_as_code_dtn = Worksheet()
    ws_code_baza_itsm = Worksheet()
    ws_oprosnik = Worksheet()



    script_dir_all = os.path.abspath(os.path.dirname(__file__))
    mas_files_list = []
    keywords = ["Опросник", "ЕДК", "ЕРО", "База_знаний"]

    file_list = glob.glob(os.path.join(script_dir_all, '*.xlsx'))
    for file in file_list:

        if any(keyword in file for keyword in keywords):
            mas_files_list.append(file)

    oprosnik_files = [os.path.basename(file) for file in mas_files_list if 'Опросник' in os.path.basename(file)]
    baza_znanii_files = [os.path.basename(file) for file in mas_files_list if 'База_знаний' in os.path.basename(file)]
    edk_files = [os.path.basename(file) for file in mas_files_list if 'ЕДК' in file or 'ЕРО' in os.path.basename(file)]

    # Теперь формируем полные пути к файлам
    baza_ero = "".join([os.path.join(script_dir_all, file) for file in edk_files])  # адрес ЕРО
    opros = "".join([os.path.join(script_dir_all, file) for file in oprosnik_files])  # адрес опросника
    znaniya = "".join([os.path.join(script_dir_all, file) for file in baza_znanii_files])  # адрес Базы знаний







    script_dir = os.path.abspath(os.path.dirname(__file__))
    adress_as = os.path.join(script_dir, 'AS')
    print()
    print('-------- Формирую отчет --------')
    print()
    warnings.simplefilter("ignore")

    # Обработка файлов ЕДК/ЕРО
    if edk_files:
        ws.xlsx_to_dict(path=baza_ero, select_sheet='Единая дорожная карта')
    else:
        print("⚠ Ошибка: Нет файлов ЕДК/ЕРО для обработки!")


    def check_version(filepath):
        workbook = openpyxl.load_workbook(filepath)
        worksheet = workbook['Надёжность']
        first_row = worksheet[1]
        mas_ver_dtn = []
        for cell in first_row:
            if cell.value is not None:
                mas_ver_dtn.append(cell.value)
        # Получение ключа страховщика
        insurer_key = next(item for item in mas_ver_dtn if item not in ['ОПРОСНИК\nпо оценке зрелости ITSM-процессов', 'версия опросника:'])

        # Получение ключа версии опросника
        pattern = r'\d+\.\d+\.\d+\.\d+\.\d+'
        version_key = next(item for item in mas_ver_dtn if re.match(pattern, str(item)))
        # Создание нового словаря
        result = {
            'ОПРОСНИК\nпо оценке зрелости ITSM-процессов': insurer_key,
            'версия опросника:': version_key
        }

        expected_version = '6.1.'

        if expected_version not in result['версия опросника:']:
            print(f"Ожидаемая версия должна быть: {expected_version}, а в Опроснике используется версия: {result['версия опросника:']}. Работа скрипта прерывается.")
            exit()

    if __name__ == "__main__":
        filepath = opros
        check_version(filepath)

    workbook_name = openpyxl.load_workbook(opros)
    worksheet = workbook_name['Надёжность']
    first_row = worksheet[1]
    mas_ver_dtn = []
    for cell in first_row:
        if cell.value is not None:
            mas_ver_dtn.append(cell.value)
    # Получение ключа страховщика
    insurer_key = next(item for item in mas_ver_dtn if item not in ['ОПРОСНИК\nпо оценке зрелости ITSM-процессов', 'версия опросника:'])
    # Получение ключа версии опросника
    pattern = r'\d+\.\d+\.\d+\.\d+\.\d+'
    version_key = next(item for item in mas_ver_dtn if re.match(pattern, str(item)))
    # Создание нового словаря
    inn_index = mas_ver_dtn.index('ИНН:')
    result = {
        'ОПРОСНИК\nпо оценке зрелости ITSM-процессов': insurer_key,
        'версия опросника:': version_key,
        'ИНН:': mas_ver_dtn[inn_index + 1]
    }

    #Вытаскиваю критичность из Excel
    second_row = worksheet[3]
    mas_profile_dtn = []
    for cell in second_row:
        if cell.value is not None:
            mas_profile_dtn.append(cell.value)

    # Получение ключа страховщика
    insurer_key_profile = next((item for item in mas_profile_dtn if item not in ['Каким максимальным уровнем критичности обладает ИТ-сервис в компании?', 'ФИО ответственного:']), None)

    # Получение значения ФИО
    fio_dtn = mas_profile_dtn[-1]
    result_profile = {
        'Каким максимальным уровнем критичности обладает ИТ-сервис в компании?': insurer_key_profile,
        'ФИО ответственного:': fio_dtn
    }

    if result_profile['Каким максимальным уровнем критичности обладает ИТ-сервис в компании?'] not in ['Нет подобных ИТ-сервисов', 'OP', 'BO', 'BC и выше']:
        print(f'В файле не указана критичность ИТ-сервис в компании, скрипт работать не будет!!!')
        sys.exit()
    else:
        for k, v in dict_inn_company.items():
            if str(result['ИНН:']) == k:
                name_dtn = v[1]
                name_company = v[1] + '_' +result_profile['Каким максимальным уровнем критичности обладает ИТ-сервис в компании?'] + '_ЕРО.xlsx'
                name_coms = name_company.replace('"', '').replace('«', '').replace('»', '')  #Если в файле  Опросника название компании в кавычках «»
                save = os.path.join(os.path.dirname(__file__), name_coms)
                break

            else:
                name_dtn = result['ОПРОСНИК\nпо оценке зрелости ITSM-процессов']
                name_company = result['ОПРОСНИК\nпо оценке зрелости ITSM-процессов'] + '_' +result_profile['Каким максимальным уровнем критичности обладает ИТ-сервис в компании?'] + '_ЕРО.xlsx'
                name_coms = name_company.replace('"', '').replace('«', '').replace('»', '')  #Если в файле  Опросника название компании в кавычках «»
                save = os.path.join(os.path.dirname(__file__), name_coms)


    workbook_dtn = openpyxl.load_workbook(baza_ero)
    sheet = workbook_dtn.active
    first_column = sheet["A"]
    count_dtn_code = 0
    for cell in first_column:
        if cell.value != 'Код задачи/отклонения':
            count_dtn_code += 1
        else:
            break

    print()
    print('-------- Формирую и расчитываю новую анкету по ЕДК --------')
    print()



    ws_itsm.xlsx_to_dict(path=baza_ero, select_sheet='Единая дорожная карта', data_only=True, skiprows=count_dtn_code)
    ws_as_code_dtn.xlsx_to_dict(path=znaniya, select_sheet='База АС')
    ws_code_baza_itsm.xlsx_to_dict(path=znaniya, select_sheet='База')
    ws_oprosnik.xlsx_to_dict(path=opros, select_sheet='Надёжность', skiprows=3, data_only=True)

    mas_key_dtn = [] #Добавляем все коды по АС из файла Базы Знаний
    for key_dtn in ws_as_code_dtn.sheet_items:
        mas_key_dtn.append(key_dtn.get('Номер вопроса'))


    book_stack = {}
    book_stata_uroven = {}
    book_itsm = {}
    book_profile = {}
    book_dtn_prof = {}
    #Словари для Опросника
    book_opros = {}


    dict_itsm = {'ITSM-INC': 'Управление технологическими инцидентами',
                 'ITSM-CATM': 'Управление каталогом сервисов',
                 'ITSM-SLM': 'Управление уровнями сервисов',
                 'ITSM-PRB': 'Управление технологическими проблемами',
                 'ITSM-ERM': 'Управление технологическими рисками',
                 'ITSM-ICHG': 'Управление инфраструктурными изменениями',
                 'ITSM-CONT': 'Управление непрерывностью технологий',
                 'ITSM-AVL': 'Управление доступностью технологий',
                 'ITSM-M&E': 'Управление технологическими событиями и мониторингом',
                 'ITSM-RFF': 'Управление технологическими запросами на обслуживание',
                 'ITSM-CFG': 'Управление конфигурациями',
                 'ITSM-UPD': 'Управление обновлением ПО',
                 'ITSM-SD': 'Управление технологическими обращениями',
                 'ITSM-CAP': 'Управление мощностями технологий',
                 'ITSM-MNTW': 'Управление регламентными работами'
                 }

    dict_itsm_dtn = {'ITSM-CFG': '5492.CFG.',
                     'ITSM-CATM': '5492.CATM.',
                     'ITSM-SLM': '5492.SLM.',
                     'ITSM-INC': '5492.INC.',
                     'ITSM-PRB': '5492.PRB.',
                     'ITSM-ERM': '5492.ERM.',
                     'ITSM-RFF': '5492.RFF.',
                     'ITSM-ICHG': '5492.ICHG.',
                     'ITSM-UPD': '5492.UPD.',
                     'ITSM-AVL': '5492.AVL.',
                     'ITSM-CONT': '5492.CONT.',
                     'ITSM-M&E': '5492.M&E.',
                     'ITSM-CAP': '5492.CAP.',
                     'ITSM-SD': '5492.SD.',
                     'ITSM-MNTW': '5492.MNTW.'

                     }
    #Собираю информацию из Базы Знаний
    for mq in ws_code_baza_itsm.sheet_items:
        juice_price = mq.setdefault('Код вопроса')
        juice_two = mq.setdefault('Название процесса')
        book_profile.setdefault(juice_price, [])
        book_profile[juice_price].append(juice_two)


   # ================= Собираю информацию из ЕДК и добавляем префикс DTN- к ключам ITSM=================
    # Ключи, к значениям которых нужно добавить префикс DTN-
    # keys_to_update = ['Код вопроса ЕОИТ', 'Код объекта']
    keys_to_update = 'Код вопроса ЕОИТ'
    for item in ws_itsm.sheet_items:
        for key in keys_to_update:
            value = item.get(key, '')
            if 'ITSM-' in value:
                if value and not value.startswith('DTN-'):
                    item[key] = 'DTN-' + value
    #====================================================================================================

    #Собираю информацию из ЕДК
    for rq in ws_itsm.sheet_items:
        juice_zero = rq.setdefault('Код объекта')
        juice_one = rq.setdefault('Код вопроса ЕОИТ')
        juice_two = rq.setdefault('Задача / выявленное отклонение')
        juice_four = rq.setdefault('Рекомендованные мероприятия')
        juice_six = rq.setdefault('Уровень принятия решения')
        juice_three = rq.setdefault('Вхождение в уровни готовности')
        juice_seven = rq.setdefault('Уровень критичности')
        juice_eight = rq.setdefault('Запланированное мероприятие')
        juice_night = rq.setdefault('Код мероприятия из трекера компании Группы (при использовании трекера)')
        juice_ten = rq.setdefault('Дата начала мероприятий')
        juice_eleven = rq.setdefault('Дата окончания мероприятий')
        juice_twelv = rq.setdefault('Ответственный')
        juice_thirteen = rq.setdefault('Статус')
        juice_fourteen = rq.setdefault('Комментарии компании Группы')
        juice_fiveteen = rq.setdefault('Вид объекта 1 уровень')
        juce_sixteen = rq.setdefault('Объект')
        juice_seventeen = rq.setdefault('Комментарии инициатора')

        #Если есть код вопроса в словаре Базы знаний, то забираю эти поля
        if juice_one in book_profile.keys():
            book_itsm.setdefault(juice_one, [])
            book_itsm[juice_one].append(juice_zero)
            book_itsm[juice_one].append(juice_two)
            book_itsm[juice_one].append(juice_three)
            book_itsm[juice_one].append(juice_four)
            book_itsm[juice_one].append(juice_six)
            book_itsm[juice_one].append(juice_seven)
            book_itsm[juice_one].append(juice_eight)
            book_itsm[juice_one].append(juice_night)
            book_itsm[juice_one].append(juice_ten)
            book_itsm[juice_one].append(juice_eleven)
            book_itsm[juice_one].append(juice_twelv)
            book_itsm[juice_one].append(juice_thirteen)
            book_itsm[juice_one].append(juice_fourteen)
            book_itsm[juice_one].append(juice_seventeen)
            book_itsm[juice_one].append(juce_sixteen)
            book_itsm[juice_one].append(juice_fiveteen)

        #Так как в общей ЕДК будут АС не только в ДТН но и в других группазх, то забираю только те коды которые соответствуют нашей Базе знаниий
        if juice_fiveteen == 'АС' and juice_one in mas_key_dtn:
            #Словарь для добавления информации по всем АС для ЕДК
            book_stack.setdefault(juce_sixteen, [])
            book_stack[juce_sixteen].append(juice_one)
            book_stack[juce_sixteen].append(juice_zero)
            book_stack[juce_sixteen].append(juice_two)
            book_stack[juce_sixteen].append(juice_three)
            book_stack[juce_sixteen].append(juice_four)
            book_stack[juce_sixteen].append(juice_six)
            book_stack[juce_sixteen].append(juice_seven)
            book_stack[juce_sixteen].append(juice_eight)
            book_stack[juce_sixteen].append(juice_night)
            book_stack[juce_sixteen].append(juice_ten)
            book_stack[juce_sixteen].append(juice_eleven)
            book_stack[juce_sixteen].append(juice_twelv)
            book_stack[juce_sixteen].append(juice_thirteen)
            book_stack[juce_sixteen].append(juice_fourteen)
            book_stack[juce_sixteen].append(juice_fiveteen)
            book_stack[juce_sixteen].append(juice_seventeen)

            #Словарь для расчета АС
            book_stata_uroven.setdefault(juce_sixteen, [])
            book_stata_uroven[juce_sixteen].append(juice_one)
            book_stata_uroven[juce_sixteen].append(juice_thirteen)
            book_stata_uroven[juce_sixteen].append(juice_seven)
            book_stata_uroven[juce_sixteen].append(juice_zero)

    #------------------------------------------------------------------------------------------------------------------

        # ----- При парсинге из ЕДК в случае если отсутствует префикс DTN, то его добавляем в словарь-----
    # Создаём новый словарь с нужными ключами
    # updated_book_opros_edk = {}
    # for key, value in book_itsm.items():
    #     if not key.startswith('DTN-'):
    #         new_key = 'DTN-' + key
    #     else:
    #         new_key = key
    #     updated_book_opros_edk[new_key] = value
    #
    # # Перезаписываем исходный словарь (если нужно)
    # book_itsm = updated_book_opros_edk
    # print(book_itsm)


    dict_itsm = {'ITSM-INC': 'Управление технологическими инцидентами',
                 'ITSM-CATM': 'Управление каталогом сервисов',
                 'ITSM-SLM': 'Управление уровнями сервисов',
                 'ITSM-PRB': 'Управление технологическими проблемами',
                 'ITSM-ERM': 'Управление технологическими рисками',
                 'ITSM-ICHG': 'Управление инфраструктурными изменениями',
                 'ITSM-CONT': 'Управление непрерывностью технологий',
                 'ITSM-AVL': 'Управление доступностью технологий',
                 'ITSM-M&E': 'Управление технологическими событиями и мониторингом',
                 'ITSM-RFF': 'Управление технологическими запросами на обслуживание',
                 'ITSM-CFG': 'Управление конфигурациями',
                 'ITSM-UPD': 'Управление обновлением ПО',
                 'ITSM-SD': 'Управление технологическими обращениями',
                 'ITSM-CAP': 'Управление мощностями технологий',
                 'ITSM-MNTW': 'Управление регламентными работами'
                 }

    dict_itsm_dtn = {'ITSM-CFG': '5492.CFG.',
                     'ITSM-CATM': '5492.CATM.',
                     'ITSM-SLM': '5492.SLM.',
                     'ITSM-INC': '5492.INC.',
                     'ITSM-PRB': '5492.PRB.',
                     'ITSM-ERM': '5492.ERM.',
                     'ITSM-RFF': '5492.RFF.',
                     'ITSM-ICHG': '5492.ICHG.',
                     'ITSM-UPD': '5492.UPD.',
                     'ITSM-AVL': '5492.AVL.',
                     'ITSM-CONT': '5492.CONT.',
                     'ITSM-M&E': '5492.M&E.',
                     'ITSM-CAP': '5492.CAP.',
                     'ITSM-SD': '5492.SD.',
                     'ITSM-MNTW': '5492.MNTW.'

                     }

    #TODO: Собираем информацию с Опросника согласно полям!
    for b in ws_oprosnik.sheet_items:
        juice_id = b.setdefault('ID\nвопроса')
        juice_urov_ur_got = b.setdefault('Уровень готовности в зависимости от критичности АС')
        juice_ur_otklon = b.setdefault('Уровень отклонения')
        juice_answer = b.setdefault('Ответ с учётом артефакта \n(авто)')

        book_opros.setdefault(juice_id, [])
        book_opros[juice_id].append(juice_urov_ur_got)
        book_opros[juice_id].append(juice_ur_otklon)
        book_opros[juice_id].append(juice_answer)
    removed_value = book_opros.pop('-')  # Удаляем ключ '-'

    #======================МАПИМ ВСЁ ИЗ ОПРОСНИКА И ДОБАВЛЯЕМ ПРЕФИКС DTN- ЕСЛИ В ОПРОСНИКЕ ЕГО НЕТ==============================
    # Создаём новый словарь с нужными ключами
    updated_book_opros = {}
    for key, value in book_opros.items():
        if not key.startswith('DTN-'):
            new_key = 'DTN-' + key
        else:
            new_key = key
        updated_book_opros[new_key] = value

    # Перезаписываем исходный словарь (если нужно)
    book_opros = updated_book_opros
    #=============================================================================================================================

    oprosnik_dict = copy.deepcopy(book_opros) #Делаем копию Опросника(для будущего сравнения с ЕДК по отклонениям и количеству общего к выполненным ---> Пример Критичный: 120/80)


    # Создаем defaultdict, где значения по умолчанию — пустые словари
    result_dicts = defaultdict(dict)
    result_dicts_edk = defaultdict(dict)

    #Расчитываем сколько компания выполнила за период отклонений первые цифры по критичности
    # должны совпадать с первым ЕРО(к примеру 107 / 0, было то тут должна остаться цифра 107)
    itsm_dtn_all = {'Критичный': 0, 'Высокий': 0, 'Умеренный': 0, 'Рекомендация на развитие': 0} # Этот словарь для подсчета общего числа отклонений
    itsm_dtn_ikrement = {'Критичный': 0, 'Высокий': 0, 'Умеренный': 0, 'Рекомендация на развитие': 0}  # Этот словарь для подсчета выполненных отклонений
    for i, d in book_itsm.items():
        if i in book_profile:
            itsm_dtn_all[d[5]] += 1
        if '(ПАО) Устранено - подтверждено ПАО' in d[11] or '(ПАО) Исключение из реестра - подтверждено ПАО' in d[11]:
            itsm_dtn_ikrement[d[5]] += 1


    prefixes = [
        'ITSM-INC', 'ITSM-CATM', 'ITSM-CFG', 'ITSM-SLM', 'ITSM-PRB', 'ITSM-ERM',
        'ITSM-RFF', 'ITSM-ICHG', 'ITSM-UPD', 'ITSM-AVL', 'ITSM-M&E', 'ITSM-CONT',
        'ITSM-CAP', 'ITSM-SD', 'ITSM-MNTW']

    for key, value in book_opros.items():
        if value[0] != 'ЦТК':
            for prefix in prefixes:
                if prefix in key:
                    result_dicts[prefix][key] = value


                    break


    # ОБНОВЛЯЕМ ДАТЫ В СЛОВАРЕ С ITSM, Т.Е. МЕНЯЕМ ЕСЛИ НЕВЕРНЫЙ ФОРМАТ ПРИМЕР ----> 14.02.2025 ----> 102025
    mas_start_kvartal = []
    mas_end_kvartal = []

    # Массив для записи будущих кварталов
    mas_data_list = []

    for i, v in book_itsm.items():
        new_v = list(v)  # Создаем копию списка значений, так как v может быть tuple
        if v[8] != "":
            result_start = date_to_quarter_year(v[8])
            new_v[8] = result_start
            mas_start_kvartal.append(result_start)
        else:
            pass
            # print(f"У {i} Пустая дата начала")
        if v[9] != "":
            result_end = date_to_quarter_year(v[9])
            new_v[9] = result_end
            mas_end_kvartal.append(result_end)
        else:
            pass
            # print(f"У {i} Пустая дата окончания")

        book_itsm[i] = tuple(new_v) if isinstance(v, tuple) else new_v  # Сохраняем обратно в словарь
    #===============================================================================================================

    for key, value in book_itsm.items():
        for prefix in prefixes:
            if prefix in key:
                result_dicts_edk[prefix][key] = value
                break



    result_stack = {}
    for key, value_list in book_stack.items():
        result_stack[key] = [value_list[i:i+16] for i in range(0, len(value_list), 16)]


    # ОБНОВЛЯЕМ ДАТЫ В СЛОВАРЕ С ИТ-УСЛУГАМИ, Т.Е. МЕНЯЕМ ЕСЛИ НЕВЕРНЫЙ ФОРМАТ ПРИМЕР ----> 14.02.2025 ----> 102025
    for process_name, process_data in result_stack.items():
        updated_data = []
        for item in process_data:
            new_item = item.copy()  # Создаем копию элемента

            # Обновляем дату начала (индекс 9)
            if new_item[9]:
                new_item[9] = date_to_quarter_year(new_item[9]) or new_item[9]
            # Обновляем дату окончания (индекс 10)
            if new_item[10]:
                new_item[10] = date_to_quarter_year(new_item[10]) or new_item[10]

            updated_data.append(new_item)

        result_stack[process_name] = updated_data
    #===============================================================================================================




    # ==============================================================================
    # Функция для нахождения следующих кварталов после текущего квартала ДЛЯ ВСЕХ ITSM
    for category, items in result_dicts_edk.items():
        for item_id, item_data in items.items():
            status = item_data[11] if len(item_data) > 11 else ""
            status_excluded = (
                    status == "(ПАО) Устранено - подтверждено ПАО" or
                    status == "(ПАО) Исключение из реестра - подтверждено ПАО" or
                    status == "(ПАО) Принятие рисков - утверждено ПАО"
            )

            if not status_excluded and item_data[9] != "":
                mas_data_list.append(item_data[9])


    # ===============================================================================
    # Функция для нахождения следующих кварталов после текущего квартала ДЛЯ ВСЕХ AS
    for date_str, date_tuple in result_stack.items():
        for i in date_tuple:


            # print(date_str, i)
            status_excluded = (
                    i[12] == "(ПАО) Устранено - подтверждено ПАО" or
                    i[12] == "(ПАО) Исключение из реестра - подтверждено ПАО"
            )
            if not status_excluded and i[10] != "":
                mas_data_list.append(i[10])

    future_quarters = find_future_quarters(current_quarter, mas_data_list)


    # ==============================================================================

    # Сортировка дат от меньшего к большему
    # =====================================
    def sort_key(q):
        quarter = int(q[0])  # Первый символ — номер квартала
        year = int(q[2:])  # Год — после "Q"
        return (year, quarter)


    # Удаляем дубликаты через set(), затем сортируем
    unique_quarters = list(set(future_quarters))  # Оставляем только уникальные значения
    sorted_quarters = sorted(unique_quarters, key=sort_key)


    # ==============================================================================



    # Теперь у нас есть словари для каждого префикса взятого из ОПРОСНИКА
    dict_inc = result_dicts['ITSM-INC']
    dict_inc_edk = result_dicts_edk['ITSM-INC']
    dict_catm = result_dicts['ITSM-CATM']
    dict_catm_edk = result_dicts_edk['ITSM-CATM']
    dict_cfg = result_dicts['ITSM-CFG']
    dict_cfg_edk = result_dicts_edk['ITSM-CFG']
    dict_slm = result_dicts['ITSM-SLM']
    dict_slm_edk = result_dicts_edk['ITSM-SLM']
    dict_prb = result_dicts['ITSM-PRB']
    dict_prb_edk = result_dicts_edk['ITSM-PRB']
    dict_erm = result_dicts['ITSM-ERM']
    dict_erm_edk = result_dicts_edk['ITSM-ERM']
    dict_rff = result_dicts['ITSM-RFF']
    dict_rff_edk = result_dicts_edk['ITSM-RFF']
    dict_ichg = result_dicts['ITSM-ICHG']
    dict_ichg_edk = result_dicts_edk['ITSM-ICHG']
    dict_upd = result_dicts['ITSM-UPD']
    dict_upd_edk = result_dicts_edk['ITSM-UPD']
    dict_avl = result_dicts['ITSM-AVL']
    dict_avl_edk = result_dicts_edk['ITSM-AVL']
    dict_me = result_dicts['ITSM-M&E']
    dict_me_edk = result_dicts_edk['ITSM-M&E']
    dict_cont = result_dicts['ITSM-CONT']
    dict_cont_edk = result_dicts_edk['ITSM-CONT']
    dict_cap = result_dicts['ITSM-CAP']
    dict_cap_edk = result_dicts_edk['ITSM-CAP']
    dict_mntw = result_dicts['ITSM-MNTW']
    dict_mntw_edk = result_dicts_edk['ITSM-MNTW']
    dict_sd = result_dicts['ITSM-SD']
    dict_sd_edk = result_dicts_edk['ITSM-SD']

    #Собираем массив с теми отклонениями которые получили ИТОГОВЫЙ ответ "НЕТ" нужны для формирования информации из базы знаний
    dict_all_answer_no = []

    dict_inc_vnd = {'ITSM-INC': {'Всего требований': 0,
                                 'Всего требований на 3': 0,
                                 'Всего выполнено требований на 3': 0,
                                 'Всего требований на 4': 0,
                                 'Всего выполнено требований на 4': 0}}

    dict_catm_vnd = {'ITSM-CATM': {'Всего требований': 0,
                                   'Всего требований на 3': 0,
                                   'Всего выполнено требований на 3': 0,
                                   'Всего требований на 4': 0,
                                   'Всего выполнено требований на 4': 0}}

    dict_cfg_vnd = {'ITSM-CFG': {'Всего требований': 0,
                                 'Всего требований на 3': 0,
                                 'Всего выполнено требований на 3': 0,
                                 'Всего требований на 4': 0,
                                 'Всего выполнено требований на 4': 0}}


    dict_slm_vnd = {'ITSM-SLM': {'Всего требований': 0,
                                 'Всего требований на 3': 0,
                                 'Всего выполнено требований на 3': 0,
                                 'Всего требований на 4': 0,
                                 'Всего выполнено требований на 4': 0}}

    dict_prb_vnd = {'ITSM-PRB': {'Всего требований': 0,
                                 'Всего требований на 3': 0,
                                 'Всего выполнено требований на 3': 0,
                                 'Всего требований на 4': 0,
                                 'Всего выполнено требований на 4': 0}}


    dict_erm_vnd = {'ITSM-ERM': {'Всего требований': 0,
                                 'Всего требований на 3': 0,
                                 'Всего выполнено требований на 3': 0,
                                 'Всего требований на 4': 0,
                                 'Всего выполнено требований на 4': 0}}


    dict_rff_vnd = {'ITSM-RFF': {'Всего требований': 0,
                                 'Всего требований на 3': 0,
                                 'Всего выполнено требований на 3': 0,
                                 'Всего требований на 4': 0,
                                 'Всего выполнено требований на 4': 0}}

    dict_ichg_vnd = {'ITSM-ICHG': {'Всего требований': 0,
                                   'Всего требований на 3': 0,
                                   'Всего выполнено требований на 3': 0,
                                   'Всего требований на 4': 0,
                                   'Всего выполнено требований на 4': 0}}

    dict_upd_vnd = {'ITSM-UPD': {'Всего требований': 0,
                                 'Всего требований на 3': 0,
                                 'Всего выполнено требований на 3': 0,
                                 'Всего требований на 4': 0,
                                 'Всего выполнено требований на 4': 0}}


    dict_avl_vnd = {'ITSM-AVL': {'Всего требований': 0,
                                 'Всего требований на 3': 0,
                                 'Всего выполнено требований на 3': 0,
                                 'Всего требований на 4': 0,
                                 'Всего выполнено требований на 4': 0}}


    dict_me_vnd = {'ITSM-M&E': {'Всего требований': 0,
                                'Всего требований на 3': 0,
                                'Всего выполнено требований на 3': 0,
                                'Всего требований на 4': 0,
                                'Всего выполнено требований на 4': 0}}

    dict_cont_vnd = {'ITSM-CONT': {'Всего требований': 0,
                                   'Всего требований на 3': 0,
                                   'Всего выполнено требований на 3': 0,
                                   'Всего требований на 4': 0,
                                   'Всего выполнено требований на 4': 0}}

    dict_cap_vnd = {'ITSM-CAP': {'Всего требований': 0,
                                 'Всего требований на 3': 0,
                                 'Всего выполнено требований на 3': 0,
                                 'Всего требований на 4': 0,
                                 'Всего выполнено требований на 4': 0}}

    dict_sd_vnd = {'ITSM-SD': {'Всего требований': 0,
                               'Всего требований на 3': 0,
                               'Всего выполнено требований на 3': 0,
                               'Всего требований на 4': 0,
                               'Всего выполнено требований на 4': 0}}

    dict_mntw_vnd = {'ITSM-MNTW': {'Всего требований': 0,
                                   'Всего требований на 3': 0,
                                   'Всего выполнено требований на 3': 0,
                                   'Всего требований на 4': 0,
                                   'Всего выполнено требований на 4': 0}}




    # ********************************************** ITSM-INC ********************************************************************

    for k, v in dict_inc.items():
        if k in dict_inc_edk:

            if '(ПАО) Устранено - подтверждено ПАО' in dict_inc_edk[k][11] or '(ПАО) Исключение из реестра - подтверждено ПАО' in dict_inc_edk[k][11]:
                dict_inc[k][2] = 'да'
        if v[0] == '3':
            dict_inc_vnd['ITSM-INC']['Всего требований на 3'] += 1
        if v[0] == '3' and v[2] == 'да':
            dict_inc_vnd['ITSM-INC']['Всего выполнено требований на 3'] += 1
        if v[0] == '4':
            dict_inc_vnd['ITSM-INC']['Всего требований на 4'] += 1
        if v[0] == '4' and v[2] == 'да':
            dict_inc_vnd['ITSM-INC']['Всего выполнено требований на 4'] += 1

    dict_inc_vnd['ITSM-INC']['Всего требований'] = dict_inc_vnd['ITSM-INC']['Всего требований на 3'] + dict_inc_vnd['ITSM-INC']['Всего требований на 4']

    # ********************************************** ITSM-CATM ********************************************************************

    for k, v in dict_catm.items():
        if k in dict_catm_edk:
            if '(ПАО) Устранено - подтверждено ПАО' in dict_catm_edk[k][11] or '(ПАО) Исключение из реестра - подтверждено ПАО' in dict_catm_edk[k][11]:
                dict_catm[k][2] = 'да'
        if v[0] == '3':
            dict_catm_vnd['ITSM-CATM']['Всего требований на 3'] += 1
        if v[0] == '3' and v[2] == 'да':
            dict_catm_vnd['ITSM-CATM']['Всего выполнено требований на 3'] += 1
        if v[0] == '4':
            dict_catm_vnd['ITSM-CATM']['Всего требований на 4'] += 1
        if v[0] == '4' and v[2] == 'да':
            dict_catm_vnd['ITSM-CATM']['Всего выполнено требований на 4'] += 1

    dict_catm_vnd['ITSM-CATM']['Всего требований'] = dict_catm_vnd['ITSM-CATM']['Всего требований на 3'] + dict_catm_vnd['ITSM-CATM']['Всего требований на 4']

    # ********************************************** ITSM-CFG ********************************************************************

    for k, v in dict_cfg.items():
        if k in dict_cfg_edk:
            if '(ПАО) Устранено - подтверждено ПАО' in dict_cfg_edk[k][11] or '(ПАО) Исключение из реестра - подтверждено ПАО' in dict_cfg_edk[k][11]:
                dict_cfg[k][2] = 'да'
        if v[0] == '3':
            dict_cfg_vnd['ITSM-CFG']['Всего требований на 3'] += 1
        if v[0] == '3' and v[2] == 'да':
            dict_cfg_vnd['ITSM-CFG']['Всего выполнено требований на 3'] += 1
        if v[0] == '4':
            dict_cfg_vnd['ITSM-CFG']['Всего требований на 4'] += 1
        if v[0] == '4' and v[2] == 'да':
            dict_cfg_vnd['ITSM-CFG']['Всего выполнено требований на 4'] += 1

    dict_cfg_vnd['ITSM-CFG']['Всего требований'] = dict_cfg_vnd['ITSM-CFG']['Всего требований на 3'] + dict_cfg_vnd['ITSM-CFG']['Всего требований на 4']

    # ********************************************** ITSM-SLM ********************************************************************

    for k, v in dict_slm.items():
        if k in dict_slm_edk:
            if '(ПАО) Устранено - подтверждено ПАО' in dict_slm_edk[k][11] or '(ПАО) Исключение из реестра - подтверждено ПАО' in dict_slm_edk[k][11]:
                dict_slm[k][2] = 'да'
        if v[0] == '3':
            dict_slm_vnd['ITSM-SLM']['Всего требований на 3'] += 1
        if v[0] == '3' and v[2] == 'да':
            dict_slm_vnd['ITSM-SLM']['Всего выполнено требований на 3'] += 1
        if v[0] == '4':
            dict_slm_vnd['ITSM-SLM']['Всего требований на 4'] += 1
        if v[0] == '4' and v[2] == 'да':
            dict_slm_vnd['ITSM-SLM']['Всего выполнено требований на 4'] += 1

    dict_slm_vnd['ITSM-SLM']['Всего требований'] = dict_slm_vnd['ITSM-SLM']['Всего требований на 3'] + dict_slm_vnd['ITSM-SLM']['Всего требований на 4']

    # ********************************************** ITSM-PRB ********************************************************************

    for k, v in dict_prb.items():
        if k in dict_prb_edk:
            if '(ПАО) Устранено - подтверждено ПАО' in dict_prb_edk[k][11] or '(ПАО) Исключение из реестра - подтверждено ПАО' in dict_prb_edk[k][11]:
                dict_prb[k][2] = 'да'
        if v[0] == '3':
            dict_prb_vnd['ITSM-PRB']['Всего требований на 3'] += 1
        if v[0] == '3' and v[2] == 'да':
            dict_prb_vnd['ITSM-PRB']['Всего выполнено требований на 3'] += 1
        if v[0] == '4':
            dict_prb_vnd['ITSM-PRB']['Всего требований на 4'] += 1
        if v[0] == '4' and v[2] == 'да':
            dict_prb_vnd['ITSM-PRB']['Всего выполнено требований на 4'] += 1

    dict_prb_vnd['ITSM-PRB']['Всего требований'] = dict_prb_vnd['ITSM-PRB']['Всего требований на 3'] + dict_prb_vnd['ITSM-PRB']['Всего требований на 4']

    # ********************************************** ITSM-ERM ********************************************************************
    for k, v in dict_erm.items():
        if k in dict_erm_edk:
            if '(ПАО) Устранено - подтверждено ПАО' in dict_erm_edk[k][11] or '(ПАО) Исключение из реестра - подтверждено ПАО' in dict_erm_edk[k][11]:
                dict_erm[k][2] = 'да'
        if v[0] == '3':
            dict_erm_vnd['ITSM-ERM']['Всего требований на 3'] += 1
        if v[0] == '3' and v[2] == 'да':
            dict_erm_vnd['ITSM-ERM']['Всего выполнено требований на 3'] += 1
        if v[0] == '4':
            dict_erm_vnd['ITSM-ERM']['Всего требований на 4'] += 1
        if v[0] == '4' and v[2] == 'да':
            dict_erm_vnd['ITSM-ERM']['Всего выполнено требований на 4'] += 1

    dict_erm_vnd['ITSM-ERM']['Всего требований'] = dict_erm_vnd['ITSM-ERM']['Всего требований на 3'] + dict_erm_vnd['ITSM-ERM']['Всего требований на 4']

    # ********************************************** ITSM-RFF ********************************************************************
    for k, v in dict_rff.items():
        if k in dict_rff_edk:
            if '(ПАО) Устранено - подтверждено ПАО' in dict_rff_edk[k][11] or '(ПАО) Исключение из реестра - подтверждено ПАО' in dict_rff_edk[k][11]:
                dict_rff[k][2] = 'да'
        if v[0] == '3':
            dict_rff_vnd['ITSM-RFF']['Всего требований на 3'] += 1
        if v[0] == '3' and v[2] == 'да':
            dict_rff_vnd['ITSM-RFF']['Всего выполнено требований на 3'] += 1
        if v[0] == '4':
            dict_rff_vnd['ITSM-RFF']['Всего требований на 4'] += 1
        if v[0] == '4' and v[2] == 'да':
            dict_rff_vnd['ITSM-RFF']['Всего выполнено требований на 4'] += 1

    dict_rff_vnd['ITSM-RFF']['Всего требований'] = dict_rff_vnd['ITSM-RFF']['Всего требований на 3'] + dict_rff_vnd['ITSM-RFF']['Всего требований на 4']

    # ********************************************** ITSM-ICHG ********************************************************************
    for k, v in dict_ichg.items():
        if k in dict_ichg_edk:
            if '(ПАО) Устранено - подтверждено ПАО' in dict_ichg_edk[k][11] or '(ПАО) Исключение из реестра - подтверждено ПАО' in dict_ichg_edk[k][11]:
                dict_ichg[k][2] = 'да'
        if v[0] == '3':
            dict_ichg_vnd['ITSM-ICHG']['Всего требований на 3'] += 1
        if v[0] == '3' and v[2] == 'да':
            dict_ichg_vnd['ITSM-ICHG']['Всего выполнено требований на 3'] += 1
        if v[0] == '4':
            dict_ichg_vnd['ITSM-ICHG']['Всего требований на 4'] += 1
        if v[0] == '4' and v[2] == 'да':
            dict_ichg_vnd['ITSM-ICHG']['Всего выполнено требований на 4'] += 1

    dict_ichg_vnd['ITSM-ICHG']['Всего требований'] = dict_ichg_vnd['ITSM-ICHG']['Всего требований на 3'] + dict_ichg_vnd['ITSM-ICHG']['Всего требований на 4']

    # ********************************************** ITSM-UPD ********************************************************************

    for k, v in dict_upd.items():
        if k in dict_upd_edk:
            if '(ПАО) Устранено - подтверждено ПАО' in dict_upd_edk[k][11] or '(ПАО) Исключение из реестра - подтверждено ПАО' in dict_upd_edk[k][11]:
                dict_upd[k][2] = 'да'
        if v[0] == '3':
            dict_upd_vnd['ITSM-UPD']['Всего требований на 3'] += 1
        if v[0] == '3' and v[2] == 'да':
            dict_upd_vnd['ITSM-UPD']['Всего выполнено требований на 3'] += 1
        if v[0] == '4':
            dict_upd_vnd['ITSM-UPD']['Всего требований на 4'] += 1
        if v[0] == '4' and v[2] == 'да':
            dict_upd_vnd['ITSM-UPD']['Всего выполнено требований на 4'] += 1

    dict_upd_vnd['ITSM-UPD']['Всего требований'] = dict_upd_vnd['ITSM-UPD']['Всего требований на 3'] + dict_upd_vnd['ITSM-UPD']['Всего требований на 4']

    # ********************************************** ITSM-AVL ********************************************************************
    for k, v in dict_avl.items():
        if k in dict_avl_edk:
            if '(ПАО) Устранено - подтверждено ПАО' in dict_avl_edk[k][11] or '(ПАО) Исключение из реестра - подтверждено ПАО' in dict_avl_edk[k][11]:
                dict_avl[k][2] = 'да'
        if v[0] == '3':
            dict_avl_vnd['ITSM-AVL']['Всего требований на 3'] += 1
        if v[0] == '3' and v[2] == 'да':
            dict_avl_vnd['ITSM-AVL']['Всего выполнено требований на 3'] += 1
        if v[0] == '4':
            dict_avl_vnd['ITSM-AVL']['Всего требований на 4'] += 1
        if v[0] == '4' and v[2] == 'да':
            dict_avl_vnd['ITSM-AVL']['Всего выполнено требований на 4'] += 1

    dict_avl_vnd['ITSM-AVL']['Всего требований'] = dict_avl_vnd['ITSM-AVL']['Всего требований на 3'] + dict_avl_vnd['ITSM-AVL']['Всего требований на 4']

    # ********************************************** ITSM-M&E ********************************************************************
    for k, v in dict_me.items():
        if k in dict_me_edk:
            if '(ПАО) Устранено - подтверждено ПАО' in dict_me_edk[k][11] or '(ПАО) Исключение из реестра - подтверждено ПАО' in dict_me_edk[k][11]:
                dict_me[k][2] = 'да'
        if v[0] == '3':
            dict_me_vnd['ITSM-M&E']['Всего требований на 3'] += 1
        if v[0] == '3' and v[2] == 'да':
            dict_me_vnd['ITSM-M&E']['Всего выполнено требований на 3'] += 1
        if v[0] == '4':
            dict_me_vnd['ITSM-M&E']['Всего требований на 4'] += 1
        if v[0] == '4' and v[2] == 'да':
            dict_me_vnd['ITSM-M&E']['Всего выполнено требований на 4'] += 1

    dict_me_vnd['ITSM-M&E']['Всего требований'] = dict_me_vnd['ITSM-M&E']['Всего требований на 3'] + dict_me_vnd['ITSM-M&E']['Всего требований на 4']

    # ********************************************** ITSM-CONT ********************************************************************
    for k, v in dict_cont.items():
        if k in dict_cont_edk:
            if '(ПАО) Устранено - подтверждено ПАО' in dict_cont_edk[k][11] or '(ПАО) Исключение из реестра - подтверждено ПАО' in dict_cont_edk[k][11]:
                dict_cont[k][2] = 'да'
        if v[0] == '3':
            dict_cont_vnd['ITSM-CONT']['Всего требований на 3'] += 1
        if v[0] == '3' and v[2] == 'да':
            dict_cont_vnd['ITSM-CONT']['Всего выполнено требований на 3'] += 1
        if v[0] == '4':
            dict_cont_vnd['ITSM-CONT']['Всего требований на 4'] += 1
        if v[0] == '4' and v[2] == 'да':
            dict_cont_vnd['ITSM-CONT']['Всего выполнено требований на 4'] += 1

    dict_cont_vnd['ITSM-CONT']['Всего требований'] = dict_cont_vnd['ITSM-CONT']['Всего требований на 3'] + dict_cont_vnd['ITSM-CONT']['Всего требований на 4']

    # ********************************************** ITSM-CAP ********************************************************************
    for k, v in dict_cap.items():
        if k in dict_cap_edk:
            if '(ПАО) Устранено - подтверждено ПАО' in dict_cap_edk[k][11] or '(ПАО) Исключение из реестра - подтверждено ПАО' in dict_cap_edk[k][11]:
                dict_cap[k][2] = 'да'
        if v[0] == '3':
            dict_cap_vnd['ITSM-CAP']['Всего требований на 3'] += 1
        if v[0] == '3' and v[2] == 'да':
            dict_cap_vnd['ITSM-CAP']['Всего выполнено требований на 3'] += 1
        if v[0] == '4':
            dict_cap_vnd['ITSM-CAP']['Всего требований на 4'] += 1
        if v[0] == '4' and v[2] == 'да':
            dict_cap_vnd['ITSM-CAP']['Всего выполнено требований на 4'] += 1

    dict_cap_vnd['ITSM-CAP']['Всего требований'] = dict_cap_vnd['ITSM-CAP']['Всего требований на 3'] + dict_cap_vnd['ITSM-CAP']['Всего требований на 4']

    # ********************************************** ITSM-MNTW ********************************************************************
    for k, v in dict_mntw.items():
        if k in dict_mntw_edk:
            if '(ПАО) Устранено - подтверждено ПАО' in dict_mntw_edk[k][11] or '(ПАО) Исключение из реестра - подтверждено ПАО' in dict_mntw_edk[k][11]:
                dict_mntw[k][2] = 'да'
        if v[0] == '3':
            dict_mntw_vnd['ITSM-MNTW']['Всего требований на 3'] += 1
        if v[0] == '3' and v[2] == 'да':
            dict_mntw_vnd['ITSM-MNTW']['Всего выполнено требований на 3'] += 1
        if v[0] == '4':
            dict_mntw_vnd['ITSM-MNTW']['Всего требований на 4'] += 1
        if v[0] == '4' and v[2] == 'да':
            dict_mntw_vnd['ITSM-MNTW']['Всего выполнено требований на 4'] += 1

    dict_mntw_vnd['ITSM-MNTW']['Всего требований'] = dict_mntw_vnd['ITSM-MNTW']['Всего требований на 3'] + dict_mntw_vnd['ITSM-MNTW']['Всего требований на 4']

    # ********************************************** ITSM-SD ********************************************************************
    for k, v in dict_sd.items():
        if k in dict_sd_edk:
            if '(ПАО) Устранено - подтверждено ПАО' in dict_sd_edk[k][11] or '(ПАО) Исключение из реестра - подтверждено ПАО' in dict_sd_edk[k][11]:
                dict_sd[k][2] = 'да'
        if v[0] == '3':
            dict_sd_vnd['ITSM-SD']['Всего требований на 3'] += 1
        if v[0] == '3' and v[2] == 'да':
            dict_sd_vnd['ITSM-SD']['Всего выполнено требований на 3'] += 1
        if v[0] == '4':
            dict_sd_vnd['ITSM-SD']['Всего требований на 4'] += 1
        if v[0] == '4' and v[2] == 'да':
            dict_sd_vnd['ITSM-SD']['Всего выполнено требований на 4'] += 1

    dict_sd_vnd['ITSM-SD']['Всего требований'] = dict_sd_vnd['ITSM-SD']['Всего требований на 3'] + dict_sd_vnd['ITSM-SD']['Всего требований на 4']

    # *********************************************************************************************************************************************

    # Все словари в переменной result_dicts_itsm
    result_dicts_itsm = {'ITSM-INC': dict_inc,
                         'ITSM-CATM': dict_catm,
                         'ITSM-CFG': dict_cfg,
                         'ITSM-SLM': dict_slm,
                         'ITSM-PRB': dict_prb,
                         'ITSM-ERM': dict_erm,
                         'ITSM-RFF': dict_rff,
                         'ITSM-ICHG': dict_ichg,
                         'ITSM-UPD': dict_upd,
                         'ITSM-AVL': dict_avl,
                         'ITSM-M&E': dict_me,
                         'ITSM-CONT': dict_cont,
                         'ITSM-CAP': dict_cap,
                         'ITSM-MNTW': dict_mntw,
                         'ITSM-SD': dict_sd
                         }

    # Объединяем все словари в один
    combined_dict = {}
    for key, value in result_dicts_itsm.items():
        combined_dict.update(value)

    # Создаем пустой словарь для хранения результатов
        # Создаем новый словарь, который будет объединять все словари - нужно для вывода на листе Надежности
        combined_dict_all_proc = {}
        combined_dict_all_proc.update(dict_inc_vnd)
        combined_dict_all_proc.update(dict_catm_vnd)
        combined_dict_all_proc.update(dict_slm_vnd)
        combined_dict_all_proc.update(dict_prb_vnd)
        combined_dict_all_proc.update(dict_erm_vnd)
        combined_dict_all_proc.update(dict_ichg_vnd)
        combined_dict_all_proc.update(dict_cont_vnd)
        combined_dict_all_proc.update(dict_avl_vnd)
        combined_dict_all_proc.update(dict_me_vnd)
        combined_dict_all_proc.update(dict_rff_vnd)
        combined_dict_all_proc.update(dict_cfg_vnd)
        combined_dict_all_proc.update(dict_upd_vnd)
        combined_dict_all_proc.update(dict_cap_vnd)
        combined_dict_all_proc.update(dict_mntw_vnd)
        combined_dict_all_proc.update(dict_sd_vnd)


    combined_dict_itsm = {}

    # Проходимся по внутренним словарям и суммируем значения
    for key in dict_inc_vnd['ITSM-INC']:
        value = (
                dict_inc_vnd['ITSM-INC'].get(key, 0) +
                dict_catm_vnd['ITSM-CATM'].get(key, 0) +
                dict_cfg_vnd['ITSM-CFG'].get(key, 0) +
                dict_slm_vnd['ITSM-SLM'].get(key, 0) +
                dict_prb_vnd['ITSM-PRB'].get(key, 0) +
                dict_erm_vnd['ITSM-ERM'].get(key, 0) +
                dict_rff_vnd['ITSM-RFF'].get(key, 0) +
                dict_ichg_vnd['ITSM-ICHG'].get(key, 0) +
                dict_upd_vnd['ITSM-UPD'].get(key, 0) +
                dict_avl_vnd['ITSM-AVL'].get(key, 0) +
                dict_me_vnd['ITSM-M&E'].get(key, 0) +
                dict_cont_vnd['ITSM-CONT'].get(key, 0))
                # dict_cap_vnd['ITSM-CAP'].get(key, 0) +
                # dict_mntw_vnd['ITSM-MNTW'].get(key, 0) +
                # dict_sd_vnd['ITSM-SD'].get(key, 0))

        combined_dict_itsm[key] = value  #Подсчет всех требований для 12 - ти ПРОЦЕССОВ для расчета индекса надежности и 3, 4-го уровней надежности!!!
    # print(combined_dict_itsm)




 # ******************************************** Этот код лучше исключить так как он берет информацию из Опросника и т.е.
    # все цифры по степени критичности - Это не нужно!!!

    # Инициализируем счетчик с нулевыми значениями
    counts_iz_oprosnika = {'Критичный': 0, 'Высокий': 0, 'Умеренный': 0, 'Рекомендация на развитие': 0}
    counts_iz_edk = {'Критичный': 0, 'Высокий': 0, 'Умеренный': 0, 'Рекомендация на развитие': 0}
    # Перебираем все записи в словаре
    for key in oprosnik_dict:
        # Второй элемент в списке - это уровень критичности
        criticality = oprosnik_dict[key][1]
        # Убираем возможные пробелы и приводим к строке
        criticality = str(criticality).strip()
        # Увеличиваем счетчик, если значение есть в нашем словаре
        if criticality in counts_iz_oprosnika:
            counts_iz_oprosnika[criticality] += 1

    for key in book_opros:
        # Второй элемент в списке - это уровень критичности
        criticality = book_opros[key][1]
        vipolnenie = book_opros[key][2]
        # Убираем возможные пробелы и приводим к строке
        criticality = str(criticality).strip()
        # Увеличиваем счетчик, если значение есть в нашем словаре
        if vipolnenie == 'да':
            if criticality in counts_iz_edk:
                counts_iz_edk[criticality] += 1
    #===============================================================================================================================================

    # Вычисляем вообще сколько требований в Опроснике первоначальном



    result_opros_itsm_proc = {'Всего требований': 0, 'Всего требований на 3': 0, 'Всего выполнено требований на 3': 0, 'Всего требований на 4': 0, 'Всего выполнено требований на 4': 0}
    for key, value_list in book_opros.items():
        if value_list[0] != 'ЦТК':
            if value_list[1] != 'Рекомендация на развитие':
                result_opros_itsm_proc['Всего требований'] += 1
                if value_list[0] == '3':
                    result_opros_itsm_proc['Всего требований на 3'] += 1
                if value_list[0] == '3' and value_list[2] == 'да':
                    result_opros_itsm_proc['Всего выполнено требований на 3'] += 1
                    # Добавьте здесь логику для подсчёта выполненных требований на 3, если нужно
                if value_list[0] == '4':
                    result_opros_itsm_proc['Всего требований на 4'] += 1
                if value_list[0] == '4' and value_list[2] == 'да':
                    result_opros_itsm_proc['Всего выполнено требований на 4'] += 1


    # print(combined_dict_itsm)

    # print(counts_iz_oprosnika) # ---> {'Критичный': 121, 'Высокий': 117, 'Умеренный': 26, 'Рекомендация на развитие': 39}
    # print(counts_iz_edk) # ---> {'Критичный': 37, 'Высокий': 25, 'Умеренный': 6, 'Рекомендация на развитие': 11}
    # print(combined_dict_itsm) # ---> {'Всего требований': 263, 'Всего требований на 3': 138, 'Всего выполнено требований на 3': 47, 'Всего требований на 4': 125, 'Всего выполнено требований на 4': 21}
    # print(combined_dict_all_proc) # ---> Для листа Надежности по всем процессам


    #Считаем из отчёта общего по ЕОИТ все отклонения(Технологии (ИТ-услуги)) + статусы если они будут в дальнейшем
    result_stack_as = {}
    for key, value_list in book_stata_uroven.items():
        result_stack_as[key] = [value_list[i:i+4] for i in range(0, len(value_list), 4)]


    def my_as_stata(result_stack_as):
        counts = {}
        counts_old = {}
        dic_itog = {}

        #Фильтрация по уровням отклонения среди всех АС, исключает уровни отклонения от ЦТК, т.к. коды вопросов не цифровые
        output_dictionary = defaultdict(list)
        # Фильтрация по первому элементу списка
        for key, value_list in result_stack_as.items():
            filtered_value_list = list(filter(lambda x: x[0].isdigit(), value_list))
            output_dictionary[key] = filtered_value_list


        risk_levels = ('Критичный','Высокий','Умеренный')
        for key, value_list in output_dictionary.items():

            level_counts = Counter([item[-2] for item in output_dictionary[key]])
            counts_old[key] = {risk_level: level_counts[risk_level] for risk_level in risk_levels}

        # print('===================================================================================')
        for key in output_dictionary:
            level_counts = Counter([item[-2] for item in output_dictionary[key] if
                                    not any(phrase in item for phrase in (
                                        'Устранено - подтверждено ПАО',
                                        'Исключение из реестра - подтверждено ПАО',
                                        '(ПАО) Устранено - подтверждено ПАО',
                                        '(ПАО) Исключение из реестра - подтверждено ПАО'))])

            # total_count = sum(level_counts.values())
            counts[key] = {risk_level: level_counts[risk_level]for risk_level in risk_levels}
        # print(counts)
        #Результат {'aos.smkt': ['11/6', '3/2', '1/0']}
        for key in counts_old:
            if key in counts:
                dic_itog[key] = []
                for level in counts_old[key]:
                    # Считаем разницу между значениями
                    diff = counts_old[key][level] - counts[key][level]

                    # Проверяем, равна ли разница нулю
                    if counts[key][level] == 0:
                        dic_itog[key].append(f"{counts_old[key][level]}/{counts_old[key][level]}")  # Используем diff, если разница равна 0


                    else:
                        dic_itog[key].append(f"{counts_old[key][level]}/{counts_old[key][level]-counts[key][level]}")

        return dic_itog


    new_as_statistic = my_as_stata(result_stack_as)

    # print(new_as_statistic) #---->'Shiptor ERP': ['83/71', '10/10', '0/0']}



    # myfunc_as -Собирает из ОТЧЕТА по ЕОИТ все номера отклонений, у которых статус пустой, для того чтобы сгенерить отчет и вытащить из базы знаний отклонения
    # genomy = ['Устранено - подтверждено ПАО', 'Исключение из реестра - подтверждено ПАО', '(ПАО) Устранено - подтверждено ПАО', '(ПАО) Исключение из реестра - подтверждено ПАО']
    # myfunc_as = {key: [item[0] for item in value if len(item) > 1 and item[1] not in genomy] for key, value in real_result_stack.items() if any(len(item) > 1 and item[1] not in genomy for item in value)}

    new_book_stac = [book_stack]


    # print(result_stack_as)



    # reeeez = (sublst for lst in result_stack_as.values() for sublst in lst)
    # if 'ПС-3' in (k[-1] for k in reeeez):
    #     print('ssdfsdfsdf')

    #Добавляем словарь, содержащий все АС, которые не попали в ЕДК, т.к. до формирования ЕРО были закрыты
    mas_it_uslug = {} # Словарь куда будут попадать АС-ки, которые не попали в ЕДК

    def find_and_open_excel(adress_as, result_stack_as, mas_it):
        as_dict = {}
        if not os.path.isdir(adress_as):
            print(f"Ошибка: папка '{adress_as}' не найдена. Возможно отсутствуют анкеты по ИТ-услугам")
            return

        for filename in os.listdir(adress_as):
            if filename.endswith(".xlsx"):
                file_path = os.path.join(adress_as, filename)
                try:
                    wb = openpyxl.load_workbook(file_path, data_only=True)
                    # print(f'Открываем файл - {filename}')
                    ws = wb['Титульный']
                    print(f'Работаем с файлом - {filename} ---> АС - {ws["C4"].value}')

                    cell_value_ps = ws['C8'].value
                    #Распарсим значения для поиска по прикладному сервису (ПС-1, ПС-2, ПС-3....)
                    real_as_relize = (sublst for lst in result_stack_as.values() for sublst in lst)
                    if cell_value_ps not in (k[-1] for k in real_as_relize):


                        #НУЖНО ПРАВИТЬ, т.к. ОШИБКА всего требований на 3 равно должно быть выполненных на 3 ну и на 4 равна 4
                        mas_it[ws['C4'].value] = [ws['C5'].value, round(ws['D20'].value, 2), round(ws['G20'].value, 2)]

                        print(f'АС -  {ws["C4"].value} отсутствует в файле ЕДК, возможно по ней нет отклонений')
                        # sys.exit()
                        continue

                    elif ws['C4'].value != 0:
                        as_dict[ws['C4'].value] = [ws['C5'].value, round(ws['D21'].value, 2), round(ws['F21'].value, 2)]
                    else:
                        print(f'Не заполненно поле -- Название ИТ-услуги -- в файле {filename}  ')
                        break

                except Exception as e:
                    print(f"Ошибка при открытии файла '{file_path}': {e}")



        return as_dict, mas_it

    rezult_as, mas_it_as = find_and_open_excel(adress_as, result_stack_as, mas_it_uslug)


    #Грузим ВСЕ АС КОТОРЫЕ ЕСТЬ В ПАПКЕ AS, они потом понадобяися для

    def all_func_as(adress_as):

        # new_book_stac - По этому словарю происходит поиск  элементов в папке AS, т.е. всех анкет
        # result_stack - Это то что загружается с Отчета_по_ITSM, получается словарь: {Технология: [[Код вопроса,Статус]]}
        # adress_as - Это адрес который передается в словарь для поиска анкет в папке АS
        #rezult_as - Это словарь где выводится АС: {'1С': ['MC', 2.85, 3.8], 'АС Аналитика': ['BO', 2.52, 3.45], 'АС Брокер': ['BC', 2.75, 3.65]}
        # new_result - Это словарь получаемый из файлов в папке AS, где Номер вопроса: 'Уровень готовности надёжности', 'Вес требования (ВНД 5492)'

        new_result = {}
        new_result_result_treb = {}

        if not os.path.isdir(adress_as):
            print(f"Ошибка: папка '{adress_as}' не найдена.")
            return

        for filename in os.listdir(adress_as):
            if filename.endswith(".xlsx"):
                file_path = os.path.join(adress_as, filename)
                try:
                    wb = openpyxl.load_workbook(file_path, data_only=True)
                    ws = wb['Технологии']
                    wm = wb['Титульный']
                    new_result[wm['C4'].value] = []
                    new_result_result_treb[wm['C4'].value] = []

                    for row in range(2, ws.max_row + 1):
                        keys = ws.cell(row, 1).value  # Столбец A
                        value_keys = wm['C8'].value  # Столбец F
                        value_l = ws.cell(row, 14).value  # Столбец N
                        new_result[wm['C4'].value].append((keys, value_keys, value_l, wm['C5'].value))
                except Exception as e:
                    print(f"Ошибка при открытии файла '{file_path}': {e}")

        clean_new_result = {}

        for key, value in new_result.items():
            if key is not None and (not isinstance(value, list) or not any(item is None or item == '-' for item in value)):
                clean_new_result[key] = value


        return clean_new_result

    if __name__ == "__main__":
        all_func_as(adress_as)

    #---------------------------------------------------------Вносим  и меняем значения УРОВЕНЬ ОТКЛОНЕНИЯ согласно загруженному словарю из общего списка---------------
    new_data = all_func_as(adress_as).copy()

    # Проходим по каждому ключу словаря
    for key in new_data:
        # Создаем новый список для хранения отфильтрованных и модифицированных данных
        new_list = []

        # Проходим по каждому элементу списка, соответствующего данному ключу
        for item in new_data[key]:
            # Если элемент содержит хотя бы одно значение None, пропускаем его
            if any(x is None for x in item):
                continue

            # Если второй элемент равен "Нет риска" и третий элемент равен 3, меняем третий элемент на "Критичный"

            if item[2] == 3:
                new_item = (item[0], item[1], 'Надежность 3', item[3])
            # Если второй элемент равен "Нет риска" и третий элемент равен 4, меняем третий элемент на "Умеренный"
            elif item[2] == 4:
                new_item = (item[0], item[1], 'Надежность 4', item[3])
            else:
                new_item = item
            # Добавляем модифицированный элемент в новый список
            new_list.append(new_item)
        # Обновляем значение в словаре новым списком
        new_data[key] = new_list


    # Функция для обновления словаря result_stack на основе данных из new_data
    def update_result_stack(result_stack, new_data):
        for key in result_stack.keys():
            if key in new_data:
                # Проходимся по всем значениям для данного ключа в result_stack
                for index, value in enumerate(result_stack[key]):
                    id_value = value[0]
                    # Преобразуем id_value в int, чтобы сравнить с ключами из new_data
                    if int(id_value) in {item[0] for item in new_data[key]}:
                        # Находим соответствующее значение в new_data
                        matching_item = next(item for item in new_data[key] if item[0] == int(id_value))
                        # Добавляем новые элементы в конец списка value
                        value.extend(matching_item[1:])

        return result_stack


    # Обновленный результат
    updated_result_stack = update_result_stack(result_stack, new_data)

    # print(result_opros_itsm_proc)
    # print(oprosnik_dict)

    #Выясняем сколько вообще требований было выполнено в Опроснике, до того как попасть в реестр ЕДК
    opros_treb_dtn = {'Всего выполнено требований на 3': 0, 'Всего выполнено требований на 4': 0}
    for key, value in oprosnik_dict.items():
        if value[0] == '3' and value[2] == 'да':
            opros_treb_dtn['Всего выполнено требований на 3'] += 1
        if value[0] == '4' and value[2] == 'да':
            opros_treb_dtn['Всего выполнено требований на 4'] += 1


    # print(opros_treb_dtn)
    #
    # print(combined_dict_itsm)
    # print(result_opros_itsm_proc)

    # Обновляем combined_dict_itsm значениями из opros_treb_dtn
    # combined_dict_itsm.update(opros_treb_dtn)



    def func_as_artefakt(result_stack, adress_as, rezult_as, realudate):

        # new_book_stac - По этому словарю происходит поиск  элементов в папке AS, т.е. всех анкет
        # result_stack - Это то что загружается с Отчета_по_ITSM, получается словарь: {Технология: [[Код вопроса,Статус]]}
        # adress_as - Это адрес который передается в словарь для поиска анкет в папке АS
        #rezult_as - Это словарь где выводится АС: {'1С': ['MC', 2.85, 3.8], 'АС Аналитика': ['BO', 2.52, 3.45], 'АС Брокер': ['BC', 2.75, 3.65]}
        # new_result - Это словарь получаемый из файлов в папке AS, где Номер вопроса: 'Уровень готовности надёжности', 'Вес требования (ВНД 5492)'

        new_result = {}
        new_result_result_treb = {}
        new_rez_as = rezult_as.copy()
        result_as_new = result_stack.copy()
        total_sum_all_treb_3 = 0
        total_sum_all_treb_4 = 0
        total_sum_release_treb_3 = 0
        total_sum_release_treb_4 = 0



        #aga_three и aga_four необходимы для добавления значений по уровням готовности для каждой АС
        aga_three = {}
        aga_four = {}

        if not os.path.isdir(adress_as):
            print(f"Ошибка: папка '{adress_as}' не найдена.")
            return

        for filename in os.listdir(adress_as):
            if filename.endswith(".xlsx"):
                file_path = os.path.join(adress_as, filename)
                try:
                    wb = openpyxl.load_workbook(file_path, data_only=True)
                    ws = wb['Технологии']
                    wm = wb['Титульный']
                    new_result[wm['C4'].value] = []
                    new_result_result_treb[wm['C4'].value] = []

                    for row in range(2, ws.max_row + 1):
                        keys = ws.cell(row, 1).value  # Столбец A

                        value_l = ws.cell(row, 14).value  # Столбец N
                        value_n = ws.cell(row, 16).value  # Столбец P

                        # new_result[keys] = [value_l, value_n]
                        #Вывод количество выполненных требований к общему числу из Титульной страницы
                        # Определяем диапазоны ячеек
                        if wm['E18'].value == 0 or wm['G18'] == 0:
                            range_1 = 'E10:E18'
                            range_3 = 'G10:G18'
                            range_2 = 'F10:F18'
                            range_4 = 'H10:H18'
                        else:
                            range_1 = 'E10:E17'
                            range_3 = 'G10:G17'
                            range_2 = 'F10:F17'
                            range_4 = 'H10:H17'

                        # Функция для получения списка значений из диапазона
                        def get_values_from_range(wm, range_str):
                            values = []
                            for row in wm[range_str]:
                                for cell in row:
                                    # Обрабатываем три случая: '-', None и другие значения
                                    if cell.value == '-' or cell.value is None:
                                        values.append(0)
                                    else:
                                        values.append(cell.value)
                            return values

                        # Получаем списки значений для обоих диапазонов
                        values_1 = get_values_from_range(wm, range_1)
                        values_2 = get_values_from_range(wm, range_3)
                        values_3 = get_values_from_range(wm, range_2)
                        values_4 = get_values_from_range(wm, range_4)

                        new_result[wm['C4'].value].append((keys, value_l, value_n))

                        new_result_result_treb[wm['C4'].value].append((sum(values_1),  sum(values_3), sum(values_2), sum(values_4)))


                except Exception as e:
                    print(f"Ошибка при открытии файла '{file_path}': {e}")




        #Получили реестр требований по всем АС для дальнейшего расчета и подсчета индекса надежности. Очень ВАЖНО!!!
        reestr_trebovan = {key: list(set(value)) for key, value in new_result_result_treb.items()}

        #Это нужно, чтобы удалить те АС-ки, которые не вошли в реестр ЕДК, чтобы не считать их данные мы их исключим!!!
        for id_key, vial_list in mas_it_as.items():
            if id_key in reestr_trebovan:
                del reestr_trebovan[id_key]

        clean_new_result = {}

        for key, value in new_result.items():
            if key is not None and (not isinstance(value, list) or not any(item is None or item == '-' for item in value)):
                clean_new_result[key] = value
        for keys in new_result:
            if keys in result_stack:
                for item in new_result[keys]:
                    # Проверяем, что item не None и имеет хотя бы один элемент
                    if item is not None and len(item) > 0:
                        # Проверяем, что первый элемент item существует и является строкой
                        # if isinstance(item[0], str) and item[0].isdigit():
                        if isinstance(item[0], (int, float)):
                            for item2 in result_stack[keys]:
                                # Проверяем, что item2 не None и имеет хотя бы один элемент
                                if item2 is not None and len(item2) > 0 and item2[0].isdigit():
                                    # Проверяем первый элемент item2
                                    if item[0] == int(item2[0]):


                                        if (item2[12] == 'Устранено - подтверждено ПАО' or
                                                item2[12] == 'Исключение из реестра - подтверждено ПАО' or
                                                item2[12] == '(ПАО) Устранено - подтверждено ПАО' or
                                                item2[12] == '(ПАО) Исключение из реестра - подтверждено ПАО'):
                                            if item[1] == 3:

                                                aga_three.setdefault(keys, []).append(round(item[2], 6))

                                            if item[1] == 4:

                                                aga_four.setdefault(keys, []).append(round(item[2], 6))


        #
        # print(aga_three)
        for key_three in aga_three:
            aga_three[key_three] = [sum(aga_three[key_three])]
        # print(aga_four)
        for key_four in aga_four:
            aga_four[key_four] = [sum(aga_four[key_four])]
        # print('3 ---->',aga_three)
        # print('4 ---->',aga_four)
        for new, val in new_rez_as.items():

            if new in aga_three:
                val[1] = val[1] + round(aga_three[new][0], 2)
            if new in aga_four:
                val[2] = val[2] + round(aga_four[new][0], 2)


        # print('3 ---->',aga_three)
        # print('4 ---->',aga_four)


        reality_dict_all = {}

        for ars_key, ars_val in reestr_trebovan.items():
            rs_val_three = aga_three.get(ars_key, None)
            rs_val_four = aga_four.get(ars_key, None)

            if ars_key in aga_three.keys() or ars_key in aga_four.keys():
                total_requirements = sum([float(x) for x in (ars_val[0][0], ars_val[0][2])])
                total_on_three = ars_val[0][0]
                total_on_four = ars_val[0][2]

                if rs_val_three is None:
                    completed_on_three = ars_val[0][1]
                else:
                    # print(ars_val[0][0], rs_val_three, ars_val[0][1])
                    if isinstance(rs_val_three[0], int):
                        completed_on_three = rs_val_three[0] + ars_val[0][1]
                    else:
                        completed_on_three = ars_val[0][0] * sum(rs_val_three) + ars_val[0][1]

                if rs_val_four is None:
                    completed_on_four = ars_val[0][3]
                else:
                    if isinstance(rs_val_four[0], int):
                        completed_on_four = rs_val_four[0] + ars_val[0][3]
                    else:
                        completed_on_four = ars_val[0][2] * sum(rs_val_four) + ars_val[0][3]



                reality_dict_all.setdefault(ars_key, []).append([
                    'Всего требований', total_requirements,
                    'Всего на 3-ку', total_on_three,
                    'Выполненных на 3-ку', completed_on_three,
                    'Всего на 4-ку', total_on_four,
                    'Выполненных на 4-ку', completed_on_four
                ])
                # print(reality_dict_all)
            else:

                reality_dict_all.setdefault(ars_key, []).append([
                    'Всего требований', sum([float(x) for x in (ars_val[0][0], ars_val[0][2])]),
                    'Всего на 3-ку', ars_val[0][0],
                    'Выполненных на 3-ку', ars_val[0][0] * sum(rs_val_three or []) + ars_val[0][1],
                    'Всего на 4-ку', ars_val[0][2],
                    'Выполненных на 4-ку', ars_val[0][2] * sum(rs_val_four or []) + ars_val[0][3]])

        #Собрали весь пул требований для расчета из словаря по АС(reality_dict_all) и по Процессам (updated_dict1)
        #В функции будем расчитывать индекс надежности и готовности по АС и по Процессам
        #1. Индекс надежности = (Сумма выполненных требований на 3-ку и 4-ку по процессам и АС)/ Сумма всех требований на 3-ку и 4-ку по процессам и АС



        index_nadejnost_proc = {'Всего требований по ITSM':(int(realudate['Всего требований'])),
                                                            'Всего требований по ITSM на 3': (int(realudate['Всего требований на 3'])),
                                                            'Выполненно требований по ITSM на 3': (int(realudate['Всего выполнено требований на 3'])),
                                                            'Всего требований по ITSM на 4':(int(realudate['Всего требований на 4'])),
                                                            'Выполненно требований по ITSM на 4': (int(realudate['Всего выполнено требований на 4']))
        }



        # for key_pr, val_pr in realudate.items():
        #
        #     index_nadejnost_proc.setdefault('Всего требований по ITSM', []).append(int(realudate['Всего требований']))
        #     index_nadejnost_proc.setdefault('Всего требований по ITSM на 3', []).append(int(realudate['Всего требований на 3']))
        #     index_nadejnost_proc.setdefault('Выполненно требований по ITSM на 3', []).append(int(realudate['Всего выполнено требований на 3']))
        #     index_nadejnost_proc.setdefault('Всего требований по ITSM на 4', []).append(int(realudate['Всего требований на 4']))
        #     index_nadejnost_proc.setdefault('Выполненно требований по ITSM на 4', []).append(int(realudate['Всего выполнено требований на 4']))




        index_nadejnost_as = {}
        for key_as, val_as in reality_dict_all.items():


            #Это блок отвечающий за АС
            index_nadejnost_as.setdefault('Всего требований по АС', []).append((val_as[0][1]))
            index_nadejnost_as.setdefault('Всего требования по АС на 3', []).append((val_as[0][3]))
            index_nadejnost_as.setdefault('Выполненно требований по АС на 3', []).append(val_as[0][5])
            index_nadejnost_as.setdefault('Всего требований по АС на 4', []).append((val_as[0][7]))
            index_nadejnost_as.setdefault('Выполненно требований по АС на 4', []).append(val_as[0][9])



        sum_index_nadejnost_proc = index_nadejnost_proc
        sum_index_nadejnost_as = {
            'Всего требований по АС': 0,
            'Всего требования по АС на 3': 0,
            'Выполненно требований по АС на 3': 0,
            'Всего требований по АС на 4': 0,
            'Выполненно требований по АС на 4': 0
        }
        # for key in index_nadejnost_proc:
        #
        #     sum_index_nadejnost_proc[key] = sum(index_nadejnost_proc[key])
        # print(index_nadejnost_as)
        for key_as in index_nadejnost_as:
            sum_index_nadejnost_as[key_as] = sum(index_nadejnost_as[key_as])


        # Извлекаем все значения из словаря
        values_mas_as = list(mas_it_as.values())

        # Оставляем только числовые значения
        filtered_values_mas_as = [[v for v in value if isinstance(v, (int, float))] for value in values_mas_as]

        # Рассчитываем суммы
        total_requirements = sum(sum(value) for value in filtered_values_mas_as)
        total_requirements_on_3 = sum(value[0] for value in filtered_values_mas_as)
        executed_requirements_on_3 = total_requirements_on_3
        total_requirements_on_4 = sum(value[1] for value in filtered_values_mas_as)
        executed_requirements_on_4 = total_requirements_on_4

        # Создаем новый словарь
        result_dict_close_as = {'Всего требований по АС': total_requirements,
                                'Всего требования по АС на 3': total_requirements_on_3,
                                'Выполненно требований по АС на 3': executed_requirements_on_3,
                                'Всего требований по АС на 4': total_requirements_on_4,
                                'Выполненно требований по АС на 4': executed_requirements_on_4
                                }

        # Обновление значений в словаре sum_index_nadejnost_as со словарем result_dict_close_as(это те AS которые не попали в ЕДК)
        for key in result_dict_close_as:
            sum_index_nadejnost_as[key] += result_dict_close_as[key]



        all_dict_as_itsm = {'Всего требований по ITSM и АС': [sum_index_nadejnost_proc['Всего требований по ITSM'], sum_index_nadejnost_as['Всего требований по АС']],
                            'Всего требований на 3 по ITSM и АС': [sum_index_nadejnost_proc['Всего требований по ITSM на 3'], sum_index_nadejnost_as['Всего требования по АС на 3']],
                            'Выполнено требований на 3 по ITSM и АС': [sum_index_nadejnost_proc['Выполненно требований по ITSM на 3'], round(sum_index_nadejnost_as['Выполненно требований по АС на 3'])],
                            'Всего требований на 4 по ITSM и АС': [sum_index_nadejnost_proc['Всего требований по ITSM на 4'], sum_index_nadejnost_as['Всего требований по АС на 4']],
                            'Выполнено требований на 4 по ITSM и АС': [sum_index_nadejnost_proc['Выполненно требований по ITSM на 4'], round(sum_index_nadejnost_as['Выполненно требований по АС на 4'])]
                            }

        # print(all_dict_as_itsm)
        # print(index_nadejnost_as)
        # print(index_nadejnost_proc)
        # print(realudate)


        return new_rez_as, all_dict_as_itsm

        # new_as_rezult = func_as_artefakt(result_stack, adress_as, rezult_as, updated_dict1)

    new_as_rezult, realudate = func_as_artefakt(result_stack, adress_as, rezult_as, result_opros_itsm_proc)

    # print(result_opros_itsm_proc) ---> {'Всего требований': 264, 'Всего требований на 3': 138, 'Всего выполнено требований на 3': 31, 'Всего требований на 4': 126, 'Всего выполнено требований на 4': 25}
    # print(combined_dict_itsm)
    #+++++++++++Это нужно для точно расчета, чтобы добавить эти АС в Единый реестр отклонений+++++++++++++
    # print(result_stack)   ---> {'Address Processing': [['1002001', 'ПС-2', 'Отсутствует локальный резерв ИТ-Услуги', 'Надежность 3', 'Предоставить ответ на вопрос.
    # print(rezult_as) #---> 'Shiptor ERP': ['MC', 2.02, 3]
    # print(new_as_rezult) #---> 'Shiptor ERP': ['MC', 2.02, 3]
    #Теперь нужно циклом перебрать все АС, которые не попали в ЕДК и добавить количество требований для этих АС

    # print(realudate)



    #=============================================== ---  Индекс надежности --- ===========================================================================

    index_dtn_for_ztk = {key: sum(value) for key, value in realudate.items()}
    #=============================================== ---  ---------- --- ===========================================================================

    # print(book_itsm)





    wb = Workbook()
    ws = wb.active
    sheet = wb.active

    ws.title = 'Единая дорожная карта'
    ws_help = wb.create_sheet(title='Справочник')
    second_workshet = wb.create_sheet(title='Показатели надежности')
    five_workshet = wb.create_sheet(title='Детальная статистика по АС')
    six_workshet = wb.create_sheet(title='Тепловая карта для 3-го уровня')
    seven_workshet = wb.create_sheet(title='Тепловая карта для 4-го уровня')



    # Определяем стиль границы
    thin_border = Border(left=Side(style='thin'),
                         right=Side(style='thin'),
                         top=Side(style='thin'),
                         bottom=Side(style='thin'))

    #Изменяем ячейки по высоте на втором активном листе
    second_workshet.row_dimensions[1].height = 50

    #Изменяем ячейки по высоте на первом листе
    ws.row_dimensions[1].height = 60


    #Статусы для листа "СПРАВОЧНИК"----------------------------------------------------------------------------------------------
    stata_help = {'Статус отклонения': ['(ДЗО) В работе - анализ',
                                        '(ДЗО) В работе – план устранения на проверку ПАО',
                                        '(ДЗО) В работе - реализация',
                                        '(ДЗО) Устранено - на проверку ПАО',
                                        '(ПАО) Устранено - подтверждено ПАО',
                                        '(ДЗО) Исключение из реестра - на проверку ПАО',
                                        '(ПАО) Исключение из реестра - подтверждено ПАО',
                                        '(ДЗО) Принятие рисков - на утверждение ПАО',
                                        '(ПАО) Принятие рисков - утверждено ПАО']}

    about_stata = {'Описание статуса':['ДЗО анализирует найденное отклонение и рекомендации',
                                       'ДЗО проработало необходимые мероприятия для устранения отклонения, ставит ДЗО',
                                       'ПАО подтверждает мероприятия, проработанные ДЗО для устранения отклонения, ДЗО реализует мероприятия',
                                       'ДЗО устранило отклонение, статус ставит ДЗО',
                                       'ПАО согласно с тем, что отклонение считается устраненным, статус ставит ДЗО',
                                       'ДЗО несогласно с отклонением и предлагает его исключить из списка на устранение, статус ставит ДЗО',
                                       'ПАО согласно с тем, что данное отклонение исключается из списка на устранение, статус ставит ПАО',
                                       'ДЗО принимает риски, связанные с отклонением и не будет его устранять, статус ставит ДЗО',
                                       'ПАО согласно с тем, что ДЗО принимает риски, связанные с отклонением и не будет его устранять, статус ставит ПАО']}


    # Заголовки
    ws_help['P1'].value = 'Статус отклонения'
    ws_help['Q1'].value = 'Описание статуса'

    # Форматируем заголовки
    ws_help['P1'].font = Font(name='Calibri', size=11, bold=True)
    ws_help['Q1'].font = Font(name='Calibri', size=11, bold=True)
    ws_help['P1'].font = Font(name='Calibri', size=12, bold=True)
    ws_help.column_dimensions['P'].width = 35
    ws_help['Q1'].font = Font(name='Calibri', size=12, bold=True)


    # Заполняем статусы отклонений
    start_row_zero = 2
    for status in stata_help['Статус отклонения']:
        cell = ws_help.cell(row=start_row_zero, column=16)
        cell.value = status
        cell.font = Font(name='Calibri', size=11, bold=False)
        start_row_zero += 1

    # Заполняем описания статусов
    start_row_zero = 2
    for description in about_stata['Описание статуса']:
        cell = ws_help.cell(row=start_row_zero, column=17)
        cell.value = description
        cell.font = Font(name='Calibri', size=11, bold=False)
        start_row_zero += 1

    #----------------------------------------Конец вывода статусов отклонений----------------------------------------

    # Скрываем второй лист
    ws_help.sheet_state = 'hidden'  # Возможные состояния: visible, hidden, veryHidden
    dv = DataValidation(type="list", formula1='=Справочник!$P$2:$P$10', allow_blank=True)
    ws.add_data_validation(dv)

    #Изменяем ячейки по ширине во второй вкладке ("Показатели надежности")

    second_workshet['A1'].font = Font(name='Calibri', size=12, bold=True)
    second_workshet.column_dimensions['A'].width = 70



    second_workshet['B1'].font = Font(name='Calibri',size=12, bold=True)
    second_workshet.column_dimensions['B'].width = 30

    second_workshet['C1'].font = Font(name='Calibri',size=12, bold=True)
    second_workshet.column_dimensions['C'].width = 30

    second_workshet['D1'].font = Font(name='Calibri',size=12, bold=True)
    second_workshet.column_dimensions['D'].width = 20


    second_workshet['E1'].font = Font(name='Calibri',size=12, bold=True)
    second_workshet.column_dimensions['E'].width = 20

    second_workshet['F1'].font = Font(name='Calibri',size=12, bold=True)
    second_workshet.column_dimensions['F'].width = 20

    second_workshet['G1'].font = Font(name='Calibri',size=12, bold=True)
    second_workshet.column_dimensions['G'].width = 20

    second_workshet['H1'].font = Font(name='Calibri',size=12, bold=True)
    second_workshet.column_dimensions['H'].width = 30

    second_workshet['I1'].font = Font(name='Calibri',size=12, bold=True)
    second_workshet.column_dimensions['I'].width = 30



    ws.column_dimensions['A'].width = 25
    ws['A1'].font = Font(name='Calibri', size=8, bold=True)
    ws.column_dimensions['B'].width = 24
    ws['B1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['C'].width = 20
    ws['C1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['D'].width = 25
    ws['D1'].font = Font(name='Calibri', size=8, bold=True)
    ws.column_dimensions['E'].width = 31
    ws['E1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['F'].width = 18
    ws['F1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['G'].width = 10
    ws['G1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['H'].width = 10
    ws['H1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['I'].width = 16
    ws['I1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['J'].width = 16
    ws['J1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['K'].width = 18
    ws['K1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['L'].width = 25
    ws['L1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['M'].width = 50
    ws['M1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['N'].width = 50
    ws['N1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['O'].width = 14
    ws['O1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['P'].width = 17
    ws['P1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['Q'].width = 20
    ws['Q1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['R'].width = 22
    ws['R1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['S'].width = 17
    ws['S1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['T'].width = 17
    ws['T1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['U'].width = 14
    ws['U1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['V'].width = 14
    ws['V1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['W'].width = 35
    ws['W1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['X'].width = 14
    ws['X1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['Y'].width = 22
    ws['Y1'].font = Font(name='Calibri',size=8, bold=True)
    ws.column_dimensions['Z'].width = 22
    ws['Z1'].font = Font(name='Calibri',size=8, bold=True)



    #Изменяем ячейки по ширине в 4-й вкладке ("Детальная статитстика по АС")

    five_workshet['A1'].font = Font(name='Calibri', size=12, bold=True)
    five_workshet.column_dimensions['A'].width = 70
    five_workshet['B1'].font = Font(name='Calibri',size=12, bold=True)
    five_workshet.column_dimensions['B'].width = 30
    five_workshet['C1'].font = Font(name='Calibri',size=12, bold=True)
    five_workshet.column_dimensions['C'].width = 30
    five_workshet['D1'].font = Font(name='Calibri',size=12, bold=True)
    five_workshet.column_dimensions['D'].width = 30
    five_workshet['E1'].font = Font(name='Calibri',size=12, bold=True)
    five_workshet.column_dimensions['E'].width = 30


    six_workshet['A1'].font = Font(name='Calibri', size=12, bold=True)
    six_workshet.column_dimensions['A'].width = 17
    six_workshet['B1'].font = Font(name='Calibri',size=12, bold=True)
    six_workshet.column_dimensions['B'].width = 17
    six_workshet['C1'].font = Font(name='Calibri',size=12, bold=True)
    six_workshet.column_dimensions['C'].width = 17
    six_workshet['D1'].font = Font(name='Calibri',size=12, bold=True)
    six_workshet.column_dimensions['D'].width = 17
    six_workshet['E1'].font = Font(name='Calibri',size=12, bold=True)
    six_workshet.column_dimensions['E'].width = 17
    six_workshet['F1'].font = Font(name='Calibri', size=12, bold=True)
    six_workshet.column_dimensions['F'].width = 17
    six_workshet['G1'].font = Font(name='Calibri',size=12, bold=True)
    six_workshet.column_dimensions['G'].width = 17
    six_workshet['H1'].font = Font(name='Calibri',size=12, bold=True)
    six_workshet.column_dimensions['H'].width = 17
    six_workshet['I1'].font = Font(name='Calibri',size=12, bold=True)
    six_workshet.column_dimensions['I'].width = 17
    six_workshet['J1'].font = Font(name='Calibri',size=12, bold=True)
    six_workshet.column_dimensions['J'].width = 17
    six_workshet['K1'].font = Font(name='Calibri', size=12, bold=True)
    six_workshet.column_dimensions['K'].width = 17
    six_workshet['L1'].font = Font(name='Calibri',size=12, bold=True)
    six_workshet.column_dimensions['L'].width = 17
    six_workshet['M1'].font = Font(name='Calibri',size=12, bold=True)
    six_workshet.column_dimensions['M'].width = 17
    six_workshet['N1'].font = Font(name='Calibri',size=12, bold=True)
    six_workshet.column_dimensions['N'].width = 17
    six_workshet['O1'].font = Font(name='Calibri',size=12, bold=True)
    six_workshet.column_dimensions['O'].width = 17
    six_workshet['P1'].font = Font(name='Calibri',size=12, bold=True)
    six_workshet.column_dimensions['P'].width = 17
    six_workshet['Q1'].font = Font(name='Calibri',size=12, bold=True)
    six_workshet.column_dimensions['Q'].width = 17
    six_workshet['R1'].font = Font(name='Calibri',size=12, bold=True)
    six_workshet.column_dimensions['R'].width = 17
    six_workshet['S1'].font = Font(name='Calibri',size=12, bold=True)
    six_workshet.column_dimensions['S'].width = 17
    six_workshet['T1'].font = Font(name='Calibri',size=12, bold=True)
    six_workshet.column_dimensions['T'].width = 17







    seven_workshet['A1'].font = Font(name='Calibri', size=12, bold=True)
    seven_workshet.column_dimensions['A'].width = 17
    seven_workshet['B1'].font = Font(name='Calibri',size=12, bold=True)
    seven_workshet.column_dimensions['B'].width = 17
    seven_workshet['C1'].font = Font(name='Calibri',size=12, bold=True)
    seven_workshet.column_dimensions['C'].width = 17
    seven_workshet['D1'].font = Font(name='Calibri',size=12, bold=True)
    seven_workshet.column_dimensions['D'].width = 17
    seven_workshet['E1'].font = Font(name='Calibri',size=12, bold=True)
    seven_workshet.column_dimensions['E'].width = 17
    seven_workshet['F1'].font = Font(name='Calibri', size=12, bold=True)
    seven_workshet.column_dimensions['F'].width = 17
    seven_workshet['G1'].font = Font(name='Calibri',size=12, bold=True)
    seven_workshet.column_dimensions['G'].width = 17
    seven_workshet['H1'].font = Font(name='Calibri',size=12, bold=True)
    seven_workshet.column_dimensions['H'].width = 17
    seven_workshet['I1'].font = Font(name='Calibri',size=12, bold=True)
    seven_workshet.column_dimensions['I'].width = 17
    seven_workshet['J1'].font = Font(name='Calibri',size=12, bold=True)
    seven_workshet.column_dimensions['J'].width = 17
    seven_workshet['K1'].font = Font(name='Calibri', size=12, bold=True)
    seven_workshet.column_dimensions['K'].width = 17
    seven_workshet['L1'].font = Font(name='Calibri',size=12, bold=True)
    seven_workshet.column_dimensions['L'].width = 17
    seven_workshet['M1'].font = Font(name='Calibri',size=12, bold=True)
    seven_workshet.column_dimensions['M'].width = 17
    seven_workshet['N1'].font = Font(name='Calibri',size=12, bold=True)
    seven_workshet.column_dimensions['N'].width = 17
    seven_workshet['O1'].font = Font(name='Calibri',size=12, bold=True)
    seven_workshet.column_dimensions['O'].width = 17
    seven_workshet['P1'].font = Font(name='Calibri',size=12, bold=True)
    seven_workshet.column_dimensions['P'].width = 17
    seven_workshet['Q1'].font = Font(name='Calibri',size=12, bold=True)
    seven_workshet.column_dimensions['Q'].width = 17
    seven_workshet['R1'].font = Font(name='Calibri',size=12, bold=True)
    seven_workshet.column_dimensions['R'].width = 17
    seven_workshet['S1'].font = Font(name='Calibri',size=12, bold=True)
    seven_workshet.column_dimensions['S'].width = 17
    seven_workshet['T1'].font = Font(name='Calibri',size=12, bold=True)
    seven_workshet.column_dimensions['T'].width = 17

    second_workshet['A1'] = 'Процесс'
    second_workshet['B1'] = 'Всего требований'
    second_workshet['C1'] = 'Требований на 3-й уровень'
    second_workshet['D1'] = 'Выполнено требований на 3-й уровень'
    second_workshet['E1'] = '%, достижения 3-го уровня'
    second_workshet['F1'] = 'Требований на 4-й уровень'
    second_workshet['G1'] = 'Выполнено требований на 4-й уровень'
    second_workshet['H1'] = '%, достижения 4-го уровня'
    # second_workshet['I1'] = 'Примечание'
    # Объединение ячеек F16:H16
    second_workshet.merge_cells('F16:H16')
    # Запись текста в объединенную ячейку
    second_workshet['F16'] = 'Прогноз уровня надежности'

    # Цвет для Критического уровня риска(подкрашиваем ячейки)
    high_risk_fill = PatternFill(start_color="FF9900", end_color="FF9900", fill_type="solid")
    high_profile_fill = PatternFill(start_color="9bbb59", end_color="9bbb59", fill_type="solid")


    five_workshet['A1'] = 'Название АС'
    five_workshet['B1'] = 'Критичность'
    five_workshet['C1'] = 'Уровень отклонения: Критичный'
    five_workshet['D1'] = 'Уровень отклонения: Высокий'
    five_workshet['E1'] = 'Уровень отклонения: Умеренный'









    # print(new_as_statistic)
    # print(rezult_as)


    detal_sttistics_as = copy.deepcopy(new_as_statistic)  # Создаем копию словаря

    for key in detal_sttistics_as:
        if key in rezult_as:  # Проверяем, есть ли ключ в dic2 (на всякий случай)
            detal_sttistics_as[key].insert(0, rezult_as[key][0])  # Вставляем первый элемент из dic2 в начало списка dic1



    result_as_dtn = {}

    for key, values in detal_sttistics_as.items():
        new_values = []
        for item in values:
            if isinstance(item, str) and '/' in item:  # Если элемент содержит '/'
                left, right = item.split('/')
                new_values.append((int(left), int(right)))  # Преобразуем в кортеж чисел
            else:
                new_values.append(item)  # Иначе оставляем как есть ('BC', 'MC' и т. д.)
        result_as_dtn[key] = new_values



    # Словарь для подсчета и вывода информации на странице === --- Детальная статистика по АС --- ===


    #Заполняем ЛИСТ ИТ-ТЕХНОЛОГИЙ
    count_as = 0
    for i, (key, values) in enumerate(result_as_dtn.items(), start=2):
        count_as += 1

        five_workshet[f'A{i}'] = key
        five_workshet[f'B{i}'] = values[0]
        five_workshet[f'C{i}'] = values[1][0]
        five_workshet[f'D{i}'] = values[2][0]
        five_workshet[f'E{i}'] = values[3][0]

        five_workshet[f"B{i}"].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        five_workshet[f"C{i}"].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        five_workshet[f"D{i}"].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        five_workshet[f"E{i}"].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)



    five_workshet[f'B{count_as + 8}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    five_workshet[f'A{count_as + 4}'] = 'Динамические показатели'
    five_workshet[f'A{count_as + 5}'] = 'АС'
    five_workshet[f'A{count_as + 7}'] = 'Доп. показатели / инфо'
    five_workshet[f'A{count_as + 8}'] = 'На слайд, в графу "Уровень технологической надежности" ='


    #Убрал, проверку пока все компании не дойдут до 3-го уровня надежности!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!
    try:
        if realudate['Выполнено требований на 3 по ITSM и АС'][1] == realudate['Всего требований на 3 по ITSM и АС'][1] and realudate['Всего требований на 3 по ITSM и АС'][1] != 0:
            value_four = round(realudate['Выполнено требований на 4 по ITSM и АС'][1] / realudate['Всего требований на 4 по ITSM и АС'][1], 2) + 3
            five_workshet[f'B{count_as + 8}'] = f"{value_four:.2f}".replace('.', ',')
        else:
            value_three = round(realudate['Выполнено требований на 3 по ITSM и АС'][1] / realudate['Всего требований на 3 по ITSM и АС'][1], 2) + 2
            five_workshet[f'B{count_as + 8}'] = f"{value_three:.2f}".replace('.', ',')
    except ZeroDivisionError:
        five_workshet[f'B{count_as + 8}'] = "3"  # Обработка деления на ноль

    five_workshet[f'B{count_as + 4}'] = 'Критичное'
    five_workshet[f'C{count_as + 4}'] = 'Высокое'
    five_workshet[f'D{count_as + 4}'] = 'Умеренное'
    five_workshet[f'A{count_as + 8}'].font = Font(name='Calibri', size=11, bold=True)
    five_workshet[f'B{count_as + 8}'].font = Font(name='Calibri', size=11, bold=True)


    # Эта часть кода  для Подсчета ОБЩЕГО КОЛИЧЕСТВА ОТКЛОНЕНИЙ ПО ВСЕМ АС ===================> на странице Детальная Статистика
    # Итоговый словарь

    new_statistics_as = {}
    # Инициализация счетчиков
    critical = [0]
    high = [0]
    moderate = [0]
    critical_result = [0]
    high_result = [0]
    moderate_result = [0]


    # Проходим по каждому элементу исходного словаря
    for key, values in result_as_dtn.items():
        # Пропускаем элементы, которые не содержат '/' в первом элементе
        # Критичные
        critical.append(values[1][0])
        high.append(values[2][0])
        moderate.append(values[3][0])
        # Выполненные
        critical_result.append(values[1][1])
        high_result.append(values[2][1])
        moderate_result.append(values[3][1])

    # Заполнение итогового словаря
    new_statistics_as['Критичное'] = sum(critical)
    new_statistics_as['Высокое'] = sum(high)
    new_statistics_as['Умеренное'] = sum(moderate)
    new_statistics_as['Критичное выполненно'] = sum(critical_result)
    new_statistics_as['Высокое выполненно'] = sum(high_result)
    new_statistics_as['Умеренное выполненно'] = sum(moderate_result)

    five_workshet[f"B{count_as + 5}"] = str(new_statistics_as['Критичное']) + ' / ' + str(new_statistics_as['Критичное выполненно'])
    five_workshet[f"C{count_as + 5}"] = str(new_statistics_as['Высокое']) + ' / ' + str(new_statistics_as['Высокое выполненно'])
    five_workshet[f"D{count_as + 5}"] = str(new_statistics_as['Умеренное']) + ' / ' + str(new_statistics_as['Умеренное выполненно'])



    five_workshet[f"B{count_as + 5}"].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    five_workshet[f"C{count_as + 5}"].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    five_workshet[f"D{count_as + 5}"].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    five_workshet[f"E{count_as + 5}"].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    five_workshet[f"B{count_as + 5}"].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    five_workshet[f"B{count_as + 4}"].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    five_workshet[f"C{count_as + 4}"].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    five_workshet[f"D{count_as + 4}"].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    five_workshet[f"E{count_as + 4}"].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)





    six_workshet['A1'] = 'Компания'
    six_workshet['B1'] = 'Общий уровень надежности'
    six_workshet['C1'] = 'Уровень технологической надежности'
    six_workshet['D1'] = 'Управление технологическими инцидентами'
    six_workshet['E1'] = 'Управление каталогом сервисов'
    six_workshet['F1'] = 'Управление уровнями сервисов'
    six_workshet['G1'] = 'Управление технологическими проблемами'
    six_workshet['H1'] = 'Управление технологическими рисками'
    six_workshet['I1'] = 'Управление инфраструктурными изменениями'
    six_workshet['J1'] = 'Управление непрерывностью технологий'
    six_workshet['K1'] = 'Управление доступностью технологий'
    six_workshet['L1'] = 'Управление технологическими событиями и мониторингом'
    six_workshet['M1'] = 'Управление технологическими запросами на обслуживание'
    six_workshet['N1'] = 'Управление конфигурациями'
    six_workshet['O1'] = 'Управление обновлением ПО'
    six_workshet['P1'] = 'Ключевые проблемы'
    six_workshet['Q1'] = 'Динамические показатели ITSM'
    six_workshet['R1'] = 'Динамические показатели по технологиям'
    six_workshet['S1'] = 'Уровень технологической надежности'
    six_workshet['T1'] = 'Индекс надежности(%)'


    seven_workshet['A1'] = 'Компания'
    seven_workshet['B1'] = 'Общий уровень надежности'
    seven_workshet['C1'] = 'Уровень технологической надежности'
    seven_workshet['D1'] = 'Управление технологическими инцидентами'
    seven_workshet['E1'] = 'Управление каталогом сервисов'
    seven_workshet['F1'] = 'Управление уровнями сервисов'
    seven_workshet['G1'] = 'Управление технологическими проблемами'
    seven_workshet['H1'] = 'Управление технологическими рисками'
    seven_workshet['I1'] = 'Управление инфраструктурными изменениями'
    seven_workshet['J1'] = 'Управление непрерывностью технологий'
    seven_workshet['K1'] = 'Управление доступностью технологий'
    seven_workshet['L1'] = 'Управление технологическими событиями и мониторингом'
    seven_workshet['M1'] = 'Управление технологическими запросами на обслуживание'
    seven_workshet['N1'] = 'Управление конфигурациями'
    seven_workshet['O1'] = 'Управление обновлением ПО'
    seven_workshet['P1'] = 'Ключевые проблемы'
    seven_workshet['Q1'] = 'Динамические показатели ITSM'
    seven_workshet['R1'] = 'Динамические показатели по технологиям'
    seven_workshet['S1'] = 'Уровень технологической надежности'
    seven_workshet['T1'] = 'Индекс надежности(%)'



    ws['A1'] = 'Код задачи/отклонения'
    ws['B1'] = 'Вид объекта 1 уровень'
    ws['C1'] = 'Вид объекта 2 уровень'
    ws['D1'] = 'Вид объекта 3 уровень'
    ws['E1'] = 'Категория задачи / отклонения'
    ws['F1'] = 'Дополнительный комментарий по задаче / выявленному отклонению'
    ws['G1'] = '№'
    ws['H1'] = 'Инициатор'
    ws['I1'] = 'Основание'
    ws['J1'] = 'Код вопроса ЕОИТ'
    ws['K1'] = 'Код объекта'
    ws['L1'] = 'Объект'
    ws['M1'] = 'Задача / выявленное отклонение'
    # ws['N1'] = 'Негативные последствия'
    ws['N1'] = 'Рекомендованные мероприятия'
    ws['O1'] = 'Уровень принятия решения'
    ws['P1'] = 'Вхождение в уровни готовности'
    ws['Q1'] = 'Уровень критичности'

    ws['R1'] = 'Запланированное мероприятие'
    ws['S1'] = 'Код мероприятия из трекера компании Группы (при использовании трекера)'
    ws['T1'] = 'Дата начала мероприятий'
    ws['U1'] = 'Дата окончания мероприятий'
    ws['V1'] = 'Ответственный'
    ws['W1'] = 'Статус'
    ws['X1'] = 'Комментарии компании Группы'
    ws['Y1'] = 'Комментарии инициатора'


    # Записываем данные на первую страницу отклонения по AS
    rezult_book_stack = copy.deepcopy(book_stack)
    chunk_size = 16  # Каждая запись содержит 18 элементов
    for key in rezult_book_stack:
        original_list = rezult_book_stack[key]
        # Разбиваем список на подсписки по chunk_size элементов
        split_list = [original_list[i:i + chunk_size]
                      for i in range(0, len(original_list), chunk_size)]
        rezult_book_stack[key] = split_list  # Заменяем исходный список разбитым


    start_i = 0
    count = 0

    for i, (key, values) in enumerate(book_itsm.items(), start=2):
        ws[f'L{i}'] = values[14]
        ws[f'K{i}'] = values[0]
        ws[f'I{i}'] = "ДТН"
        ws[f'O{i}'] = "Уровень 1"
        ws[f'B{i}'] = values[15]
        ws[f'W{i}'] = values[11]
        dv.add(ws[f'W{i}']) # Добаввляем стату в колонку W
        ws[f'G{i}'] = (int(f'{i}') - 1)
        ws[f'J{i}'] = key     #Начинаем писать данные со столбца --- F --- в ЕОИТ

        ws[f'M{i}'] = ''.join(values[1])
        ws[f'N{i}'] = ''.join(values[3])
        ws[f'P{i}'] = values[2]
        ws[f'Q{i}'] = values[5]


        # Если встречаем дату с временем, удаляем время '00:00:00'
        if isinstance(values[8], str) and '00:00:00' in values[8]:
            ws[f'T{i}'] = values[8].split()[0]
        else:
            ws[f'T{i}'] = values[8]
        # Если встречаем дату с временем, удаляем время '00:00:00'
        if isinstance(values[9], str) and '00:00:00' in values[9]:
            ws[f'U{i}'] = values[9].split()[0]
        else:
            ws[f'U{i}'] = values[9]

        ws[f'V{i}'] = values[10]
        ws[f'X{i}'] = values[12]
        ws[f'Y{i}'] = values[13]


        #Изменяем шрифт у активных полей

        ws[f'B{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'C{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'E{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'F{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'G{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'H{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'I{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'J{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'K{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'L{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'M{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'N{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'O{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'P{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'Q{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'R{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'S{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'T{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'U{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'V{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'X{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'Y{i}'].font = Font(name='Calibri', size=8, bold=False)
        ws[f'W{i}'].font = Font(name='Calibri', size=8, bold=False)

        #Изменяем расположение активных полей

        ws[f'B{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'C{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'E{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'F{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'G{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'H{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'I{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'J{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'K{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'L{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'M{i}'].alignment = Alignment(horizontal="left", vertical="top", wrapText=True)
        ws[f'N{i}'].alignment = Alignment(horizontal="left", vertical="top", wrapText=True)
        ws[f'O{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'P{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'S{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'T{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'Q{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'R{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'U{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'V{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'X{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'Y{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        ws[f'W{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        start_i = i

    #updated_result_stack  #---> Заменил rezult_book_stack
    for i, (key, values) in enumerate(updated_result_stack.items(), start=start_i):
        for j in values:
            count += 1
            ws[f'G{start_i+count}'] = (int(f'{start_i+count}') - 1)
            ws[f'I{start_i+count}'] = "ДТН"
            ws[f'B{start_i+count}'] = 'АС'
            ws[f'O{start_i+count}'] = "Уровень 1"
            ws[f'J{start_i+count}'] = j[0]

            ws[f'W{start_i+count}'] = j[12]
            dv.add(ws[f'W{start_i+count}']) # Добаввляем стату в колонку W
            ws[f'K{start_i+count}'] = j[1]
            ws[f'L{start_i+count}'] = key
            ws[f'M{start_i+count}'] = j[2]
            # ws[f'N{start_i+count}'] = j[3] ----> Негативные последствия убрали из ЕОИТ
            ws[f'N{start_i+count}'] = j[4]
            ws[f'P{start_i+count}'] = j[3]
            ws[f'Q{start_i+count}'] = j[6]


            # Если встречаем дату с временем, удаляем время '00:00:00'
            ws[f'T{start_i + count}'] = j[9]
            ws[f'U{start_i + count}'] = j[10]

            # if isinstance(j[9], str) and '00:00:00' in j[9]:
            #     ws[f'T{start_i+count}'] = j[9].split()[0]
            # else:
            #     ws[f'T{start_i + count}'] = j[9]
            # # Если встречаем дату с временем, удаляем время '00:00:00'
            # if isinstance(j[10], str) and '00:00:00' in j[10]:
            #     ws[f'U{start_i+count}'] = j[10].split()[0]
            # else:
            #     ws[f'U{start_i + count}'] = j[11]

            ws[f'V{start_i+count}'] = j[11]
            ws[f'X{start_i+count}'] = j[13]
            ws[f'Y{start_i+count}'] = j[15]


            #Изменяем шрифт у активных полей
            ws[f'B{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'C{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'E{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'F{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'G{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'H{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'I{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'J{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'K{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'L{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'M{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'N{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'O{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'P{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'Q{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'R{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'S{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'T{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'U{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'V{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'Y{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'X{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)
            ws[f'W{start_i+count}'].font = Font(name='Calibri', size=8, bold=False)

            #Изменяем расположение активных полей
            ws[f'B{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'C{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'E{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'F{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'G{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'H{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'I{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'J{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'K{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'L{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'M{start_i+count}'].alignment = Alignment(horizontal="left", vertical="top", wrapText=True)
            ws[f'N{start_i+count}'].alignment = Alignment(horizontal="left", vertical="top", wrapText=True)
            ws[f'O{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'P{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'S{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'Q{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'R{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'T{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'U{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'V{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'Y{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'X{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            ws[f'W{start_i+count}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)


    # Первая таблица
    for col in 'ABCDEFGHIJKLMNOPQRSTUVWXYZ':
        for row in ['1']:
            cell = f'{col}{row}'
            ws[cell].alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)




    for i, (key, values) in enumerate(combined_dict_all_proc.items(), start=2):
        if key not in ['ITSM-CAP', 'ITSM-SD', 'ITSM-MNTW']:

            second_workshet[f'A{i}'] = dict_itsm[key]
            second_workshet[f'B{i}'] = combined_dict_all_proc[key]['Всего требований']
            second_workshet[f'C{i}'] = combined_dict_all_proc[key]['Всего требований на 3']
            second_workshet[f'D{i}'] = combined_dict_all_proc[key]['Всего выполнено требований на 3']

            if combined_dict_all_proc[key]['Всего требований на 3'] == 0:
                second_workshet[f'E{i}'] = 0
            else:
                second_workshet[f'E{i}'] = f"{round(second_workshet[f'D{i}'].value / second_workshet[f'C{i}'].value, 2) * 100:.0f}%"
            second_workshet[f'F{i}'] = combined_dict_all_proc[key]['Всего требований на 4']
            second_workshet[f'G{i}'] = combined_dict_all_proc[key]['Всего выполнено требований на 4']
            if combined_dict_all_proc[key]['Всего требований на 4'] == 0:
                second_workshet[f'H{i}'] = 0
            else:

                second_workshet[f'H{i}'] = f"{round(second_workshet[f'G{i}'].value / second_workshet[f'F{i}'].value, 2) * 100:.0f}%"


        second_workshet[f'B{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        second_workshet[f'C{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        second_workshet[f'D{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        second_workshet[f'E{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        second_workshet[f'F{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        second_workshet[f'G{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        second_workshet[f'H{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
        second_workshet[f'I{i}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)

    # Цвет для Критического уровня риска(подкрашиваем ячейки)
    high_risk_fill = PatternFill(start_color="FF9900", end_color="FF9900", fill_type="solid")
    high_profile_fill = PatternFill(start_color="9bbb59", end_color="9bbb59", fill_type="solid")

    # Динамические показатели
    second_workshet['A16'] = 'Выбран профиль'
    second_workshet['B16'].fill = high_risk_fill
    second_workshet['B16'] = insurer_key_profile
    second_workshet['A17'] = 'Динамические показатели'
    second_workshet['A18'] = 'ITSM'
    second_workshet['B17'] = 'Критичное'

    second_workshet['B18'] = str(itsm_dtn_all['Критичный']) + ' / ' + str(itsm_dtn_ikrement['Критичный'])
    second_workshet['C17'] = 'Высокое'
    second_workshet['C18'] = str(itsm_dtn_all['Высокий']) + ' / ' + str(itsm_dtn_ikrement['Высокий'])
    second_workshet['D17'] = 'Умеренное'
    second_workshet['D18'] = str(itsm_dtn_all['Умеренный']) + ' / ' + str(itsm_dtn_ikrement['Умеренный'])
    second_workshet['A20'] = 'На слайд, в графу "Уровень надёжности" ='
    second_workshet['A22'] = 'Индекс надежности для ЦТК и тепловой карты'


    index_dtn_for_ztk = {key: sum(value) for key, value in realudate.items()}
    index_nadejnost_all = (index_dtn_for_ztk['Выполнено требований на 3 по ITSM и АС'] +
                           index_dtn_for_ztk['Выполнено требований на 4 по ITSM и АС']) / (index_dtn_for_ztk['Всего требований на 3 по ITSM и АС'] +
                                                                                           index_dtn_for_ztk['Всего требований на 4 по ITSM и АС']) * 100

    second_workshet['B22'] = str(round(index_nadejnost_all, 0)) + '%'
    three_uroven_nadejnosti = index_dtn_for_ztk['Выполнено требований на 3 по ITSM и АС']/index_dtn_for_ztk['Всего требований на 3 по ITSM и АС'] + 2
    four_uroven_nadejnosti = index_dtn_for_ztk['Выполнено требований на 4 по ITSM и АС']/index_dtn_for_ztk['Всего требований на 4 по ITSM и АС'] + 3
    second_workshet['A25'] = '3-й уровень готовности  надежности'
    second_workshet['A26'] = '4-й уровень готовности  надежности'

    second_workshet['B25'] = float(round(three_uroven_nadejnosti, 2))
    second_workshet['B26'] = float(round(four_uroven_nadejnosti, 2))
    second_workshet['B20'] = float(second_workshet['B25'].value)



    # Шестая таблица
    for col in 'ABCDEFGHIJKLMNOP':
        for row in ['1', '2']:
            cell = f'{col}{row}'
            six_workshet[cell].alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)

    # Седьмая таблица
    for col in 'ABCDEFGHIJKLMNOP':
        for row in ['1', '2']:
            cell = f'{col}{row}'
            seven_workshet[cell].alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)







        # Функция для замены % в значении, если оно есть
    def remove_percent(value):
        if isinstance(value, str):  # Проверяем, является ли значение строкой
            return value.replace('%', '')  # Удаляем символ %
        return value  # Если это не строка, возвращаем значение как есть

    six_workshet['A2'] = name_dtn
    six_workshet['B2'] = second_workshet['B25'].value
    six_workshet['C2'] = five_workshet[f'B{count_as + 8}'].value
    if realudate['Выполнено требований на 3 по ITSM и АС'][1] == realudate['Всего требований на 3 по ITSM и АС'][1]:
        # six_workshet['C2'] = five_workshet[f'B{count_as + 8}'].value
        value = round(realudate['Выполнено требований на 3 по ITSM и АС'][1] / realudate['Всего требований на 3 по ITSM и АС'][1], 2) + 3
    else:
        value = round(realudate['Выполнено требований на 3 по ITSM и АС'][1] / realudate['Всего требований на 3 по ITSM и АС'][1], 2) + 2
        six_workshet['C2'] = f"{value:.2f}".replace('.', ',')  # Заменяем запятую на точку

    six_workshet['D2'] = remove_percent(second_workshet['E2'].value)
    six_workshet['E2'] = remove_percent(second_workshet['E3'].value)
    six_workshet['F2'] = remove_percent(second_workshet['E4'].value)
    six_workshet['G2'] = remove_percent(second_workshet['E5'].value)


    if dict_erm_vnd['ITSM-ERM']['Всего требований на 3'] != 0:
        six_workshet['H2'] = remove_percent(second_workshet['E6'].value)
    else:
        six_workshet['H2'] = '-'

    six_workshet['I2'] = remove_percent(second_workshet['E7'].value)

    if dict_cont_vnd['ITSM-CONT']['Всего требований на 3'] != 0:
        six_workshet['J2'] = remove_percent(second_workshet['E8'].value)
    else:
        six_workshet['J2'] = '-'

    if dict_avl_vnd['ITSM-AVL']['Всего требований на 3'] != 0:
        six_workshet['K2'] = remove_percent(second_workshet['E9'].value)
    else:
        six_workshet['K2'] = '-'
    six_workshet['L2'] = remove_percent(second_workshet['E10'].value)
    six_workshet['M2'] = remove_percent(second_workshet['E11'].value)
    six_workshet['N2'] = remove_percent(second_workshet['E12'].value)

    if dict_upd_vnd['ITSM-UPD']['Всего требований на 3'] != 0:
        six_workshet['O2'] = remove_percent(second_workshet['E13'].value)
    else:
        six_workshet['O2'] = '-'





    seven_workshet['A2'] = name_dtn
    seven_workshet['B2'] = second_workshet['B26'].value

    if realudate['Выполнено требований на 4 по ITSM и АС'][1] == realudate['Всего требований на 4 по ITSM и АС'][1]:
        seven_workshet['C2'] = '4'
    else:
        seven_workshet['C2'] = five_workshet[f'B{count_as + 8}'].value
        # value = round(realudate['Выполнено требований на 4 по ITSM и АС'][1] / realudate['Всего требований на 4 по ITSM и АС'][1], 2) + 3
        # seven_workshet['C2'] = f"{value:.2f}".replace('.', ',')  # Заменяем запятую на точку

    seven_workshet['D2'] = remove_percent(second_workshet['H2'].value)
    seven_workshet['E2'] = remove_percent(second_workshet['H3'].value)
    seven_workshet['F2'] = remove_percent(second_workshet['H4'].value)
    seven_workshet['G2'] = remove_percent(second_workshet['H5'].value)
    seven_workshet['H2'] = remove_percent(second_workshet['H6'].value)
    seven_workshet['I2'] = remove_percent(second_workshet['H7'].value)
    seven_workshet['J2'] = remove_percent(second_workshet['H8'].value)
    seven_workshet['K2'] = remove_percent(second_workshet['H9'].value)
    seven_workshet['L2'] = remove_percent(second_workshet['H10'].value)
    seven_workshet['M2'] = remove_percent(second_workshet['H11'].value)
    seven_workshet['N2'] = remove_percent(second_workshet['H12'].value)
    seven_workshet['O2'] = remove_percent(second_workshet['H13'].value)



    #==================== !!!!!!!!  ПРОГНОЗИРОВАНИЕ  !!!!!!!!!!!!!!! ===========================================================================================

    # print(sorted_quarters)
    # print(combined_dict_all_proc)
    # print(combined_dict_itsm)
    # print(result_opros_itsm_proc)
    # print(sorted_quarters)
    # print(result_opros_itsm_proc)
    # uroven_tech_nadejnosti_prognoz = copy.deepcopy(result_opros_itsm_proc)
    # Инициализация структуры для прогноза


    # Статусы для исключения
    excluded_statuses = {
        '(ПАО) Устранено - подтверждено ПАО',
        '(ПАО) Исключение из реестра - подтверждено ПАО'
    }

    #ОБРАБОТКА ПО ITSM===========================================
    # 1. Собираем выполненные требования по кварталам
    quarterly_counts = defaultdict(lambda: {'3': 0, '4': 0})

    for req_id, values in book_itsm.items():

        quarter = values[9]  # Квартал выполнения
        status = values[11]  # Статус
        recomended = values[5] # Уровень отклонения(Критичный, Высокий, Рекомендация на развитие)
        nadejnost = values[2]  # "Надежность 3" или "Надежность 4"


        # Пропускаем, если квартал не в списке
        if quarter not in sorted_quarters:
            continue  # ← просто игнорируем, ничего не делаем

        # Пропускаем, если статус исключён
        if status in excluded_statuses:
            continue

        # Пропускаем, если критичность в исключении
        if recomended == 'Рекомендация на развитие':
            continue

        # Учитываем только если квартал в списке и статус не исключён
        # if quarter in sorted_quarters and status not in excluded_statuses:
        #     if nadejnost == 'Надежность 3':
        #         quarterly_counts[quarter]['3'] += 1
        #     elif nadejnost == 'Надежность 4':
        #         quarterly_counts[quarter]['4'] += 1

        if nadejnost == 'Надежность 3':
            quarterly_counts[quarter]['3'] += 1
        elif nadejnost == 'Надежность 4':
            quarterly_counts[quarter]['4'] += 1

    # 2. Формируем прогноз с контролем лимитов
    current_state = copy.deepcopy(result_opros_itsm_proc)
    uroven_tech_nadejnosti_prognoz = {}

    for quarter in sorted_quarters:
        # Получаем выполнения для текущего квартала
        counts = quarterly_counts[quarter]

        # Рассчитываем сколько МОЖНО добавить, не превышая лимиты
        add_3 = min(counts['3'], current_state['Всего требований на 3'] - current_state['Всего выполнено требований на 3'])
        add_4 = min(counts['4'], current_state['Всего требований на 4'] - current_state['Всего выполнено требований на 4'])

        # Обновляем состояния
        current_state['Всего выполнено требований на 3'] += add_3
        current_state['Всего выполнено требований на 4'] += add_4

        # Сохраняем результат
        uroven_tech_nadejnosti_prognoz[quarter] = copy.deepcopy(current_state)
    # 3. Вывод результатов
    # for quarter in sorted_quarters:
    #     print(f"{quarter}: {uroven_tech_nadejnosti_prognoz[quarter]}")
     #ОБРАБОТКА ПО ИТ-УСЛУГАМ===========================================



    versioned_data_it_sm = {}
    versioned_data_it_as = {}
    current_state_it_sm = copy.deepcopy(book_itsm)  # Текущее состояние данных
    current_state_it_as = copy.deepcopy(updated_result_stack)  # Текущее состояние данных

    for quarter in sorted_quarters:
        # Создаем копию для текущего квартала
        quarter_state = copy.deepcopy(current_state_it_sm)
        quarter_state_as = copy.deepcopy(current_state_it_as)

        for key, value in quarter_state.items():
            # Проверяем, что это запись текущего квартала
            if value[9] == quarter:
                value[11] = "(ПАО) Устранено - подтверждено ПАО"


        for i, v in quarter_state_as.items():
            for entry in v:
                if entry[10] == quarter:
                    entry[12] = "(ПАО) Устранено - подтверждено ПАО"

        # Сохраняем состояние после обработки текущего квартала
        versioned_data_it_sm[quarter] = quarter_state
        versioned_data_it_as[quarter] = quarter_state_as
        # Обновляем current_state_it_sm для следующей итерации
        current_state_it_sm = copy.deepcopy(quarter_state)
        current_state_it_as = copy.deepcopy(quarter_state_as)




    # print(versioned_data_it_as)
    count_pr = 0
    for quarter in sorted_quarters:
        count_pr += 1
        def func_as_artefakt_prognoz(result_stack, adress_as, rezult_as, realudate, quarter):

            # print(quarter, '----')
            # new_book_stac - По этому словарю происходит поиск  элементов в папке AS, т.е. всех анкет
            # result_stack - Это то что загружается с Отчета_по_ITSM, получается словарь: {Технология: [[Код вопроса,Статус]]}
            # adress_as - Это адрес который передается в словарь для поиска анкет в папке АS
            #rezult_as - Это словарь где выводится АС: {'1С': ['MC', 2.85, 3.8], 'АС Аналитика': ['BO', 2.52, 3.45], 'АС Брокер': ['BC', 2.75, 3.65]}
            # new_result - Это словарь получаемый из файлов в папке AS, где Номер вопроса: 'Уровень готовности надёжности', 'Вес требования (ВНД 5492)'

            new_result = {}
            new_result_result_treb = {}
            new_rez_as = rezult_as.copy()
            result_as_new = result_stack.copy()
            total_sum_all_treb_3 = 0
            total_sum_all_treb_4 = 0
            total_sum_release_treb_3 = 0
            total_sum_release_treb_4 = 0



            # Изменяем инициализацию словарей
            aga_three = {quarter: {}}  # Структура: {quarter: {system: [values]}}
            aga_four = {quarter: {}}


            if not os.path.isdir(adress_as):
                print(f"Ошибка: папка '{adress_as}' не найдена.")
                return

            for filename in os.listdir(adress_as):
                if filename.endswith(".xlsx"):
                    file_path = os.path.join(adress_as, filename)
                    try:
                        wb = openpyxl.load_workbook(file_path, data_only=True)
                        ws = wb['Технологии']
                        wm = wb['Титульный']
                        new_result[wm['C4'].value] = []
                        new_result_result_treb[wm['C4'].value] = []

                        for row in range(2, ws.max_row + 1):
                            keys = ws.cell(row, 1).value  # Столбец A

                            value_l = ws.cell(row, 14).value  # Столбец N
                            value_n = ws.cell(row, 16).value  # Столбец P

                            # new_result[keys] = [value_l, value_n]
                            #Вывод количество выполненных требований к общему числу из Титульной страницы
                            # Определяем диапазоны ячеек
                            if wm['E18'].value == 0 or wm['G18'] == 0:
                                range_1 = 'E10:E18'
                                range_3 = 'G10:G18'
                                range_2 = 'F10:F18'
                                range_4 = 'H10:H18'
                            else:
                                range_1 = 'E10:E17'
                                range_3 = 'G10:G17'
                                range_2 = 'F10:F17'
                                range_4 = 'H10:H17'

                            # Функция для получения списка значений из диапазона
                            def get_values_from_range(wm, range_str):
                                values = []
                                for row in wm[range_str]:
                                    for cell in row:
                                        # Обрабатываем три случая: '-', None и другие значения
                                        if cell.value == '-' or cell.value is None:
                                            values.append(0)
                                        else:
                                            values.append(cell.value)
                                return values

                            # Получаем списки значений для обоих диапазонов
                            values_1 = get_values_from_range(wm, range_1)
                            values_2 = get_values_from_range(wm, range_3)
                            values_3 = get_values_from_range(wm, range_2)
                            values_4 = get_values_from_range(wm, range_4)

                            new_result[wm['C4'].value].append((keys, value_l, value_n))

                            new_result_result_treb[wm['C4'].value].append((sum(values_1),  sum(values_3), sum(values_2), sum(values_4)))


                    except Exception as e:
                        print(f"Ошибка при открытии файла '{file_path}': {e}")


            #Получили реестр требований по всем АС для дальнейшего расчета и подсчета индекса надежности. Очень ВАЖНО!!!
            reestr_trebovan = {key: list(set(value)) for key, value in new_result_result_treb.items()}

            #Это нужно, чтобы удалить те АС-ки, которые не вошли в реестр ЕДК, чтобы не считать их данные мы их исключим!!!
            for id_key, vial_list in mas_it_as.items():
                if id_key in reestr_trebovan:
                    del reestr_trebovan[id_key]

            clean_new_result = {}

            for key, value in new_result.items():
                if key is not None and (not isinstance(value, list) or not any(item is None or item == '-' for item in value)):
                    clean_new_result[key] = value
            for keys in new_result:

                if keys in result_stack[quarter]:
                    for item in new_result[keys]:
                        # Проверяем, что item не None и имеет хотя бы один элемент
                        if item is not None and len(item) > 0:

                            # Проверяем, что первый элемент item существует и является строкой
                            # if isinstance(item[0], str) and item[0].isdigit():
                            if isinstance(item[0], (int, float)):
                                for item2 in result_stack[quarter][keys]:


                                    # Проверяем, что item2 не None и имеет хотя бы один элемент
                                    if item2 is not None and len(item2) > 0 and item2[0].isdigit():

                                        # Проверяем первый элемент item2
                                        if item[0] == int(item2[0]):

                                            if item2[12] in ['Устранено - подтверждено ПАО',
                                                                  'Исключение из реестра - подтверждено ПАО',
                                                                  '(ПАО) Устранено - подтверждено ПАО',
                                                                  '(ПАО) Исключение из реестра - подтверждено ПАО']:

                                                if item[1] == 3:
                                                    if keys not in aga_three[quarter]:
                                                        aga_three[quarter][keys] = []
                                                    if item[2] != '-':
                                                        aga_three[quarter][keys].append(round(float(item[2]), 6))


                                                # print(aga_three)
                                                if item[1] == 4:
                                                    if keys not in aga_four[quarter]:
                                                        aga_four[quarter][keys] = []
                                                    if item[2] != '-':
                                                        aga_four[quarter][keys].append(round(float(item[2]), 6))

            # print(aga_three)
            # print(aga_four)


            for key_three in aga_three[quarter]:
                aga_three[quarter][key_three] = [sum(aga_three[quarter][key_three])]

            for key_four in aga_four[quarter]:
                aga_four[quarter][key_four] = [sum(aga_four[quarter][key_four])]
                # print('3 ---->',aga_three)
                # print('4 ---->',aga_four)

            for new, val in new_rez_as.items():

                if new in aga_three:
                    val[1] = val[1] + round(aga_three[quarter][new][0], 2)
                if new in aga_four:
                    val[2] = val[2] + round(aga_four[quarter][new][0], 2)

                # print('3 ---->',aga_three)
                # print('4 ---->',aga_four)

            reality_dict_all = {quarter: {}}

            for ars_key, ars_val in reestr_trebovan.items():

                rs_val_three = aga_three[quarter].get(ars_key, None)
                rs_val_four = aga_four[quarter].get(ars_key, None)

                if ars_key not in reality_dict_all[quarter]:
                    reality_dict_all[quarter][ars_key] = []

                if ars_key in aga_three[quarter].keys() or ars_key in aga_four[quarter].keys():
                    total_requirements = sum([float(x) for x in (ars_val[0][0], ars_val[0][2])])
                    total_on_three = ars_val[0][0]
                    total_on_four = ars_val[0][2]

                    if rs_val_three is None:
                        completed_on_three = ars_val[0][1]
                    else:
                        real_val_three = int(rs_val_three[0])
                        if real_val_three > 0:
                            completed_on_three = real_val_three + ars_val[0][1]  # Используем преобразованное значение
                        else:
                            completed_on_three = ars_val[0][0] * sum(rs_val_three) + ars_val[0][1]

                    if rs_val_four is None:
                        completed_on_four = ars_val[0][3]
                    else:
                        real_val_four = int(rs_val_four[0])
                        if real_val_four > 0:
                            completed_on_four = real_val_four + ars_val[0][1]  # Используем преобразованное значение
                        else:
                            completed_on_four = ars_val[0][2] * sum(rs_val_four) + ars_val[0][3]



                    reality_dict_all[quarter][ars_key].append([
                        'Всего требований', total_requirements,
                        'Всего на 3-ку', total_on_three,
                        'Выполненных на 3-ку', completed_on_three,
                        'Всего на 4-ку', total_on_four,
                        'Выполненных на 4-ку', completed_on_four
                    ])
                else:
                    reality_dict_all[quarter][ars_key].append([
                        'Всего требований', sum([float(x) for x in (ars_val[0][0], ars_val[0][2])]),
                        'Всего на 3-ку', ars_val[0][0],
                        'Выполненных на 3-ку', ars_val[0][0] * sum(rs_val_three or []) + ars_val[0][1],
                        'Всего на 4-ку', ars_val[0][2],
                        'Выполненных на 4-ку', ars_val[0][2] * sum(rs_val_four or []) + ars_val[0][3]])



                # print(reality_dict_all)
            # Собрали весь пул требований для расчета из словаря по АС(reality_dict_all) и по Процессам (updated_dict1)
            # В функции будем расчитывать индекс надежности и готовности по АС и по Процессам
            # 1. Индекс надежности = (Сумма выполненных требований на 3-ку и 4-ку по процессам и АС)/ Сумма всех требований на 3-ку и 4-ку по процессам и АС

            # print(reality_dict_all)

            index_nadejnost_proc = {quarter: {
                'Всего требований по ITSM': [],
                'Всего требований по ITSM на 3': [],
                'Выполненно требований по ITSM на 3': [],
                'Всего требований по ITSM на 4': [],
                'Выполненно требований по ITSM на 4': []
            } for quarter in realudate}

            for quarter in realudate:
                data = realudate[quarter]

                # Всего требований по ITSM
                index_nadejnost_proc[quarter]['Всего требований по ITSM'].append(data['Всего требований'])
                # Всего требований по ITSM на 3
                index_nadejnost_proc[quarter]['Всего требований по ITSM на 3'].append(data['Всего требований на 3'])
                # Выполнено требований по ITSM на 3
                index_nadejnost_proc[quarter]['Выполненно требований по ITSM на 3'].append(data['Всего выполнено требований на 3'])
                # Всего требований по ITSM на 4
                index_nadejnost_proc[quarter]['Всего требований по ITSM на 4'].append(data['Всего требований на 4'])
                # Выполнено требований по ITSM на 4
                index_nadejnost_proc[quarter]['Выполненно требований по ITSM на 4'].append(data['Всего выполнено требований на 4'])

            # print(index_nadejnost_proc)

            index_nadejnost_as = {
                quarter: {
                    'Всего требований по АС': [],
                    'Всего требования по АС на 3': [],
                    'Выполненно требований по АС на 3': [],
                    'Всего требований по АС на 4': [],
                    'Выполненно требований по АС на 4': []
                } for quarter in reality_dict_all
            }

            for quarter in reality_dict_all:
                for as_name, as_data in reality_dict_all[quarter].items():
                    # as_data - это список, содержащий один элемент - другой список с данными
                    data_list = as_data[0]

                    # Создаем словарь для удобного доступа к данным
                    data_dict = {data_list[i]: data_list[i+1] for i in range(0, len(data_list), 2)}

                    # Заполняем index_nadejnost_as
                    index_nadejnost_as[quarter]['Всего требований по АС'].append(data_dict.get('Всего требований'))
                    index_nadejnost_as[quarter]['Всего требования по АС на 3'].append(data_dict.get('Всего на 3-ку', 0))
                    index_nadejnost_as[quarter]['Выполненно требований по АС на 3'].append(round(data_dict.get('Выполненных на 3-ку', 0) , 0))
                    index_nadejnost_as[quarter]['Всего требований по АС на 4'].append(data_dict.get('Всего на 4-ку', 0))
                    index_nadejnost_as[quarter]['Выполненно требований по АС на 4'].append(round(data_dict.get('Выполненных на 4-ку', 0) , 0))

            # print(index_nadejnost_as)












            sum_index_nadejnost_proc = index_nadejnost_proc


            #     {quarter: {
            #     'Всего требований по ITSM': 0,
            #     'Всего требований по ITSM на 3': 0,
            #     'Выполненно требований по ITSM на 3': 0,
            #     'Всего требований по ITSM на 4': 0,
            #     'Выполненно требований по ITSM на 4': 0
            # }}

            sum_index_nadejnost_as = {}

            for key_as in index_nadejnost_as[quarter]:
                sum_index_nadejnost_as[key_as] = sum(index_nadejnost_as[quarter][key_as])


            # print(sum_index_nadejnost_as)

            #print(sum_index_nadejnost_proc)
            # print(sum_index_nadejnost_as)






            # Извлекаем все значения из словаря
            values_mas_as = list(mas_it_as.values())

            # Оставляем только числовые значения
            filtered_values_mas_as = [[v for v in value if isinstance(v, (int, float))] for value in values_mas_as]

            # Рассчитываем суммы
            total_requirements = sum(sum(value) for value in filtered_values_mas_as)
            total_requirements_on_3 = sum(value[0] for value in filtered_values_mas_as)
            executed_requirements_on_3 = total_requirements_on_3
            total_requirements_on_4 = sum(value[1] for value in filtered_values_mas_as)
            executed_requirements_on_4 = total_requirements_on_4

            # Создаем новый словарь
            result_dict_close_as = {quarter: {
                'Всего требований по АС': total_requirements,
                'Всего требования по АС на 3': total_requirements_on_3,
                'Выполненно требований по АС на 3': executed_requirements_on_3,
                'Всего требований по АС на 4': total_requirements_on_4,
                'Выполненно требований по АС на 4': executed_requirements_on_4
            }}

            # Обновление значений в словаре sum_index_nadejnost_as со словарем result_dict_close_as(это те AS которые не попали в ЕДК)


            for key in result_dict_close_as[quarter]:
                try:
                    sum_index_nadejnost_as[quarter][key] += result_dict_close_as[quarter][key]
                except KeyError:
                    pass

            all_dict_as_itsm = {quarter: {
                'Всего требований по ITSM и АС': [sum_index_nadejnost_proc[quarter]['Всего требований по ITSM'],
                                                  sum_index_nadejnost_as['Всего требований по АС']],
                'Всего требований на 3 по ITSM и АС': [sum_index_nadejnost_proc[quarter]['Всего требований по ITSM на 3'],
                                                       sum_index_nadejnost_as['Всего требования по АС на 3']],
                'Выполнено требований на 3 по ITSM и АС': [
                    sum_index_nadejnost_proc[quarter]['Выполненно требований по ITSM на 3'],
                    round(sum_index_nadejnost_as['Выполненно требований по АС на 3'])],
                'Всего требований на 4 по ITSM и АС': [sum_index_nadejnost_proc[quarter]['Всего требований по ITSM на 4'],
                                                       sum_index_nadejnost_as['Всего требований по АС на 4']],
                'Выполнено требований на 4 по ITSM и АС': [
                    sum_index_nadejnost_proc[quarter]['Выполненно требований по ITSM на 4'],
                    round(sum_index_nadejnost_as['Выполненно требований по АС на 4'])]
            }}

            # print(all_dict_as_itsm)
            # print(index_nadejnost_as)
            # print(index_nadejnost_proc)
            # print(realudate)

            return new_rez_as, all_dict_as_itsm

            # new_as_rezult = func_as_artefakt(result_stack, adress_as, rezult_as, updated_dict1)

        new_as_rezult_prognoz, realudate_prognoz = func_as_artefakt_prognoz(versioned_data_it_as, adress_as, rezult_as, uroven_tech_nadejnosti_prognoz, quarter)


        # print(realudate_prognoz)

        cleaned_prognoz = {
            quarter: {
                key: [value[0][0] if isinstance(value[0], list) else value[0], value[1]]
                for key, value in data.items()
            }
            for quarter, data in realudate_prognoz.items()
        }

        # print(cleaned_prognoz)


        nadejnost_prognoz = {
            quarter: {
                'Уровень готовности 3 (Надежность)': None,
                'Уровень готовности 4 (Надежность)': None
            }
        }

        for quar, quarter_data in cleaned_prognoz.items():
            if 'Выполнено требований на 3 по ITSM и АС' in quarter_data:
                nadejnost_prognoz[quar]['Уровень готовности 3 (Надежность)'] = round(
                    (sum(quarter_data['Выполнено требований на 3 по ITSM и АС']) /
                     (sum(quarter_data['Всего требований на 3 по ITSM и АС'])) + 2), 2)

            if 'Выполнено требований на 4 по ITSM и АС' in quarter_data:
                nadejnost_prognoz[quar]['Уровень готовности 4 (Надежность)'] = round(
                    (sum(quarter_data['Выполнено требований на 4 по ITSM и АС']) /
                     (sum(quarter_data['Всего требований на 4 по ITSM и АС'])) + 3), 2)

            else:
                print(f"{quar}: Ключ не найден")

        for i, (key, values) in enumerate(nadejnost_prognoz.items(), start=16):
            # Проверяем наличие данных для текущего квартала
            if (values['Уровень готовности 3 (Надежность)'] is not None or
                    values['Уровень готовности 4 (Надежность)'] is not None):

                # Запись данных
                second_workshet[f'F{i+count_pr}'] = key
                second_workshet[f'G{i+count_pr}'] = values['Уровень готовности 3 (Надежность)'] or ""
                second_workshet[f'H{i+count_pr}'] = values['Уровень готовности 4 (Надежность)'] or ""

                # Форматирование ячеек для строки 16 и (16 + count_pr)
                rows_to_format = [16, 16 + count_pr]  # Список строк для форматирования

                for row in rows_to_format:
                    for col in ['F', 'G', 'H']:
                        cell = second_workshet[f'{col}{row}']
                        cell.alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
                        cell.border = Border(
                            left=Side(style='thin'),
                            right=Side(style='thin'),
                            top=Side(style='thin'),
                            bottom=Side(style='thin')
                        )
        # print(nadejnost_prognoz)




    #==================== ОКОНЧАНИЕ ПРОГНОЗИРОВАНИЯ =======================================================================================================














    for cell in second_workshet[1]:
        if cell.value:
            cell.alignment = Alignment(horizontal="center", vertical="center")
            second_workshet['B1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['C1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['D1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['E1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['F1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['G1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['H1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['A16'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['B16'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['B17'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['B18'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['B20'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['B22'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['B25'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['B26'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['C17'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['C18'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['D17'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            second_workshet['D18'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            five_workshet['A1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            five_workshet['B1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            five_workshet['C1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            five_workshet['D1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            five_workshet['E1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)


            # five_workshet[f'B{count_as + 4}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            # five_workshet[f'C{count_as + 4}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            # five_workshet[f'D{count_as + 4}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            # five_workshet[f'B{count_as + 5}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            # five_workshet[f'C{count_as + 5}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
            # five_workshet[f'D{count_as + 5}'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)


    # print(new_as_statistic)
    # {'Address Processing': ['72/0', '10/0', '0/0'], 'Engy': ['83/0', '10/0', '0/0'], 'Gateway legacy': ['83/0', '10/0', '0/0'], 'OPS': ['0/0', '77/0', '10/0'], 'Shiptor ERP': ['75/0', '10/0', '0/0']}
    # print(rezult_as)
    # {'Address Processing': ['BC', 2.16, 3.27], 'Engy': ['BC', 2.02, 3], 'Gateway legacy': ['MC', 2.02, 3], 'OPS': ['OP', 2.03, 3], 'Shiptor ERP': ['MC', 2.09, 3.27]}
    #----------------------------------Нужно для сложения всех значений из АС чтобы получилась единое число ['313/313', '117/117', '10/10']

    # Инициализируем списки для сумм числителей и знаменателей
    sum_numerators = [0, 0, 0]
    sum_denominators = [0, 0, 0]

    # Проходим по каждому элементу словаря
    for key in new_as_statistic:
        for i in range(3):
            numerator, denominator = map(int, new_as_statistic[key][i].split('/'))
            sum_numerators[i] += numerator
            sum_denominators[i] += denominator

    # Формируем результат
    result = [
        f"{sum_numerators[0]} / {sum_denominators[0]}",
        f"{sum_numerators[1]} / {sum_denominators[1]}",
        f"{sum_numerators[2]} / {sum_denominators[2]}"
    ]
    #--------------------------------------------------
    # print(result)


    #-----------------Записываю результаты в Тепловую карту показания технологической надежности для АС--------------------
    six_workshet['R2'] = ", ".join(result)
    six_workshet['Q2'] = (
            str(itsm_dtn_all['Критичный']) + ' / ' + str(itsm_dtn_ikrement['Критичный']) + ', ' +
            str(itsm_dtn_all['Высокий']) + ' / ' + str(itsm_dtn_ikrement['Высокий']) + ', ' +
            str(itsm_dtn_all['Умеренный']) + ' / ' + str(itsm_dtn_ikrement['Умеренный'])
    )
    six_workshet['T2'] = second_workshet['B22'].value

    critical_value = second_workshet['B16'].value
    # Запись с проверкой допустимых значений
    if critical_value in ['OP', 'BO', 'BC и выше']:  # Если значение из допустимого списка
        six_workshet['S2'] = critical_value
        seven_workshet['S2'] = critical_value
    elif critical_value is not None:  # Если есть значение, но недопустимое
        six_workshet['S2'] = f"Некорректное значение: {critical_value}. Введите OP, BO, BC и выше"
        seven_workshet['S2'] = f"Некорректное значение: {critical_value}. Введите OP, BO, BC и выше"
    else:  # Если значение отсутствует
        six_workshet['S2'] = "Введите Критичность вручную (OP, BO, BC и выше)"
        seven_workshet['S2'] = "Введите Критичность вручную (OP, BO, BC и выше)"


    seven_workshet['R2'] = ", ".join(result)
    seven_workshet['Q2'] = (
            str(itsm_dtn_all['Критичный']) + ' / ' + str(itsm_dtn_ikrement['Критичный']) + ', ' +
            str(itsm_dtn_all['Высокий']) + ' / ' + str(itsm_dtn_ikrement['Высокий']) + ', ' +
            str(itsm_dtn_all['Умеренный']) + ' / ' + str(itsm_dtn_ikrement['Умеренный'])
    )

    seven_workshet['T2'] = second_workshet['B22'].value

    six_workshet['P2'] = 'Ключевые проблемы (максимум 5 шт.), выбираем вручную, каждая отделяется друг от друга точкой с запятой. (Для удобства в ячейке  используем alt+Enter)'
    seven_workshet['P2'] = 'Ключевые проблемы (максимум 5 шт.), выбираем вручную, каждая отделяется друг от друга точкой с запятой. (Для удобства в ячейке  используем alt+Enter)'


    six_workshet['A1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['B1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['C1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['D1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['E1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['F1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['G1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['H1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['I1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['J1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['K1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['L1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['M1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['N1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['O1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['P1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['Q1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['R1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['S1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['T1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)

    six_workshet['A2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['B2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['C2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['D2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['E2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['F2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['G2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['H2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['I2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['J2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['K2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['L2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['M2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['N2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['O2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['P2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['P2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['Q2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['R2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['S2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    six_workshet['T2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)




    seven_workshet['A1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['B1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['C1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['D1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['E1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['F1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['G1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['H1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['I1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['J1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['K1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['L1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['M1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['N1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['O1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['P1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['A2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['B2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['C2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['D2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['E2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['F2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['G2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['H2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['I2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['J2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['K2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['L2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['M2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['N2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['O2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['P2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['P1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['Q1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['R1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['S1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['T1'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['P2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['Q2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['R2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['S2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)
    seven_workshet['T2'].alignment = Alignment(horizontal="center", vertical="center", wrapText=True)


    #Для прорисовки таблицы на листе ---===Показатели Надежности===---
    for row in second_workshet['A17':'D18']:
        for cell in row:
            cell.border = thin_border

    for row in five_workshet[f'A{count_as + 4}':f'D{count_as + 5}']:
        for cell in row:
            cell.border = thin_border

    # for row in five_workshet[f'A{count_as + 4}':f'D{count_as + 5}']:
    #     for cell in row:
    #         cell.border = thin_border



        #================================================== Cохранение Отчета ===============================================================================
    clean_name = name_company.replace("ЕРО", "").strip()
    save = os.path.join(os.path.dirname(__file__), f'Расчет_{current_quarter}_{clean_name}')
    if os.path.exists(save):
        # k = f'{name_company}'
        k = 'Показатели надежности'
        print(f"Файл '{k}' уже существует.")
        # Запрашиваем у пользователя, хочет ли он перезаписать файл
        while True:
            perazapis = input("Вы хотите перезаписать файл Показателей надежности? (y/n): ")
            if perazapis == 'n':
                print(f"Файл '{k} не перезаписан.")
                break
            elif perazapis == 'y':
                wb.save(save)
                print(f"Файл '{k} перезаписан.")
                print('=====  Отчет сформирован  =====')
                break
            else:
                print("Неверный ответ. Пожалуйста, введите 'y' для перезаписи или 'n' для отказа.")
    else:
        wb.save(save)
        print('=====  Отчет сформирован  =====')
        #================================================== Конец Cохранения Отчета =====================================================================
